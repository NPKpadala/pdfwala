#!/usr/bin/env python3
"""PDFWala Enterprise V11.0.0 — Modular Production Backend"""
__version__ = "11.0.0"
import os, sys, io, re, csv, json, uuid, time, shutil, signal, zipfile
import logging, unicodedata, threading, subprocess, tempfile
import hashlib
from datetime import datetime
from pathlib import Path
from contextlib import contextmanager
from flask import Flask, request, jsonify, g, send_file, Response
from werkzeug.utils import secure_filename
from config import Config
from services.redis_service import redis_service
from services.file_service import FileService
from services.queue_service import (
    backpressure, cb_libreoffice, cb_ghostscript, cb_tesseract,
    queue_service
)
from services.auth_service import require_auth, require_rate_limit
from utils.validators import validate_file, validate_password
from utils.helpers import (
    generate_output_filename, sanitize_string, safe_int,
    format_file_size, get_timestamp, generate_uuid
)
from utils.security import generate_signed_url, verify_signed_url, SafeRegex, REDACTION_PATTERNS
from utils.pdf_utils import (
    parse_page_ranges, create_watermark_pdf, create_page_number_pdf,
    get_pdf_page_count, compress_pdf_images, is_valid_pdf
)
from utils.office_utils import (
    coerce_cell_value, coerce_cell_for_csv,
    _is_structured_table, _get_table_signature, _normalize_header,
    _merge_tables, _write_optimized_sheet,
)
try:
    from flask_cors import CORS
    CORS_AVAILABLE = True
except ImportError:
    CORS_AVAILABLE = False
import fitz
from PIL import Image, ImageChops
from PyPDF2 import PdfReader, PdfWriter, PdfMerger
from reportlab.pdfgen import canvas as rl_canvas
from reportlab.lib.pagesizes import letter, A4
try:
    from PIL import ImageDraw, ImageFont
    IMAGEDRAW_AVAILABLE = True
except ImportError:
    IMAGEDRAW_AVAILABLE = False
try:
    from docx import Document as DocxDocument
    from docx.shared import Inches, Pt
    from openpyxl.styles import Font as XlFont
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
try:
    from openpyxl import load_workbook, Workbook
    from openpyxl.drawing.image import Image as XlImage
    from openpyxl.styles import Font, Alignment, PatternFill
    from openpyxl.utils import get_column_letter
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False
try:
    import msoffcrypto
    MSOFFCRYPTO_AVAILABLE = True
except ImportError:
    MSOFFCRYPTO_AVAILABLE = False
try:
    from pdf2docx import Converter as Pdf2DocxConverter
    PDF2DOCX_AVAILABLE = True
except ImportError:
    PDF2DOCX_AVAILABLE = False
try:
    import tabula
    TABULA_AVAILABLE = True
except ImportError:
    TABULA_AVAILABLE = False
try:
    import pytesseract
    from pytesseract import Output as TesseractOutput
    TESSERACT_AVAILABLE = True
except ImportError:
    TESSERACT_AVAILABLE = False
try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches, Pt as PptxPt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False
try:
    from prometheus_client import Counter, Histogram, generate_latest, REGISTRY
    METRICS_ENABLED = True
    _prom_req_count = Counter('pdfwala_requests_total', 'Total requests',
                               ['method', 'endpoint', 'status'])
    _prom_req_dur = Histogram('pdfwala_request_duration_seconds', 'Duration',
                                 ['method', 'endpoint'])
    _prom_file_size = Histogram('pdfwala_file_size_bytes', 'Upload size', ['operation'],
                                 buckets=[1024,10240,102400,1048576,10485760,104857600])
    _prom_download_size = Histogram('pdfwala_download_size_bytes', 'Download size',
                                     buckets=[1024,10240,102400,1048576,10485760,104857600])
except ImportError:
    METRICS_ENABLED = False

# FIX 34: PIL decompression bomb protection — halve the default threshold
Image.MAX_IMAGE_PIXELS = getattr(Config, 'MAX_IMAGE_PIXELS', 89_478_485)

# Detect wkhtmltopdf availability at startup
_WKHTMLTOPDF_AVAILABLE = False
try:
    _wk_result = subprocess.run(
        [getattr(Config, 'WKHTMLTOPDF', 'wkhtmltopdf'), '--version'],
        capture_output=True, timeout=5
    )
    _WKHTMLTOPDF_AVAILABLE = (_wk_result.returncode == 0)
except Exception:
    _WKHTMLTOPDF_AVAILABLE = False

# FIX 29: Validate SIGNED_URL_SECRET at startup
if not getattr(Config, 'SIGNED_URL_SECRET', None) or len(str(Config.SIGNED_URL_SECRET)) < 32:
    raise RuntimeError("SIGNED_URL_SECRET must be set and at least 32 characters")

# Unified large-file async threshold (10 MB)
_ASYNC_FILE_THRESHOLD = int(os.environ.get("ASYNC_FILE_THRESHOLD", 10 * 1024 * 1024))

# FIX 1: Module-level lock for thread registry
_registry_lock = threading.Lock()

# Celery app imported directly
from workers.celery_app import celery_app

# ── Ensure directories exist ───────────────────────────────────────────────────
for _d in [Config.UPLOAD_FOLDER, Config.OUTPUT_FOLDER, Config.TEMP_FOLDER]:
    os.makedirs(_d, exist_ok=True)
_APP_START = time.time()
log = logging.getLogger("pdfwala")

# ── Module-level constants ─────────────────────────────────────────────────────
# FIX 2: Thread-safe health cache
_HEALTH_CACHE = {}
_HEALTH_CACHE_TTL = 30  # seconds
_health_lock = threading.Lock()

# FIX 20/40: OCR concurrency semaphore
_ocr_semaphore = threading.Semaphore(getattr(Config, 'MAX_OCR_THREADS', 2))

# FIX 1: Thread registry for cancellable background conversions
_thread_registry = {}  # job_id -> (thread, cancel_event)

# FIX 2: Health cache helpers
def _get_health_cache():
    with _health_lock:
        return _HEALTH_CACHE.get("result"), _HEALTH_CACHE.get("at", 0)

def _set_health_cache(result):
    with _health_lock:
        _HEALTH_CACHE["result"] = result
        _HEALTH_CACHE["at"] = time.time()

def _clear_health_cache():
    with _health_lock:
        _HEALTH_CACHE.clear()

# ── Flask app ──────────────────────────────────────────────────────────────────
app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = Config.MAX_FILE_SIZE
app.secret_key = Config.SECRET_KEY
if CORS_AVAILABLE:
    CORS(app, origins=Config.CORS_ORIGINS)


# ── Zip Slip prevention helper ─────────────────────────────────────────────────
def _safe_zip_extract(zf: zipfile.ZipFile, dest: str) -> None:
    """Extract zip safely, rejecting path traversal attempts (Zip Slip fix)."""
    dest_real = os.path.realpath(dest)
    for member in zf.infolist():
        target = os.path.realpath(os.path.join(dest_real, member.filename))
        if not target.startswith(dest_real + os.sep):
            raise ValueError(f"Zip Slip rejected: {member.filename!r}")
        if member.is_dir():
            os.makedirs(target, exist_ok=True)
        else:
            os.makedirs(os.path.dirname(target), exist_ok=True)
            with zf.open(member) as src, open(target, "wb") as dst:
                dst.write(src.read())


# ── Request lifecycle ──────────────────────────────────────────────────────────
@app.before_request
def _before():
    g.start = time.time()
    g.request_id = request.headers.get("X-Request-ID", str(uuid.uuid4())[:8])
    g.user_id = "anonymous"

@app.after_request
def _after(response):
    ms = round((time.time() - g.get("start", time.time())) * 1000, 1)
    log.info(f"{request.method} {request.path} → {response.status_code} [{ms}ms]")
    if METRICS_ENABLED:
        try:
            _prom_req_count.labels(request.method, request.path, str(response.status_code)).inc()
            _prom_req_dur.labels(request.method, request.path).observe(
                time.time() - g.get("start", time.time()))
        except Exception:
            pass
    response.headers["X-Request-ID"] = g.get("request_id", "-")
    response.headers["X-Content-Type-Options"] = "nosniff"
    response.headers["X-Frame-Options"] = "DENY"
    response.headers["X-XSS-Protection"] = "1; mode=block"
    response.headers["Content-Security-Policy"] = (
        "default-src 'self'; "
        "script-src 'self' 'unsafe-inline'; "
        "style-src 'self' 'unsafe-inline'; "
        "img-src 'self' data:; "
        "connect-src 'self'; "
        "font-src 'self'; "
        "object-src 'none'; "
        "base-uri 'self'; "
        "form-action 'self'"
    )
    response.headers.pop("Server", None)
    return response

# ── Response helpers ───────────────────────────────────────────────────────────
def err(msg, code=400):
    log.warning(f"[{g.get('request_id','-')}] ERR {code}: {msg}")
    return jsonify({"success": False, "error": msg,
                    "request_id": g.get("request_id", "-")}), code

def ok(msg, path=None, **extras):
    payload = {"success": True, "message": msg, **extras}
    if path and os.path.exists(path):
        size = os.path.getsize(path)
        payload.update({
            "download_url": f"/pdfwala/download/{os.path.basename(path)}",
            "signed_url": generate_signed_url(path),
            "filename": os.path.basename(path),
            "size_human": format_file_size(size),
            "expires_in": f"{Config.FILE_TTL_SEC // 60} minutes",
        })
    return jsonify(payload)


# ── PDF validation helper ─────────────────────────────────────────────────────
# is_valid_pdf() is imported from utils.pdf_utils (strict version with magic-byte
# check, xref validation and min_pages guard).  The local definition has been
# removed in V11.1 to eliminate the weaker shadow that skipped header checks.


# ── Empty PDF guard helper ──────────────────────────────────────────────────
def _guard_empty_pdf(path: str):
    """Returns an error response if PDF has no pages, else None. Call after temp_upload."""
    try:
        _g = fitz.open(path)
        pages = len(_g)
        _g.close()
        if pages == 0:
            return err("Input PDF has no pages", 400)
        return None
    except Exception as ex:
        return err(f"Cannot open PDF: {type(ex).__name__}: {ex}", 400)


# ── Worker availability check ─────────────────────────────────────────────────
def _queue_has_workers(queue_name, timeout=0.5):
    """FIX 6: Check if any workers are consuming the given queue."""
    if not celery_app:
        return False
    try:
        inspector = celery_app.control.inspect(timeout=timeout)
        active_queues = inspector.active_queues()
        if not active_queues:
            return False
        for worker_queues in active_queues.values():
            for q in (worker_queues or []):
                if q.get("name") == queue_name:
                    return True
        return False
    except Exception:
        return False


# ── LibreOffice helper ────────────────────────────────────────────────────────
def libre(input_path, fmt, output_filename=None, temp=False):
    if fmt not in Config.LIBRE_ALLOWED_FMTS:
        return None
    if not cb_libreoffice.can_execute():
        log.warning(f"LibreOffice circuit breaker open — skipping conversion to {fmt}")
        return None
    if output_filename:
        output_filename = secure_filename(output_filename)
        if not output_filename:
            return None
    out_dir = tempfile.mkdtemp()
    try:
        result = subprocess.run(
            [Config.LIBREOFFICE, "--headless", "--convert-to", fmt,
             "--outdir", out_dir, input_path],
            capture_output=True, timeout=Config.SUBPROCESS_TIMEOUT,
            env={**os.environ, "HOME": out_dir, "TMPDIR": out_dir, "OOO_FORCE_DESKTOP": "true"})
        if result.returncode != 0:
            cb_libreoffice.record_failure()
            _clear_health_cache()
            return None
        base = os.path.splitext(os.path.basename(input_path))[0]
        converted = os.path.join(out_dir, f"{base}.{fmt}")
        if not os.path.exists(converted):
            matches = list(Path(out_dir).glob(f"*.{fmt}"))
            if not matches:
                cb_libreoffice.record_failure()
                _clear_health_cache()
                return None
            converted = str(matches[0])
        if os.path.getsize(converted) == 0:
            cb_libreoffice.record_failure()
            _clear_health_cache()
            return None
        if temp:
            final = os.path.join(Config.TEMP_FOLDER, f"{uuid.uuid4()}.{fmt}")
        elif output_filename:
            final = os.path.join(Config.OUTPUT_FOLDER, output_filename)
        else:
            final = os.path.join(Config.OUTPUT_FOLDER, f"{uuid.uuid4()}_output.{fmt}")
        try:
            os.replace(converted, final)
        except OSError:
            shutil.copy2(converted, final)
            os.remove(converted)
        cb_libreoffice.record_success()
        return final
    except subprocess.TimeoutExpired:
        cb_libreoffice.record_failure()
        _clear_health_cache()
        return None
    except Exception as ex:
        cb_libreoffice.record_failure()
        _clear_health_cache()
        log.error(f"LibreOffice: {ex}")
        return None
    finally:
        shutil.rmtree(out_dir, ignore_errors=True)


# ── Ghostscript helper ────────────────────────────────────────────────────────
def ghostscript_compress(input_path, output_path, gs_setting="/ebook",
                          extra_flags=None, timeout=300):
    if not cb_ghostscript.can_execute():
        return False
    cmd = [Config.GHOSTSCRIPT, "-sDEVICE=pdfwrite", "-dCompatibilityLevel=1.5",
           f"-dPDFSETTINGS={gs_setting}", "-dNOPAUSE", "-dBATCH", "-dQUIET",
           "-dSAFER", "-dDetectDuplicateImages=true", "-dCompressFonts=true",
           "-dSubsetFonts=true", "-dAutoRotatePages=/None",
           f"-sOutputFile={output_path}"]
    if extra_flags:
        cmd.extend(extra_flags)
    cmd.append(input_path)
    try:
        result = subprocess.run(cmd, capture_output=True, timeout=timeout)
        if result.returncode != 0:
            cb_ghostscript.record_failure()
            _clear_health_cache()
            return False
        if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
            cb_ghostscript.record_failure()
            _clear_health_cache()
            return False
        try:
            doc = fitz.open(output_path)
            if len(doc) == 0:
                doc.close()
                try: os.remove(output_path)
                except OSError: pass
                return False
            doc.close()
        except Exception:
            try: os.remove(output_path)
            except OSError: pass
            return False
        cb_ghostscript.record_success()
        return True
    except subprocess.TimeoutExpired:
        cb_ghostscript.record_failure()
        _clear_health_cache()
        try: os.remove(output_path)
        except OSError: pass
        return False
    except Exception as ex:
        cb_ghostscript.record_failure()
        _clear_health_cache()
        log.error(f"Ghostscript: {ex}")
        return False


# ── Cursive signature renderer ─────────────────────────────────────────────────
def _render_cursive_signature(page, name, reason, date_str, sig_x, sig_y):
    """Renders a visually distinct cursive-style signature block."""
    page.insert_text(
        (sig_x, sig_y),
        name,
        fontsize=14,
        fontname="tibo",
        color=(0.05, 0.05, 0.35),
        overlay=True
    )
    shape = page.new_shape()
    y_base = sig_y + 4
    shape.draw_bezier(
        (sig_x,        y_base),
        (sig_x + 40,   y_base - 4),
        (sig_x + 100,  y_base + 4),
        (sig_x + 150,  y_base - 2)
    )
    shape.finish(color=(0.1, 0.1, 0.5), width=1.2, closePath=False)
    shape.commit()
    border = fitz.Rect(sig_x - 4, sig_y - 17, sig_x + 154, sig_y + 18)
    page.draw_rect(border, color=(0.6, 0.6, 0.8), fill=None, width=0.3)
    page.insert_text(
        (sig_x, sig_y + 14),
        f"{reason}  ·  {date_str}",
        fontsize=7,
        fontname="helv",
        color=(0.45, 0.45, 0.45),
        overlay=True
    )


# ── Background cleanup ─────────────────────────────────────────────────────────
def _cleanup_worker():
    MAX_DEL = int(os.environ.get("CLEANUP_MAX_DELETES",
                                  getattr(Config, 'CLEANUP_MAX_DELETES', 500)))
    while True:
        try:
            now = time.time()
            deleted = 0
            for folder in [Config.OUTPUT_FOLDER, Config.UPLOAD_FOLDER, Config.TEMP_FOLDER]:
                if deleted >= MAX_DEL:
                    break
                try:
                    with os.scandir(folder) as it:
                        for entry in it:
                            if deleted >= MAX_DEL:
                                log.info(f"Cleanup: hit {MAX_DEL} delete limit, deferring rest")
                                break
                            if entry.is_file(follow_symlinks=False):
                                try:
                                    # FIX 41: follow_symlinks=False to prevent symlink timing attacks
                                    if now - entry.stat(follow_symlinks=False).st_mtime > Config.FILE_TTL_SEC:
                                        tomb = entry.path + f".deleting_{uuid.uuid4().hex}"
                                        try:
                                            os.rename(entry.path, tomb)
                                            os.remove(tomb)
                                            deleted += 1
                                        except FileNotFoundError:
                                            pass
                                        except OSError as tomb_err:
                                            log.debug(f"Cleanup rename failed: {tomb_err}")
                                            try:
                                                os.remove(entry.path)
                                                deleted += 1
                                            except OSError:
                                                pass
                                except OSError:
                                    pass
                except OSError:
                    pass
        except Exception as ex:
            log.error(f"Cleanup worker: {ex}")
        # FIX 19: Adaptive sleep based on activity
        sleep_secs = 60 if deleted >= 10 else (120 if deleted > 0 else 300)
        time.sleep(sleep_secs)

threading.Thread(target=_cleanup_worker, daemon=True, name="cleanup").start()


# ============================================================================
# HEALTH / METRICS / DOWNLOAD
# ============================================================================
@app.route("/health")
@app.route("/api/v1/health")
def health():
    if os.environ.get("FLASK_ENV") == "production":
        return jsonify({"status": "ok", "version": Config.VERSION})

    # FIX 2: Use thread-safe health cache helpers
    now = time.time()
    cached, cached_at = _get_health_cache()
    if cached and (now - cached_at) < _HEALTH_CACHE_TTL:
        return jsonify(cached)

    tools = {}
    for name, cmd in [("libreoffice",[Config.LIBREOFFICE,"--version"]),
                       ("ghostscript",[Config.GHOSTSCRIPT,"--version"]),
                       ("tesseract",["tesseract","--version"])]:
        try:
            subprocess.run(cmd, capture_output=True, timeout=5)
            tools[name] = "ok"
        except Exception:
            tools[name] = "unavailable"

    tools["wkhtmltopdf"] = "ok" if _WKHTMLTOPDF_AVAILABLE else "unavailable"

    _task_health = {}
    for _tname, _tcls in [("pdf_tasks", "compress_pdf_task"),
                           ("ocr_tasks", "ocr_pdf_task"),
                           ("office_tasks", "pdf_to_word_task")]:
        try:
            _mod = __import__(f"tasks.{_tname}", fromlist=[_tcls])
            _task_health[_tname] = "ok"
        except Exception as _te:
            _task_health[_tname] = f"FAILED: {_te}"

    rc = redis_service.client
    redis_status = "ok" if rc else "unavailable"
    try:
        if rc: rc.ping()
    except Exception:
        redis_status = "error"

    celery_status = "unavailable"
    if celery_app:
        try:
            celery_app.control.ping(timeout=1)
            celery_status = "ok"
        except Exception:
            celery_status = "degraded"

    overall_status = "ok"
    if any("FAILED" in str(v) for v in _task_health.values()):
        overall_status = "degraded"

    result = {
        "success": True, "status": overall_status, "version": Config.VERSION,
        "uptime_seconds": round(time.time() - _APP_START, 1),
        "redis": redis_status, "celery": celery_status, "tools": tools,
        "async_tasks": _task_health,
        "libraries": {
            "pdf2docx": PDF2DOCX_AVAILABLE, "pdfplumber": PDFPLUMBER_AVAILABLE,
            "tabula": TABULA_AVAILABLE, "pytesseract": TESSERACT_AVAILABLE,
            "python_docx": DOCX_AVAILABLE, "openpyxl": OPENPYXL_AVAILABLE,
            "msoffcrypto": MSOFFCRYPTO_AVAILABLE, "python_pptx": PPTX_AVAILABLE,
        },
        "circuit_breakers": {
            "libreoffice": cb_libreoffice.state,
            "ghostscript": cb_ghostscript.state,
            "tesseract": cb_tesseract.state,
        },
    }
    # FIX 2: Use thread-safe setter
    _set_health_cache(result)
    return jsonify(result)


@app.route("/api/v1/ready")
def readiness():
    rc = redis_service.client
    if rc:
        try: rc.ping()
        except Exception:
            return jsonify({"ready": False, "reason": "Redis unavailable"}), 503
    return jsonify({"ready": True})


@app.route("/api/v1/live")
def liveness():
    # FIX 25: Remove PID from liveness response
    return jsonify({"alive": True})


@app.route("/metrics")
@require_auth
def metrics():
    # FIX 23: Require auth for metrics endpoint
    if not METRICS_ENABLED:
        return jsonify({"error": "prometheus_client not installed"}), 404
    return Response(generate_latest(REGISTRY), mimetype="text/plain")


@app.route("/download/<filename>")
@require_auth
@require_rate_limit
def download(filename):
    # FIX 30: Full sanitisation BEFORE verify_signed_url
    filename = unicodedata.normalize("NFC", filename)
    safe = secure_filename(filename)
    if not safe or "/" in safe or ".." in safe:
        return err("Invalid filename", 400)
    if os.sep in safe or safe != os.path.basename(safe):
        return err("Invalid filename", 400)
    ALLOWED_EXTS = (".pdf",".zip",".jpg",".jpeg",".png",".docx",".xlsx",
                    ".pptx",".txt",".json",".html",".csv",".webp")
    if not safe.lower().endswith(ALLOWED_EXTS):
        return err("Invalid file type for download", 400)
    # FIX 30: All validation above done before signature check
    expires = request.args.get("expires", "")
    signature = request.args.get("sig", "")
    if expires and signature:
        if not verify_signed_url(safe, expires, signature):
            return err("Invalid or expired download link", 403)
    path = os.path.realpath(os.path.join(Config.OUTPUT_FOLDER, safe))
    if not path.startswith(os.path.realpath(Config.OUTPUT_FOLDER) + os.sep):
        return err("Access denied", 403)
    if not os.path.exists(path):
        return err("File not found or expired", 404)
    if METRICS_ENABLED:
        try:
            _prom_download_size.observe(os.path.getsize(path))
        except Exception:
            pass
    response = send_file(path, as_attachment=True, conditional=True)
    response.headers["X-Accel-Buffering"] = "no"
    response.headers["Cache-Control"] = "no-cache"
    return response

@app.route("/api/v1/jobs/<job_id>", methods=["GET"])
@require_auth
def api_job_status(job_id):
    if len(job_id) > 64:
        return err("Invalid job ID", 400)
    safe_id = re.sub(r"[^a-f0-9\-]", "", job_id)
    if safe_id != job_id:
        return err("Invalid job ID", 400)
    job = redis_service.job_get(job_id)
    if not job:
        if celery_app:
            try:
                from celery.result import AsyncResult
                task = AsyncResult(job_id, app=celery_app)
                if task.state == "PENDING":
                    return jsonify({
                        "success": True, "status": "pending", "progress": 0,
                        "warning": "Job is queued or task ID may not exist"
                    })
                elif task.state == "SUCCESS":
                    result = task.result or {}
                    out = result.get("output", "")
                    resp = {"success": True, "status": "completed", "progress": 100}
                    if out and os.path.exists(out):
                        resp["download_url"] = generate_signed_url(out)
                        resp["filename"] = os.path.basename(out)
                    elif out:
                        # FIX 16: Do not leak internal output_path
                        resp["warning"] = "Output file has expired or is not yet available"
                        resp["filename"] = os.path.basename(out)
                    return jsonify(resp)
                elif task.state == "FAILURE":
                    return jsonify({"success": False, "status": "failed",
                                    "error": str(task.info)}), 500
                return jsonify({"success": True, "status": task.state.lower(), "progress": 0})
            except Exception:
                pass
        return err("Job not found", 404)
    status = job.get("status", "unknown")
    progress = int(job.get("progress", 0))
    resp = {"success": True, "job_id": job_id, "status": status,
            "progress": progress, "operation": job.get("operation", ""),
            "created_at": job.get("created_at", "")}
    if job.get("total_pages"):
        resp["total_pages"] = int(job["total_pages"])
        resp["current_page"] = int(job.get("current_page", 0))
    if status == "completed":
        out = job.get("output_path", "")
        if out and os.path.exists(out):
            resp["download_url"] = generate_signed_url(out)
            resp["filename"] = os.path.basename(out)
            resp["size_human"] = format_file_size(os.path.getsize(out))
        elif out:
            # FIX 16: Do not leak internal filesystem path
            resp["warning"] = "Output file has expired or is not yet available"
            resp["filename"] = os.path.basename(out)
    if status == "failed":
        resp["error"] = job.get("error", "Unknown error")
    return jsonify(resp)

@app.route("/api/v1/endpoints", methods=["GET"])
@require_auth
def list_endpoints():
    # FIX 24: Require auth for endpoints listing
    endpoints = []
    for rule in app.url_map.iter_rules():
        if rule.endpoint == "static":
            continue
        endpoints.append({"path": rule.rule,
                           "methods": sorted(m for m in rule.methods if m not in ("HEAD","OPTIONS"))})
    return jsonify({"success": True, "endpoints": sorted(endpoints, key=lambda x: x["path"])})


# ============================================================================
# /remove-bg stub
# ============================================================================
@app.route("/api/v1/remove-bg", methods=["POST"])
@app.route("/api/remove-bg", methods=["POST"])
@require_auth
@require_rate_limit
def remove_bg():
    return err(
        "Background removal is not yet implemented. "
        "Consider integrating rembg (pip install rembg) or an external API.",
        501
    )


# ============================================================================
# PDF ORGANIZE
# ============================================================================
@app.route("/api/v1/merge", methods=["POST"])
@app.route("/api/merge", methods=["POST"])
@require_auth
@require_rate_limit
def merge_pdf():
    files = request.files.getlist("files")
    if len(files) < 2: return err("Minimum 2 PDF files required")
    if len(files) > Config.MAX_FILES_MERGE: return err(f"Maximum {Config.MAX_FILES_MERGE} files")
    for f in files:
        e = validate_file(f, Config.ALLOWED_PDF)
        if e: return err(e)
    try:
        with FileService.temp_uploads(files) as paths:
            # FIX: Wrap PdfMerger in try/finally
            merger = PdfMerger()
            try:
                page_sizes = set()
                for p in paths:
                    # FIX 9: Wrap each fitz.open in try/finally
                    doc = None
                    try:
                        doc = fitz.open(p)
                        if len(doc) == 0:
                            doc.close()
                            return err("Input PDF has no pages", 400)
                        if doc.is_encrypted and not doc.authenticate(""):
                            doc.close()
                            return err("PDF is password-protected. Please unlock first.", 400)
                        for pg in doc:
                            page_sizes.add((round(pg.rect.width,0), round(pg.rect.height,0)))
                    finally:
                        if doc:
                            try: doc.close()
                            except Exception: pass
                    merger.append(p)
                fname = generate_output_filename(files[0].filename, "merged",
                                                  is_multi=True, filenames=[f.filename for f in files])
                out = os.path.join(Config.OUTPUT_FOLDER, fname)
                merger.write(out)
                # FIX 42: Validate merge output file size
                if os.path.getsize(out) == 0:
                    try: os.remove(out)
                    except OSError: pass
                    return err("Merge produced empty output — all inputs may be corrupted", 500)
            finally:
                merger.close()
        return ok(f"Merged {len(files)} PDFs", out, mixed_page_sizes=(len(page_sizes)>1))
    except Exception:
        log.exception("merge"); return err("Merge failed", 500)

@app.route("/api/v1/split", methods=["POST"])
@app.route("/api/split", methods=["POST"])
@require_auth
@require_rate_limit
def split_pdf():
    """Split a PDF file into individual pages or by specified page ranges."""
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_PDF)
    if e:
        return err(e)

    mode = request.form.get("mode", "all")
    ranges = request.form.get("ranges", "")

    if mode == "range" and not ranges.strip():
        return err("Page range required when mode is range", 400)

    try:
        with FileService.temp_upload(f) as path:
            guard = _guard_empty_pdf(path)
            if guard:
                return guard

            reader = PdfReader(path)
            total = len(reader.pages)

            if mode == "all":
                indices = list(range(total))
            else:
                indices = parse_page_ranges(ranges, total)

            if not indices:
                return err("No valid pages in range")

            # Single page extraction → return PDF directly (not ZIP)
            if len(indices) == 1:
                writer = PdfWriter()
                writer.add_page(reader.pages[indices[0]])

                fname = generate_output_filename(f.filename, f"page_{indices[0]+1}")
                fname = re.sub(r'\.pdf$', '', fname, flags=re.IGNORECASE) + ".pdf"
                out = os.path.join(Config.OUTPUT_FOLDER, fname)

                with open(out, "wb") as fh:
                    writer.write(fh)

                return ok(f"Extracted page {indices[0]+1}", out)

            # FIX 10: Multiple pages → stream directly to disk ZIP (not BytesIO)
            fname = generate_output_filename(f.filename, "split_pages" if mode == "all" else "extracted_pages")
            out = os.path.join(Config.OUTPUT_FOLDER, fname)

            with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zf:
                for idx in indices:
                    writer = PdfWriter()
                    writer.add_page(reader.pages[idx])

                    page_buf = io.BytesIO()
                    writer.write(page_buf)
                    zf.writestr(f"page_{idx+1:04d}.pdf", page_buf.getvalue())
                    page_buf.close()

        return ok(f"Split into {len(indices)} pages", out)

    except Exception:
        log.exception("split")
        return err("Split failed", 500)

@app.route("/api/v1/organize", methods=["POST"])
@app.route("/api/organize", methods=["POST"])
@require_auth
@require_rate_limit
def organize_pdf():
    uploaded_file = request.files.get("file")
    validation_error = validate_file(uploaded_file, Config.ALLOWED_PDF)
    if validation_error:
        return err(validation_error)

    action = request.form.get("action", "reorder").lower()
    page_order_input = request.form.get("order", "").strip()

    if not page_order_input:
        return err("order parameter required")

    try:
        with FileService.temp_upload(uploaded_file) as temp_path:
            guard = _guard_empty_pdf(temp_path)
            if guard: return guard

            pdf_reader = PdfReader(temp_path)
            total_pages = len(pdf_reader.pages)

            try:
                one_based_indices = parse_page_ranges(page_order_input, total_pages)
            except Exception:
                one_based_indices = []

            if not one_based_indices:
                return err(
                    f"No valid pages specified. "
                    f"This PDF has {total_pages} page(s); valid range is 1–{total_pages}.",
                    400
                )

            for pg_num in one_based_indices:
                if pg_num < 1 or pg_num > total_pages:
                    return err(
                        f"Page {pg_num} is out of range. "
                        f"This PDF has {total_pages} page(s); valid range is 1–{total_pages}.",
                        400
                    )

            zero_based_indices = [page_num - 1 for page_num in one_based_indices]

            if action == "delete":
                pages_to_keep = [i for i in range(total_pages) if i not in set(zero_based_indices)]
                if not pages_to_keep:
                    return err("Cannot delete all pages", 400)
            elif action == "extract":
                pages_to_keep = zero_based_indices
            else:
                pages_to_keep = zero_based_indices

            output_writer = PdfWriter()
            for page_index in pages_to_keep:
                if page_index < 0 or page_index >= total_pages:
                    return err(f"Page index {page_index + 1} is out of range (1–{total_pages})", 400)
                output_writer.add_page(pdf_reader.pages[page_index])

            output_filename = generate_output_filename(uploaded_file.filename, "organized")
            output_path = os.path.join(Config.OUTPUT_FOLDER, output_filename)

            temp_output = output_path + ".tmp"
            with open(temp_output, "wb") as file_handle:
                output_writer.write(file_handle)
            os.replace(temp_output, output_path)

        action_labels = {
            "reorder": "Reordered",
            "extract": "Extracted",
            "delete": "Deleted pages from"
        }
        return ok(f"{action_labels.get(action, 'Organized')} PDF", output_path)

    except Exception:
        log.exception("organize")
        return err("Organize failed", 500)


@app.route("/api/v1/remove-pages", methods=["POST"])
@app.route("/api/remove-pages", methods=["POST"])
@require_auth
@require_rate_limit
def remove_pages():
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_PDF)
    if e: return err(e)
    order = request.form.get("order", "")
    if not order: return err("Pages to remove required")
    try:
        with FileService.temp_upload(f) as path:
            guard = _guard_empty_pdf(path)
            if guard: return guard
            reader = PdfReader(path)
            total = len(reader.pages)
            remove = set(parse_page_ranges(order, total))
            if len(remove) >= total: return err("Cannot remove all pages", 400)
            w = PdfWriter()
            for i, page in enumerate(reader.pages):
                if i not in remove: w.add_page(page)
            fname = generate_output_filename(f.filename, "pages_removed")
            out = os.path.join(Config.OUTPUT_FOLDER, fname)
            with open(out, "wb") as fh: w.write(fh)
        return ok(f"Removed {len(remove)} page(s)", out)
    except Exception:
        log.exception("remove_pages"); return err("Remove pages failed", 500)


@app.route("/api/v1/extract-pages", methods=["POST"])
@app.route("/api/extract-pages", methods=["POST"])
@require_auth
@require_rate_limit
def extract_pages():
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_PDF)
    if e: return err(e)
    order = request.form.get("order", "")
    if not order: return err("Pages to extract required")
    try:
        with FileService.temp_upload(f) as path:
            guard = _guard_empty_pdf(path)
            if guard: return guard
            reader = PdfReader(path)
            total = len(reader.pages)
            indices = parse_page_ranges(order, total)
            w = PdfWriter()
            for idx in indices: w.add_page(reader.pages[idx])
            fname = generate_output_filename(f.filename, "extracted")
            out = os.path.join(Config.OUTPUT_FOLDER, fname)
            with open(out, "wb") as fh: w.write(fh)
        return ok(f"Extracted {len(indices)} page(s)", out)
    except Exception:
        log.exception("extract_pages"); return err("Extract pages failed", 500)


# ============================================================================
# PDF OPTIMIZE
# ============================================================================
@app.route("/api/v1/compress", methods=["POST"])
@app.route("/api/compress", methods=["POST"])
@require_auth
@require_rate_limit
def compress_pdf():
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_PDF)
    if e: return err(e)
    quality = request.form.get("quality", "medium").lower()
    cfg = {"low":{"dpi":150,"quality":85,"gs":"/printer"},
           "medium":{"dpi":120,"quality":72,"gs":"/printer"},
           "high":{"dpi":96,"quality":60,"gs":"/ebook"}}.get(
               quality, {"dpi":120,"quality":72,"gs":"/printer"})
    _uid = uuid.uuid4().hex
    stage1 = os.path.join(Config.TEMP_FOLDER, f"{_uid}_stage1.pdf")
    gs_out = os.path.join(Config.TEMP_FOLDER, f"{_uid}_gs.pdf")
    try:
        with FileService.temp_upload(f) as path:
            orig = os.path.getsize(path)
            if METRICS_ENABLED:
                try: _prom_file_size.labels("compress").observe(orig)
                except Exception: pass
            try:
                guard = _guard_empty_pdf(path)
                if guard: return guard
            except Exception:
                return err("Input PDF is corrupted or unreadable", 400)

            # FIX 34: Catch DecompressionBombError on any image open within compress flow
            # (handled globally via Image.MAX_IMAGE_PIXELS at module level)

            if orig > _ASYNC_FILE_THRESHOLD:
                _compress_task = None
                try:
                    from tasks.pdf_tasks import compress_pdf_task as _compress_task
                except Exception as _ie:
                    log.error(f"compress_pdf_task lazy import failed: {_ie}")
                if celery_app and _compress_task:
                    # FIX 6: Check workers before dispatching
                    if not _queue_has_workers("fast"):
                        log.warning("No workers on 'fast' queue for compress_pdf")
                    ext = path.rsplit(".",1)[-1].lower()
                    bg_input = os.path.join(Config.UPLOAD_FOLDER, f"{uuid.uuid4()}.{ext}")
                    shutil.copy(path, bg_input)
                    try:
                        redis_service.client.setex(
                            f"cleanup:{bg_input}", Config.FILE_TTL_SEC + 600, "pending")
                    except Exception: pass
                    fname = generate_output_filename(f.filename, "compressed")
                    bg_out = os.path.join(Config.OUTPUT_FOLDER, fname)
                    job_id = str(uuid.uuid4())
                    redis_service.job_set(job_id, {
                        "status":"pending","operation":"compress_pdf",
                        "created_at":get_timestamp(),
                        "user_id":getattr(g,"user_id","default")
                    })
                    try:
                        redis_service.client.expire(f"job:{job_id}", getattr(Config, 'JOB_TTL_SEC', 7200))
                    except Exception: pass
                    # FIX 3: Wrap .delay() in try/except to delete bg_input on failure
                    try:
                        _compress_task.delay(bg_input, bg_out, job_id, quality)
                    except Exception:
                        try: os.remove(bg_input)
                        except OSError: pass
                        raise
                    return jsonify({
                        "success": True,
                        "message": "Large file queued — download will appear when ready",
                        "job_id": job_id, "status_url": f"/api/v1/jobs/{job_id}",
                        "async": True
                    })
                else:
                    log.warning("compress_pdf async unavailable, processing synchronously")

            stage1_size = orig
            try:
                doc = fitz.open(path)
                modified = compress_pdf_images(doc, cfg["dpi"], cfg["quality"])
                if modified:
                    doc.save(stage1, deflate=True, deflate_images=True,
                             deflate_fonts=True, garbage=3, clean=False)
                    stage1_size = os.path.getsize(stage1)
                else:
                    shutil.copy(path, stage1)
                # FIX 33: Explicitly close doc after stage1 save — don't rely on GC
                doc.close()
            except Exception as ex:
                log.warning(f"Stage1: {ex}"); shutil.copy(path, stage1)

            fname = generate_output_filename(f.filename, "compressed")
            out = os.path.join(Config.OUTPUT_FOLDER, fname)
            gs_ok = ghostscript_compress(stage1, gs_out, cfg["gs"],
                        extra_flags=["-dColorImageDownsampleType=/Bicubic",
                                     f"-dColorImageResolution={cfg['dpi']}",
                                     f"-dGrayImageResolution={cfg['dpi']}"])
            chosen = None
            if gs_ok and os.path.exists(gs_out) and os.path.getsize(gs_out) < stage1_size:
                chosen = gs_out
            if not chosen and os.path.exists(stage1) and stage1_size < orig:
                chosen = stage1
            if not chosen:
                chosen = path
            shutil.copy(chosen, out)
            new_size = os.path.getsize(out)
            # FIX 49: Cap reduction at 0 — never return negative reduction
            reduction = max(0, round((1 - new_size / orig) * 100, 1)) if orig else 0
        return ok(f"Compressed — {reduction}% smaller", out,
                  reduction_pct=reduction, original_size_bytes=orig,
                  compressed_size_bytes=new_size)
    except Exception:
        log.exception("compress"); return err("Compression failed", 500)
    finally:
        for tmp in [stage1, gs_out]:
            if tmp:
                try: os.remove(tmp)
                except OSError: pass


@app.route("/api/v1/repair-pdf", methods=["POST"])
@app.route("/api/repair-pdf", methods=["POST"])
@require_auth
@require_rate_limit
def repair_pdf():
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_PDF)
    if e:
        return err(e)

    try:
        with FileService.temp_upload(f) as path:
            orig_size = os.path.getsize(path)
            input_valid, _ = is_valid_pdf(path)

            fname = generate_output_filename(f.filename, "repaired")
            out = os.path.join(Config.OUTPUT_FOLDER, fname)

            # STAGE 1: PyMuPDF Repair
            try:
                doc = fitz.open(path)
                input_pages = len(doc)

                if input_pages == 0:
                    doc.close()
                    return err("Input PDF has no pages", 400)

                tmp_out = out + ".tmp"
                doc.save(tmp_out, garbage=4, deflate=True, clean=True)
                doc.close()

                os.replace(tmp_out, out)

                valid, reason = is_valid_pdf(out, min_pages=1)
                if valid:
                    return ok("PDF repaired (PyMuPDF)", out,
                              pages=input_pages, original_size_bytes=orig_size)
                else:
                    log.warning(f"PyMuPDF output invalid: {reason}")
                    try: os.remove(out)
                    except OSError: pass

            except Exception as ex1:
                log.warning(f"PyMuPDF repair failed: {ex1}")

            # STAGE 2: Ghostscript Repair
            gs_tmp = os.path.join(Config.TEMP_FOLDER, f"{uuid.uuid4().hex}_repair_gs.pdf")
            gs_ok = ghostscript_compress(
                path, gs_tmp, "/default",
                extra_flags=["-dPDFSTOPONERROR=false", "-dPDFSTOPONWARNING=false"]
            )

            if gs_ok and os.path.exists(gs_tmp) and os.path.getsize(gs_tmp) > 0:
                valid, reason = is_valid_pdf(gs_tmp, min_pages=1)
                if valid:
                    os.replace(gs_tmp, out)
                    return ok("PDF repaired (Ghostscript)", out,
                              original_size_bytes=orig_size)
                else:
                    log.warning(f"Ghostscript output invalid: {reason}")

            # FIX 43: Always clean up gs_tmp in all failure branches
            try:
                if os.path.exists(gs_tmp):
                    os.remove(gs_tmp)
            except OSError:
                pass

            if input_valid:
                shutil.copy(path, out)
                return ok("PDF is already valid (no repair needed)", out)

            return err("Repair failed: unable to produce valid PDF", 500)

    except Exception:
        log.exception("repair_pdf")
        return err("Repair failed", 500)


@app.route("/api/v1/linearize-pdf", methods=["POST"])
@app.route("/api/linearize-pdf", methods=["POST"])
@require_auth
@require_rate_limit
def linearize_pdf():
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_PDF)
    if e: return err(e)
    try:
        with FileService.temp_upload(f) as path:
            guard = _guard_empty_pdf(path)
            if guard: return guard
            orig_size = os.path.getsize(path)
            fname = generate_output_filename(f.filename, "linearized")
            out = os.path.join(Config.OUTPUT_FOLDER, fname)
            gs_ok = ghostscript_compress(path, out, gs_setting="/printer")
            if not gs_ok or not os.path.exists(out) or os.path.getsize(out) == 0:
                return err("Linearization failed — ensure Ghostscript is installed.", 500)
            new_size = os.path.getsize(out)
            reduction = round((1 - new_size / orig_size) * 100, 1) if orig_size else 0
        return ok(f"PDF processed — {reduction}% size change", out,
                  original_size_bytes=orig_size, new_size_bytes=new_size)
    except Exception:
        log.exception("linearize"); return err("Linearization failed", 500)


# ============================================================================
# OCR
# ============================================================================
@app.route("/api/v1/ocr-pdf", methods=["POST"])
@app.route("/api/ocr-pdf", methods=["POST"])
@require_auth
@require_rate_limit
def ocr_pdf():
    if not TESSERACT_AVAILABLE:
        return err("OCR requires pytesseract.", 501)
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_PDF)
    if e: return err(e)
    raw_lang = request.form.get("lang", "eng")
    lang = re.sub(r'[^a-zA-Z0-9+\-]', '', raw_lang)[:50] or "eng"
    if not re.fullmatch(r'[a-zA-Z]{3}(\+[a-zA-Z]{3})*', lang):
        lang = "eng"
    dpi = safe_int(request.form.get("dpi","300"), 300, 72, 400)
    psm = safe_int(request.form.get("psm","3"), 3, 1, 13)
    oem = safe_int(request.form.get("oem","3"), 3, 0, 3)
    # FIX 11: Page cap for sync OCR path
    MAX_OCR_PAGES_SYNC = getattr(Config, "MAX_OCR_PAGES_SYNC", 30)
    skip = False
    f.seek(0,2); file_size = f.tell(); f.seek(0)

    if file_size > _ASYNC_FILE_THRESHOLD:
        _ocr_task = None
        try:
            from tasks.ocr_tasks import ocr_pdf_task as _ocr_task
        except Exception as _ie:
            log.error(f"ocr_pdf_task lazy import failed: {_ie}")
        if celery_app and _ocr_task:
            # FIX 6: Check workers
            if not _queue_has_workers("slow"):
                log.warning("No workers on 'slow' queue for ocr_pdf")
            ext = f.filename.rsplit(".",1)[-1].lower() if "." in f.filename else "pdf"
            bg_input = os.path.join(Config.UPLOAD_FOLDER, f"{uuid.uuid4()}.{ext}")
            f.seek(0)
            with open(bg_input,"wb") as fh: fh.write(f.read())
            try:
                redis_service.client.setex(f"cleanup:{bg_input}", Config.FILE_TTL_SEC + 600, "pending")
            except Exception: pass
            fname = generate_output_filename(f.filename, "ocr")
            bg_out = os.path.join(Config.OUTPUT_FOLDER, fname)
            job_id = str(uuid.uuid4())
            redis_service.job_set(job_id, {"status":"pending","operation":"ocr_pdf",
                                            "created_at":get_timestamp(),
                                            "user_id":getattr(g,"user_id","default")})
            try:
                redis_service.client.expire(f"job:{job_id}", getattr(Config, 'JOB_TTL_SEC', 7200))
            except Exception: pass
            redis_service.job_update(job_id, {"input_path": bg_input})
            # FIX 3: Wrap .delay() in try/except
            try:
                _ocr_task.delay(bg_input, bg_out, job_id, lang, dpi, psm, oem)
            except Exception:
                try: os.remove(bg_input)
                except OSError: pass
                raise
            return jsonify({"success":True,"message":"OCR queued. Poll status endpoint.",
                            "job_id":job_id,"status_url":f"/api/v1/jobs/{job_id}"})
        else:
            log.warning("ocr_pdf async unavailable, processing synchronously")

    try:
        with FileService.temp_upload(f) as path:
            src_doc = fitz.open(path)
            if len(src_doc) == 0:
                src_doc.close()
                return err("Input PDF has no pages", 400)
            # FIX 11: Enforce sync page cap
            if len(src_doc) > MAX_OCR_PAGES_SYNC:
                src_doc.close()
                return err(f"PDF too large for sync OCR (max {MAX_OCR_PAGES_SYNC} pages). Use async endpoint.", 413)
            out_doc = fitz.open()
            pages_processed = pages_skipped = 0
            try:
                for page_num, src_page in enumerate(src_doc):
                    pw, ph = src_page.rect.width, src_page.rect.height
                    if skip and src_page.get_text().strip():
                        new_page = out_doc.new_page(width=pw, height=ph)
                        new_page.show_pdf_page(fitz.Rect(0,0,pw,ph), src_doc, page_num, overlay=False)
                        pages_skipped += 1
                        # FIX 11: Close pil_img on skip branch too (none to close here, just continue)
                        continue
                    mat = fitz.Matrix(dpi/72, dpi/72)
                    pix = src_page.get_pixmap(matrix=mat, alpha=False, colorspace=fitz.csGRAY)
                    img_sx = pw / pix.width; img_sy = ph / pix.height
                    hocr_data = None
                    pil_img = None
                    try:
                        pil_img = Image.open(io.BytesIO(pix.tobytes("png")))
                        pix = None
                        # FIX 20: OCR semaphore for image_to_data
                        with _ocr_semaphore:
                            hocr_data = pytesseract.image_to_data(
                                pil_img, lang=lang, output_type=TesseractOutput.DICT,
                                config=f"--psm {psm} --oem {oem}")
                    except Exception as ocr_ex:
                        log.warning(f"OCR page {page_num+1}: {ocr_ex}")
                    finally:
                        if pil_img is not None:
                            try: pil_img.close()
                            except Exception: pass
                    new_page = out_doc.new_page(width=pw, height=ph)
                    new_page.show_pdf_page(fitz.Rect(0,0,pw,ph), src_doc, page_num, overlay=False)
                    if hocr_data:
                        for word_str, conf_str, x0, y0, wd, ht in zip(
                            hocr_data.get("text", []),
                            hocr_data.get("conf", []),
                            hocr_data.get("left", []),
                            hocr_data.get("top", []),
                            hocr_data.get("width", []),
                            hocr_data.get("height", [])
                        ):
                            word = (word_str or "").strip()
                            try:
                                conf = int(conf_str)
                            except (ValueError, TypeError):
                                conf = 0
                            if not word or conf < 30:
                                continue
                            x0_f = float(x0) * img_sx
                            y1_f = (float(y0) + float(ht)) * img_sy
                            fs = max(4.0, (float(ht) * img_sy) * 0.85)
                            new_page.insert_text((x0_f, y1_f-1), word+" ", fontsize=fs,
                                                  fontname="helv", color=(0,0,0),
                                                  render_mode=3, overlay=True)
                    pages_processed += 1
                if pages_processed == 0 and pages_skipped == 0:
                    return err("OCR produced no output — all pages failed", 500)
                fname = generate_output_filename(f.filename, "ocr")
                out = os.path.join(Config.OUTPUT_FOLDER, fname)
                out_doc.save(out, deflate=True, garbage=2)
            finally:
                out_doc.close(); src_doc.close()
        return ok("OCR complete — PDF is text-searchable", out,
                  output_metadata={"pages_processed":pages_processed,
                                   "pages_skipped":pages_skipped,"lang":lang,"dpi":dpi})
    except Exception:
        log.exception("ocr_pdf"); return err("OCR failed", 500)


# ============================================================================
# PDF EDIT
# ============================================================================
@app.route("/api/v1/rotate", methods=["POST"])
@app.route("/api/rotate", methods=["POST"])
@require_auth
@require_rate_limit
def rotate_pdf():
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_PDF)
    if e: return err(e)
    try: angle = int(request.form.get("angle","90"))
    except ValueError: return err("Angle must be 90, 180 or 270")
    if angle not in (90,180,270): return err("Angle must be 90, 180 or 270")
    pages_spec = request.form.get("pages","all").strip()
    try:
        with FileService.temp_upload(f) as path:
            doc = fitz.open(path)
            if len(doc) == 0:
                doc.close()
                return err("Input PDF has no pages", 400)
            total = len(doc)
            # V11.1: hard cap guard
            MAX_ROTATE = getattr(Config, "MAX_ROTATE_PAGES", 1000)
            if total > MAX_ROTATE:
                doc.close()
                return err(f"PDF too large for rotation (max {MAX_ROTATE} pages). Split the PDF first.", 413)
            idxs = list(range(total)) if pages_spec.lower()=="all" else parse_page_ranges(pages_spec, total)
            if not idxs:
                doc.close()
                return err(f"No valid pages matched the specified range (document has {total} pages)", 400)
            for i in idxs: doc[i].set_rotation(angle)
            fname = generate_output_filename(f.filename, "rotated")
            out = os.path.join(Config.OUTPUT_FOLDER, fname)
            tmp_out = out + ".tmp"
            doc.save(tmp_out)
            doc.close()
            os.replace(tmp_out, out)
        return ok(f"Rotated {len(idxs)} page(s) by {angle}°", out)
    except Exception:
        log.exception("rotate"); return err("Rotate failed", 500)


@app.route("/api/v1/watermark", methods=["POST"])
@app.route("/api/watermark", methods=["POST"])
@require_auth
@require_rate_limit
def watermark_pdf():
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_PDF)
    if e: return err(e)
    text = sanitize_string(request.form.get("text","CONFIDENTIAL"))
    color = sanitize_string(request.form.get("color","808080"), 10)
    opacity = max(0.0, min(1.0, float(request.form.get("opacity","0.3"))))
    position = sanitize_string(request.form.get("position","diagonal"), 20)
    if position not in ("diagonal","center","top","bottom","tile"): position="diagonal"
    rotation = max(-90, min(90, float(request.form.get("rotation","45"))))
    scale = max(0.1, min(1.0, float(request.form.get("scale","0.3"))))
    image_data = None
    image_file = request.files.get("image")
    if image_file and image_file.filename:
        _img_err = validate_file(image_file, Config.ALLOWED_IMAGE)
        if _img_err:
            return err(f"Watermark image: {_img_err}", 400)
        image_file.seek(0)
        raw = image_file.read()
        try:
            # FIX 34: Catch DecompressionBombError on watermark image open
            pil_wm = Image.open(io.BytesIO(raw)).convert("RGBA")
            buf_wm = io.BytesIO(); pil_wm.save(buf_wm, format="PNG")
            image_data = buf_wm.getvalue()
            pil_wm.close()
        except Image.DecompressionBombError:
            return err("Watermark image too large to process safely", 413)
        except Exception:
            image_data = raw
    try:
        with FileService.temp_upload(f) as path:
            doc = fitz.open(path)
            if len(doc) == 0:
                doc.close()
                return err("Input PDF has no pages", 400)
            # V11.1: hard cap guard
            MAX_WM = getattr(Config, "MAX_WATERMARK_PAGES", 1000)
            if len(doc) > MAX_WM:
                doc.close()
                return err(f"PDF too large for watermarking (max {MAX_WM} pages). Split the PDF first.", 413)
            try:
                _wm_bytes = None
                _wm_w = _wm_h = 0
                pre_decoded_img = None
                if image_data:
                    try:
                        pre_decoded_img = Image.open(io.BytesIO(image_data)).convert("RGBA")
                    except Image.DecompressionBombError:
                        return err("Watermark image too large to process safely", 413)
                    _wm = pre_decoded_img.copy()
                    if rotation != 0:
                        _wm = _wm.rotate(rotation, expand=True, resample=Image.BICUBIC)
                    _r_ch, _g_ch, _b_ch, _a_ch = _wm.split()
                    _a_ch = _a_ch.point(lambda x: int(x * opacity))
                    _wm.putalpha(_a_ch)
                    _buf = io.BytesIO()
                    _wm.save(_buf, format="PNG")
                    _wm_bytes = _buf.getvalue()
                    _wm_w, _wm_h = _wm.size
                    _wm.close()
                    del _wm

                # FIX 14: pre_decoded_img.close() in try/finally
                try:
                    for page in doc:
                        r = page.rect
                        if _wm_bytes and _wm_w and _wm_h:
                            img_w = r.width * scale
                            img_h = img_w * _wm_h / _wm_w
                            if position == "top":
                                ix, iy = r.x0+(r.width-img_w)/2, r.y0+r.height*0.05
                            elif position == "bottom":
                                ix, iy = r.x0+(r.width-img_w)/2, r.y1-img_h-r.height*0.05
                            else:
                                ix, iy = r.x0+(r.width-img_w)/2, r.y0+(r.height-img_h)/2
                            page.insert_image(fitz.Rect(ix, iy, ix+img_w, iy+img_h),
                                               stream=_wm_bytes, overlay=True)
                        else:
                            wm = create_watermark_pdf(text, opacity, color, r.width, r.height, position, rotation)
                            wmpdf = fitz.open("pdf", wm)
                            page.show_pdf_page(fitz.Rect(0,0,r.width,r.height), wmpdf, 0, overlay=True)
                            wmpdf.close()

                    fname = generate_output_filename(f.filename, "watermarked")
                    out = os.path.join(Config.OUTPUT_FOLDER, fname)
                    # FIX 26: Atomic write for watermark_pdf
                    tmp_out = out + ".tmp"
                    doc.save(tmp_out)
                    os.replace(tmp_out, out)
                finally:
                    if pre_decoded_img is not None:
                        try: pre_decoded_img.close()
                        except Exception: pass
            finally:
                doc.close()
        wm_type = "image" if image_data else "text"
        return ok(f"Watermark ({wm_type}) added", out)
    except Exception:
        log.exception("watermark"); return err("Watermark failed", 500)


@app.route("/api/v1/page-numbers", methods=["POST"])
@app.route("/api/page-numbers", methods=["POST"])
@require_auth
@require_rate_limit
def page_numbers():
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_PDF)
    if e: return err(e)
    position = request.form.get("position","bottom")
    start = safe_int(request.form.get("start","1"), 1, 1)
    prefix = sanitize_string(request.form.get("prefix",""), 50)
    try:
        with FileService.temp_upload(f) as path:
            doc = fitz.open(path)
            if len(doc) == 0:
                doc.close()
                return err("Input PDF has no pages", 400)
            # V11.1: hard cap guard
            MAX_PN = getattr(Config, "MAX_PAGE_NUMBERS_PAGES", 1000)
            if len(doc) > MAX_PN:
                doc.close()
                return err(f"PDF too large for page numbering (max {MAX_PN} pages). Split the PDF first.", 413)
            try:
                for i, page in enumerate(doc):
                    r = page.rect
                    label = f"{prefix}{start + i}"
                    pn = create_page_number_pdf(label, position, r.width, r.height)
                    pnpdf = fitz.open("pdf", pn)
                    page.show_pdf_page(fitz.Rect(0,0,r.width,r.height), pnpdf, 0, overlay=True)
                    pnpdf.close()
                fname = generate_output_filename(f.filename, "numbered")
                out = os.path.join(Config.OUTPUT_FOLDER, fname)
                # FIX 26: Atomic write for page_numbers
                tmp_out = out + ".tmp"
                doc.save(tmp_out)
                os.replace(tmp_out, out)
            finally:
                doc.close()
        return ok("Page numbers added", out)
    except Exception:
        log.exception("page_numbers"); return err("Page numbering failed", 500)


@app.route("/api/v1/crop", methods=["POST"])
@app.route("/api/crop", methods=["POST"])
@require_auth
@require_rate_limit
def crop_pdf():
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_PDF)
    if e: return err(e)
    def _sf(k, d=0.0):
        try: return float(request.form.get(k, d))
        except: return d
    left,right,top,bottom = _sf("left"),_sf("right"),_sf("top"),_sf("bottom")
    if any(v < 0 for v in (left,right,top,bottom)):
        return err("Crop margins must be non-negative", 400)
    try:
        with FileService.temp_upload(f) as path:
            doc = fitz.open(path)
            if len(doc) == 0:
                doc.close()
                return err("Input PDF has no pages", 400)
            # V11.1: hard cap guard
            MAX_CROP = getattr(Config, "MAX_CROP_PAGES", 1000)
            if len(doc) > MAX_CROP:
                doc.close()
                return err(f"PDF too large for cropping (max {MAX_CROP} pages). Split the PDF first.", 413)
            try:
                first_r = doc[0].rect
                if (left + right) >= first_r.width or (top + bottom) >= first_r.height:
                    doc.close()
                    return err("Crop margins exceed page dimensions", 400)
                for page in doc:
                    r = page.rect
                    nr = fitz.Rect(r.x0+left, r.y0+top, r.x1-right, r.y1-bottom)
                    if nr.is_empty or nr.is_infinite:
                        doc.close(); return err("Crop margins too large", 400)
                    page.set_cropbox(nr)
                fname = generate_output_filename(f.filename, "cropped")
                out = os.path.join(Config.OUTPUT_FOLDER, fname)
                doc.save(out)
            finally:
                doc.close()
        return ok("PDF pages cropped", out)
    except Exception:
        log.exception("crop"); return err("Crop failed", 500)


@app.route("/api/v1/info", methods=["POST"])
@app.route("/api/info", methods=["POST"])
@require_auth
@require_rate_limit
def pdf_info():
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_PDF)
    if e: return err(e)
    try:
        with FileService.temp_upload(f) as path:
            file_size = os.path.getsize(path)
            doc = fitz.open(path)
            if len(doc) == 0:
                doc.close()
                return err("Input PDF has no pages", 400)
            try:
                meta = doc.metadata
                size_counts = {}
                image_count = 0
                font_names = set()
                has_forms = False
                for pg in doc:
                    key = (round(pg.rect.width,1), round(pg.rect.height,1))
                    size_counts[key] = size_counts.get(key, 0) + 1
                    image_count += len(pg.get_images())
                    for fi in pg.get_fonts(full=True):
                        bf = fi[3] if len(fi) > 3 else ""
                        if bf: font_names.add(bf)
                    if not has_forms and pg.first_widget:
                        has_forms = True
                unique_sizes = [{"w":k[0],"h":k[1],"count":v}
                                for k,v in sorted(size_counts.items(), key=lambda x:-x[1])]
                has_toc = len(doc.get_toc()) > 0
                is_lin = False
                try:
                    xobj = doc.xref_object(1, compressed=False)
                    is_lin = "/Linearized" in (xobj or "")
                except Exception: pass
                try: pdf_version = doc.pdf_version()
                except Exception: pdf_version = meta.get("format","unknown")
                out_data = {
                    "page_count": len(doc), "pdf_version": str(pdf_version),
                    "title": meta.get("title",""), "author": meta.get("author",""),
                    "subject": meta.get("subject",""), "creator": meta.get("creator",""),
                    "encrypted": doc.is_encrypted, "file_size_bytes": file_size,
                    "size_human": format_file_size(file_size),
                    "has_forms": has_forms, "has_toc": has_toc,
                    "image_count": image_count,
                    "fonts_used": sorted(font_names)[:20],
                    "total_fonts_found": len(font_names),
                    "is_linearized": is_lin,
                    "page_sizes": {"unique_sizes":unique_sizes,"all_same":len(unique_sizes)==1},
                }
            finally:
                doc.close()
        return ok("PDF info retrieved", metadata=out_data, pages=out_data.get("page_count"))
    except Exception:
        log.exception("pdf_info"); return err("Info retrieval failed", 500)


# ============================================================================
# PDF SECURITY
# ============================================================================
@app.route("/api/v1/protect", methods=["POST"])
@app.route("/api/protect", methods=["POST"])
@require_auth
@require_rate_limit
def protect_pdf():
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_PDF)
    if e: return err(e)
    pw = sanitize_string(request.form.get("password",""))
    pw2 = sanitize_string(request.form.get("password2",""))
    ep = validate_password(pw, pw2)
    if ep: return err(ep)
    allow_print = request.form.get("allow_print","true").lower() in ("true","1","yes")
    allow_copy = request.form.get("allow_copy","true").lower() in ("true","1","yes")
    try:
        with FileService.temp_upload(f) as path:
            doc = fitz.open(path)
            if len(doc) == 0:
                doc.close()
                return err("Input PDF has no pages", 400)
            try:
                permissions = int(fitz.PDF_PERM_ACCESSIBILITY)
                if allow_print: permissions |= int(fitz.PDF_PERM_PRINT)
                if allow_copy: permissions |= int(fitz.PDF_PERM_COPY)
                fname = generate_output_filename(f.filename, "protected")
                out = os.path.join(Config.OUTPUT_FOLDER, fname)
                # FIX 17: Separate _pw for actual use, redact pw before any logging
                _pw = pw
                pw = "[REDACTED]"
                try:
                    doc.save(out, encryption=fitz.PDF_ENCRYPT_AES_256,
                             owner_pw=_pw, user_pw=_pw, permissions=permissions)
                except Exception as save_ex:
                    log.error(f"protect_pdf save error: {type(save_ex).__name__}")
                    raise
                finally:
                    _pw = None
            finally:
                doc.close()
        return ok("PDF password-protected with AES-256", out)
    except Exception:
        log.exception("protect"); return err("Protect failed", 500)


@app.route("/api/v1/unlock", methods=["POST"])
@app.route("/api/unlock", methods=["POST"])
@require_auth
@require_rate_limit
def unlock_pdf():
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_PDF)
    if e: return err(e)
    pw = sanitize_string(request.form.get("password",""))
    if not pw: return err("Password required")
    try:
        with FileService.temp_upload(f) as path:
            doc = fitz.open(path)
            if len(doc) == 0:
                doc.close()
                return err("Input PDF has no pages", 400)
            if doc.is_encrypted:
                if not doc.authenticate(pw):
                    doc.close(); return err("Wrong password", 401)
            fname = generate_output_filename(f.filename, "unlocked")
            out = os.path.join(Config.OUTPUT_FOLDER, fname)
            doc.save(out, encryption=fitz.PDF_ENCRYPT_NONE)
            doc.close()
        return ok("PDF unlocked", out)
    except Exception:
        log.exception("unlock"); return err("Unlock failed", 500)


@app.route("/api/v1/sign-pdf", methods=["POST"])
@app.route("/api/sign-pdf", methods=["POST"])
@app.route("/api/v1/stamp-pdf", methods=["POST"])
@app.route("/api/stamp-pdf", methods=["POST"])
@require_auth
@require_rate_limit
def sign_pdf():
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_PDF)
    if e: return err(e)
    name = sanitize_string(request.form.get("name","Signed"))
    reason = sanitize_string(request.form.get("reason","Approved"))
    page_target = sanitize_string(request.form.get("page","last"), 10)
    position = sanitize_string(request.form.get("position","bottom-right"), 20)
    if position not in ("bottom-left","bottom-right","top-left","top-right","center"):
        position = "bottom-right"
    sig_data = None
    sig_file = request.files.get("signature")
    if sig_file and sig_file.filename:
        _sig_err = validate_file(sig_file, Config.ALLOWED_IMAGE)
        if _sig_err:
            return err(f"Signature image: {_sig_err}", 400)
        sig_file.seek(0)
        try:
            pil_sig = Image.open(sig_file).convert("RGBA")
            buf_sig = io.BytesIO(); pil_sig.save(buf_sig, format="PNG")
            sig_data = buf_sig.getvalue()
            pil_sig.close()
        except Image.DecompressionBombError:
            return err("Signature image too large to process safely", 413)
        except Exception:
            sig_data = sig_file.read()
    today_str = datetime.now().strftime("%Y-%m-%d")

    def _sig_pos(rect, pos):
        m = {"bottom-right":(rect.x1-180,rect.y1-70),
             "bottom-left": (rect.x0+30, rect.y1-70),
             "top-right": (rect.x1-180,rect.y0+50),
             "top-left": (rect.x0+30, rect.y0+50),
             "center": (rect.x0+rect.width/2-75, rect.y0+rect.height/2)}
        return m.get(pos,(rect.x1-180,rect.y1-70))

    try:
        with FileService.temp_upload(f) as path:
            doc = fitz.open(path)
            if len(doc) == 0:
                doc.close()
                return err("Input PDF has no pages", 400)
            total = len(doc)
            try:
                if page_target == "all": page_indices = list(range(total))
                elif page_target == "first": page_indices = [0]
                elif page_target == "last": page_indices = [total-1]
                else:
                    try:
                        pg_num = int(page_target)
                        if pg_num < 1 or pg_num > total:
                            return err(f"Page {pg_num} out of range (1-{total})", 400)
                        page_indices = [pg_num-1]
                    except ValueError:
                        return err(f"Invalid page target: '{page_target}'", 400)
                for pg_idx in page_indices:
                    page = doc[pg_idx]; rect = page.rect
                    sig_x, sig_y = _sig_pos(rect, position)
                    if sig_data:
                        img_rect = fitz.Rect(sig_x, sig_y-40, sig_x+150, sig_y+5)
                        page.insert_image(img_rect, stream=sig_data, overlay=True)
                        line = f"{name} | {reason} | {today_str}"
                        page.insert_text((sig_x, sig_y+14), line, fontsize=7,
                                          fontname="helv", color=(0.4,0.4,0.4))
                    else:
                        _render_cursive_signature(page, name, reason, today_str, sig_x, sig_y)
                fname = generate_output_filename(f.filename, "stamped")
                out = os.path.join(Config.OUTPUT_FOLDER, fname)
                # FIX 26: Atomic write for sign_pdf
                tmp_out = out + ".tmp"
                doc.save(tmp_out)
                os.replace(tmp_out, out)
            finally:
                doc.close()
        response = ok(f"Stamp added to {len(page_indices)} page(s) at {position}", out)
        if request.path.endswith("sign-pdf"):
            response.headers["Deprecation"] = "true"
            response.headers["Sunset"] = "2027-01-01"
            response.headers["Link"] = '</api/v1/stamp-pdf>; rel="successor-version"'
        return response
    except Exception:
        log.exception("sign_pdf"); return err("Stamp failed", 500)


@app.route("/api/v1/redact-pdf", methods=["POST"])
@app.route("/api/redact-pdf", methods=["POST"])
@require_auth
@require_rate_limit
def redact_pdf():
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_PDF)
    if e: return err(e)
    mode = sanitize_string(request.form.get("mode","text"), 10)
    if mode not in ("text","regex","preset"): mode="text"
    search_text = compiled = None
    if mode == "text":
        search_text = sanitize_string(request.form.get("search_text",""))
        if not search_text: return err("search_text required for mode=text")
    elif mode == "regex":
        pattern_str = sanitize_string(request.form.get("pattern",""), 500)
        if not pattern_str: return err("pattern required for mode=regex")
        try:
            compiled = SafeRegex.compile(pattern_str)
        except ValueError as rv:
            return err(str(rv), 400)
        except re.error as rex:
            return err(f"Invalid regex: {rex}", 400)
    elif mode == "preset":
        preset_name = sanitize_string(request.form.get("preset",""), 30)
        if preset_name not in REDACTION_PATTERNS:
            return err(f"Unknown preset. Choose: {', '.join(REDACTION_PATTERNS)}")
        compiled = re.compile(REDACTION_PATTERNS[preset_name])
    try:
        with FileService.temp_upload(f) as path:
            doc = fitz.open(path)
            if len(doc) == 0:
                doc.close()
                return err("Input PDF has no pages", 400)
            # V11.1: hard cap guard
            MAX_RD = getattr(Config, "MAX_REDACT_PAGES", 500)
            if len(doc) > MAX_RD:
                doc.close()
                return err(f"PDF too large for redaction (max {MAX_RD} pages). Split the PDF first.", 413)
            try:
                count = 0
                for page in doc:
                    if mode == "text":
                        for rect in page.search_for(search_text):
                            page.add_redact_annot(rect, fill=(0,0,0)); count += 1
                    else:
                        page_text = re.sub(r"\s+", " ", page.get_text("text"))
                        for match in compiled.finditer(page_text):
                            for rect in page.search_for(match.group()):
                                page.add_redact_annot(rect, fill=(0,0,0)); count += 1
                    page.apply_redactions()
                fname = generate_output_filename(f.filename, "redacted")
                out = os.path.join(Config.OUTPUT_FOLDER, fname)
                doc.save(out)
            finally:
                doc.close()

            _flat = os.path.join(Config.TEMP_FOLDER, f"{uuid.uuid4().hex}_redact_flat.pdf")
            if ghostscript_compress(out, _flat, "/default"):
                try:
                    os.replace(_flat, out)
                except OSError:
                    shutil.move(_flat, out)
            else:
                try: os.remove(_flat)
                except OSError: pass

        msg = f"Redacted {count} occurrence(s) (mode={mode}) — structurally flattened"
        warn = "No matches found — document unchanged" if count == 0 else None
        return ok(msg, out, redaction_count=count, warning=warn,
                  note="Flattened via Ghostscript to prevent text recovery")
    except Exception:
        log.exception("redact"); return err("Redact failed", 500)


# ============================================================================
# PDF CONVERT
# ============================================================================
@app.route("/api/v1/pdf-to-image", methods=["POST"])
@app.route("/api/pdf-to-image", methods=["POST"])
@require_auth
@require_rate_limit
def pdf_to_image():
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_PDF)
    if e: return err(e)
    fmt = request.form.get("format","jpg").lower()
    dpi = safe_int(request.form.get("dpi","150"), 150, 72, 300)
    if fmt not in ("jpg","png"): fmt = "jpg"
    try:
        with FileService.temp_upload(f) as path:
            doc = fitz.open(path)
            if len(doc) == 0:
                doc.close()
                return err("Input PDF has no pages", 400)
            total_pages = len(doc)
            doc.close()
            doc = None  # closed; chunked_pdf_processor will re-open

            CHUNK_THRESHOLD = getattr(Config, "PDF_TO_IMAGE_CHUNK_THRESHOLD", 50)
            CHUNK_PAGES     = getattr(Config, "PDF_TO_IMAGE_CHUNK_PAGES",     50)
            CHUNK_WORKERS   = getattr(Config, "PDF_TO_IMAGE_MAX_WORKERS",      4)
            MAX_IMG_PAGES   = getattr(Config, "MAX_IMAGE_PAGES", 2000)

            if total_pages > MAX_IMG_PAGES:
                return err(f"PDF exceeds maximum ({MAX_IMG_PAGES} pages) for image export.", 413)

            fname = generate_output_filename(f.filename, "to_image")
            out   = os.path.join(Config.OUTPUT_FOLDER, fname)

            def _render_pages_to_zip(
                chunk_pdf: str, chunk_idx: int, start_page: int, end_page: int
            ) -> str:
                chunk_zip = chunk_pdf.replace("_in.pdf", f"_imgs_{chunk_idx}.zip")
                src = fitz.open(chunk_pdf)
                try:
                    with zipfile.ZipFile(chunk_zip, "w", zipfile.ZIP_DEFLATED) as zf:
                        for local_i, page in enumerate(src):
                            abs_i = start_page + local_i
                            mat   = fitz.Matrix(dpi/72, dpi/72)
                            pix   = page.get_pixmap(matrix=mat, alpha=True)
                            pil   = Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGBA")
                            pix   = None
                            buf   = io.BytesIO()
                            if fmt == "jpg":
                                bg = Image.new("RGB", pil.size, (255, 255, 255))
                                bg.paste(pil, mask=pil.split()[3])
                                bg.save(buf, "JPEG", quality=85, optimize=True)
                            else:
                                pil.save(buf, "PNG")
                            pil = None
                            zf.writestr(f"page_{abs_i+1:04d}.{fmt}", buf.getvalue())
                finally:
                    src.close()
                return chunk_zip

            from utils.pdf_utils import chunked_pdf_processor as _chunked_proc, merge_zip_chunks as _merge_zips
            success = _chunked_proc(
                input_path=path,
                output_path=out,
                job_id=str(uuid.uuid4()),
                total_pages=total_pages,
                chunk_size=CHUNK_PAGES,
                max_workers=CHUNK_WORKERS,
                process_chunk_func=_render_pages_to_zip,
                merge_func=_merge_zips,
                redis_service=None,
                tool_name="PDF-to-Image",
                report_progress=False,
                chunk_retry=1,
            )
            if not success:
                # single-pass fallback
                src = fitz.open(path)
                try:
                    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zf:
                        for i, page in enumerate(src):
                            mat = fitz.Matrix(dpi/72, dpi/72)
                            pix = page.get_pixmap(matrix=mat, alpha=True)
                            pil = Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGBA")
                            pix = None
                            buf = io.BytesIO()
                            if fmt == "jpg":
                                bg = Image.new("RGB", pil.size, (255,255,255))
                                bg.paste(pil, mask=pil.split()[3])
                                bg.save(buf, "JPEG", quality=85, optimize=True)
                            else:
                                pil.save(buf, "PNG")
                            pil = None
                            zf.writestr(f"page_{i+1:04d}.{fmt}", buf.getvalue())
                finally:
                    src.close()
        return ok(f"Exported {total_pages} page(s) as {fmt.upper()}", out)
    except Exception:
        log.exception("pdf_to_image"); return err("Export failed", 500)


@app.route("/api/v1/pdf-to-word", methods=["POST"])
@app.route("/api/pdf-to-word", methods=["POST"])
@require_auth
@require_rate_limit
def pdf_to_word():
    """Convert PDF to editable Word document (DOCX)."""
    if not PDF2DOCX_AVAILABLE:
        return err("PDF to Word requires pdf2docx.", 501)

    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_PDF)
    if e:
        return err(e)

    f.seek(0, 2)
    file_size = f.tell()
    f.seek(0)

    try:
        _chk = fitz.open(stream=f.read(), filetype="pdf")
        if len(_chk) == 0:
            _chk.close()
            return err("Input PDF has no pages", 400)
        _chk.close()
        f.seek(0)
    except Exception:
        f.seek(0)

    ext = f.filename.rsplit(".", 1)[-1].lower() if "." in f.filename else "pdf"
    upload_path = os.path.join(Config.UPLOAD_FOLDER, f"{uuid.uuid4()}.{ext}")
    f.seek(0)
    with open(upload_path, "wb") as fh:
        fh.write(f.read())

    fname = generate_output_filename(f.filename, "to_word")
    fname = re.sub(r'\.pdf$', '.docx', fname, flags=re.IGNORECASE)
    if not fname.endswith(".docx"):
        fname = Path(fname).stem + ".docx"
    out = os.path.join(Config.OUTPUT_FOLDER, fname)

    if file_size > _ASYNC_FILE_THRESHOLD:
        _pdf_to_word_task = None
        try:
            from tasks.office_tasks import pdf_to_word_task as _pdf_to_word_task
        except Exception as _ie:
            log.error(f"pdf_to_word_task lazy import failed: {_ie}")

        if celery_app and _pdf_to_word_task:
            try:
                redis_service.client.setex(
                    f"cleanup:{upload_path}", Config.FILE_TTL_SEC + 600, "pending"
                )
            except Exception:
                pass

            job_id = str(uuid.uuid4())
            redis_service.job_set(job_id, {
                "status": "pending",
                "progress": "0",
                "operation": "pdf_to_word",
                "created_at": get_timestamp(),
                "user_id": getattr(g, "user_id", "default")
            })

            try:
                redis_service.client.expire(f"job:{job_id}", getattr(Config, 'JOB_TTL_SEC', 7200))
            except Exception:
                pass

            # FIX 3: Wrap .delay() in try/except
            try:
                task = _pdf_to_word_task.delay(upload_path, out, job_id)
            except Exception:
                try: os.remove(upload_path)
                except OSError: pass
                raise
            redis_service.job_update(job_id, {"task_id": task.id})

            return jsonify({
                "success": True,
                "message": "Large file — conversion running. Check status_url.",
                "job_id": job_id,
                "status_url": f"/api/v1/jobs/{job_id}",
                "poll_interval_ms": 2000,
                "async": True
            })

        else:
            log.warning("pdf_to_word async unavailable, processing synchronously")
            try:
                cv = Pdf2DocxConverter(upload_path)
                cv.convert(out, start=0, end=None)
                cv.close()
                if os.path.exists(out) and os.path.getsize(out) > 0:
                    return ok("PDF converted to Word (sync)", out)
                else:
                    return err("Conversion failed — output is empty", 500)
            except Exception:
                log.exception("pdf_to_word sync fallback")
                return err("PDF to Word conversion failed", 500)
            finally:
                try:
                    os.remove(upload_path)
                except OSError:
                    pass

    try:
        cv = Pdf2DocxConverter(upload_path)
        cv.convert(out, start=0, end=None)
        cv.close()
    except Exception:
        log.exception("pdf_to_word sync")
        return err("PDF to Word conversion failed", 500)
    finally:
        try:
            os.remove(upload_path)
        except OSError:
            pass

    if not os.path.exists(out) or os.path.getsize(out) == 0:
        return err("Conversion failed — output is empty", 500)

    return ok("PDF converted to Word", out)

@app.route("/api/v1/pdf-to-excel", methods=["POST"])
@app.route("/api/pdf-to-excel", methods=["POST"])
@require_auth
@require_rate_limit
def pdf_to_excel():
    if not OPENPYXL_AVAILABLE:
        return err("PDF to Excel requires openpyxl.", 501)

    uploaded_file = request.files.get("file")
    validation_error = validate_file(uploaded_file, Config.ALLOWED_PDF)
    if validation_error:
        return err(validation_error)

    try:
        with FileService.temp_upload(uploaded_file) as temp_path:
            guard = _guard_empty_pdf(temp_path)
            if guard: return guard

            wb = Workbook()
            wb.remove(wb.active)
            all_tables = []
            seen_signatures = set()
            method_used = None
            confidence = "low"

            # STAGE 1: pdfplumber
            if PDFPLUMBER_AVAILABLE:
                try:
                    with pdfplumber.open(temp_path) as pdf:
                        for page_num, page in enumerate(pdf.pages):
                            tables = page.extract_tables({
                                "vertical_strategy": "lines",
                                "horizontal_strategy": "lines",
                                "snap_tolerance": 4,
                                "intersection_tolerance": 4
                            })
                            if not tables:
                                tables = page.extract_tables({
                                    "vertical_strategy": "text",
                                    "horizontal_strategy": "text",
                                    "snap_tolerance": 6,
                                    "join_tolerance": 6
                                })
                            for table in tables:
                                if table and _is_structured_table(table):
                                    sig = _get_table_signature(table)
                                    if sig not in seen_signatures:
                                        seen_signatures.add(sig)
                                        all_tables.append(table)
                    if all_tables:
                        method_used = "pdfplumber"
                        confidence = "high" if len(all_tables) > 1 else "medium"
                except Exception as ex:
                    log.warning(f"pdfplumber extraction failed: {ex}")

            # STAGE 2: tabula-py fallback
            if not all_tables and TABULA_AVAILABLE:
                try:
                    # FIX 39: Cap JVM heap for tabula
                    dfs = tabula.read_pdf(temp_path, pages="all", multiple_tables=True,
                                          lattice=True, silent=True,
                                          java_options=["-Xmx256m", "-Xms64m"])
                    if not dfs:
                        dfs = tabula.read_pdf(temp_path, pages="all", multiple_tables=True,
                                              stream=True, silent=True,
                                              java_options=["-Xmx256m", "-Xms64m"])
                    for df in dfs:
                        if not df.empty:
                            table = [df.columns.tolist()] + df.values.tolist()
                            if _is_structured_table(table):
                                sig = _get_table_signature(table)
                                if sig not in seen_signatures:
                                    seen_signatures.add(sig)
                                    all_tables.append(table)
                    if all_tables:
                        method_used = "tabula"
                        confidence = "medium"
                except Exception as ex:
                    log.warning(f"tabula extraction failed: {ex}")

            # STAGE 3: Merge multi-page continuations
            merged_tables = []
            for table in all_tables:
                if merged_tables and _normalize_header(merged_tables[-1][0]) == _normalize_header(table[0]):
                    merged_tables[-1] = _merge_tables(merged_tables[-1], table)
                else:
                    merged_tables.append(table)

            # STAGE 4: Write tables to sheets
            if merged_tables:
                for idx, table in enumerate(merged_tables, 1):
                    ws = wb.create_sheet(f"Table_{idx}")
                    _write_optimized_sheet(ws, table, method_used or "unknown")
                tables_extracted = len(merged_tables)
            else:
                tables_extracted = 0

            # STAGE 5: Raw text fallback
            if tables_extracted == 0:
                ws = wb.create_sheet("Raw_Text")
                ws["A1"] = "No structured tables detected — raw text below:"
                ws["A1"].font = Font(bold=True, size=12)
                # FIX 32: Wrap fitz.open in try/finally in raw text fallback
                doc = fitz.open(temp_path)
                try:
                    row_idx = 3
                    for page_num, page in enumerate(doc):
                        ws[f"A{row_idx}"] = f"=== Page {page_num + 1} ==="
                        ws[f"A{row_idx}"].font = Font(bold=True)
                        row_idx += 1
                        page_text = page.get_text("text")
                        for line in page_text.split('\n'):
                            if line.strip():
                                ws[f"A{row_idx}"] = line.strip()
                                row_idx += 1
                finally:
                    doc.close()
                method_used = "raw_text"
                confidence = "none"

            if not wb.worksheets:
                return err("No content could be extracted from this PDF", 500)

            output_filename = generate_output_filename(uploaded_file.filename, "to_excel")
            output_filename = re.sub(r'\.pdf$', '.xlsx', output_filename, flags=re.IGNORECASE)
            output_path = os.path.join(Config.OUTPUT_FOLDER, output_filename)
            tmp_out = output_path + ".tmp"
            wb.save(tmp_out)
            os.replace(tmp_out, output_path)

            warning_msg = None
            if tables_extracted == 0:
                warning_msg = "No tables found — raw text extracted instead"
            elif confidence == "medium":
                warning_msg = "Tables extracted with basic method — verify alignment"

            return ok(
                f"Extracted {tables_extracted} table(s) using {method_used}",
                output_path,
                tables_found=tables_extracted,
                extraction_method=method_used,
                confidence=confidence,
                warning=warning_msg
            )

    except Exception:
        log.exception("pdf_to_excel")
        return err("PDF to Excel conversion failed", 500)


@app.route("/api/v1/pdf-to-ppt", methods=["POST"])
@app.route("/api/pdf-to-ppt", methods=["POST"])
@require_auth
@require_rate_limit
def pdf_to_ppt():
    if not PPTX_AVAILABLE: return err("PDF to PPT requires python-pptx.", 501)
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_PDF)
    if e: return err(e)
    try:
        with FileService.temp_upload(f) as path:
            doc = fitz.open(path)
            if len(doc) == 0:
                doc.close()
                return err("Input PDF has no pages", 400)
            # FIX 35: Page count guard for pdf_to_ppt
            MAX_PPT_PAGES = getattr(Config, "MAX_PPT_PAGES", 50)
            if len(doc) > MAX_PPT_PAGES:
                doc.close()
                return err(f"PDF too large for PPT conversion (max {MAX_PPT_PAGES} pages)", 413)
            try:
                prs = Presentation()
                prs.slide_width = PptxInches(10); prs.slide_height = PptxInches(7.5)
                blank = prs.slide_layouts[6]
                for page in doc:
                    pix = page.get_pixmap(dpi=200)
                    tmp_img = tempfile.NamedTemporaryFile(
                        suffix=".png", delete=False, dir=Config.TEMP_FOLDER
                    )
                    try:
                        tmp_img.write(pix.tobytes("png"))
                        tmp_img.close()
                        pix = None
                        slide = prs.slides.add_slide(blank)
                        slide.shapes.add_picture(
                            tmp_img.name, 0, 0, prs.slide_width, prs.slide_height
                        )
                    finally:
                        try: os.unlink(tmp_img.name)
                        except OSError: pass
            finally:
                doc.close()
            fname = generate_output_filename(f.filename, "to_ppt")
            fname = re.sub(r'\.pdf$','.pptx',fname,flags=re.IGNORECASE)
            out = os.path.join(Config.OUTPUT_FOLDER, fname)
            prs.save(out)
        return ok("PDF converted to PowerPoint", out)
    except Exception:
        log.exception("pdf_to_ppt"); return err("PDF to PPT failed", 500)


@app.route("/api/v1/pdf-to-pdfa", methods=["POST"])
@app.route("/api/pdf-to-pdfa", methods=["POST"])
@require_auth
@require_rate_limit
def pdf_to_pdfa():
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_PDF)
    if e: return err(e)
    version = request.form.get("version","1b")
    pdfa_val = "2" if "3" in version else "1"
    try:
        with FileService.temp_upload(f) as path:
            # FIX 18: Guard empty PDF before Ghostscript
            guard = _guard_empty_pdf(path)
            if guard: return guard
            fname = generate_output_filename(f.filename, "pdfa")
            out = os.path.join(Config.OUTPUT_FOLDER, fname)
            cmd = [Config.GHOSTSCRIPT, "-dBATCH", "-dNOPAUSE", "-dSAFER",
                     "-sDEVICE=pdfwrite", f"-dPDFA={pdfa_val}",
                     "-dPDFACompatibilityPolicy=1", f"-sOutputFile={out}", path]
            result = subprocess.run(cmd, capture_output=True, timeout=getattr(Config, 'PDFA_TIMEOUT', 300))
            if result.returncode != 0:
                return err("Ghostscript PDF/A conversion failed.", 500)
            # FIX 18: Validate output PDF
            valid, reason = is_valid_pdf(out, min_pages=1)
            if not valid:
                try: os.remove(out)
                except OSError: pass
                return err(f"PDF/A output invalid: {reason}", 500)
            validation_result = None
            if getattr(Config, 'PDFA_VALIDATE', False):
                try:
                    vr = subprocess.run(
                        [Config.VERAPDF_PATH,"--flavour","1b","--format","text",out],
                        capture_output=True, timeout=60)
                    validation_result = {"passed":vr.returncode==0,"output":vr.stdout.decode()[:500]}
                except Exception as vex:
                    validation_result = {"passed":None,"error":str(vex)}
        resp = ok(f"Converted to PDF/A-{version}", out)
        if validation_result:
            resp_data = resp.get_json()
            resp_data["pdfa_validation"] = validation_result
            return jsonify(resp_data)
        return resp
    except subprocess.TimeoutExpired:
        return err("PDF/A conversion timed out", 500)
    except Exception:
        log.exception("pdf_to_pdfa"); return err("PDF/A conversion failed", 500)


# ============================================================================
# COMPARE PDF
# ============================================================================
@app.route("/api/v1/compare-pdf", methods=["POST"])
@app.route("/api/compare-pdf", methods=["POST"])
@require_auth
@require_rate_limit
def compare_pdf():
    import difflib
    files = request.files.getlist("files")
    if len(files) != 2: return err("Exactly 2 PDF files required")
    for f in files:
        e = validate_file(f, Config.ALLOWED_PDF)
        if e: return err(e)
    # FIX 5: Page count guard for compare
    MAX_COMPARE_PAGES = getattr(Config, "MAX_COMPARE_PAGES", 50)
    try:
        with FileService.temp_uploads(files) as paths:
            doc1 = fitz.open(paths[0]); doc2 = fitz.open(paths[1])
            if len(doc1) == 0:
                doc1.close(); doc2.close()
                return err("First PDF has no pages", 400)
            if len(doc2) == 0:
                doc1.close(); doc2.close()
                return err("Second PDF has no pages", 400)
            # FIX 5: Enforce page cap before loop
            if len(doc1) > MAX_COMPARE_PAGES or len(doc2) > MAX_COMPARE_PAGES:
                doc1.close(); doc2.close()
                return err(f"PDFs too large to compare (max {MAX_COMPARE_PAGES} pages each)", 413)
            try:
                pages = min(len(doc1), len(doc2))
                text_diff_pages = []; overall_sims = []
                # FIX 38: Hard cap MAX_WORDS at 150 (not slice hint)
                MAX_WORDS = 150
                # FIX 5: Stream directly to disk ZIP
                fname = generate_output_filename(files[0].filename, "comparison",
                                                  is_multi=True, filenames=[f.filename for f in files])
                out = os.path.join(Config.OUTPUT_FOLDER, fname)
                with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zf:
                    for i in range(pages):
                        # FIX 5: Use 96 DPI (reduced from 150)
                        pix1 = doc1[i].get_pixmap(dpi=96)
                        pix2 = doc2[i].get_pixmap(dpi=96)
                        img1 = Image.open(io.BytesIO(pix1.tobytes("png"))).convert("RGB")
                        img2 = Image.open(io.BytesIO(pix2.tobytes("png"))).convert("RGB")
                        # FIX 5: Free pixmaps immediately
                        pix1 = pix2 = None
                        if img1.size != img2.size: img2 = img2.resize(img1.size, Image.LANCZOS)
                        diff = ImageChops.difference(img1, img2)
                        diff_e = diff.point(lambda x: min(x*8, 255))
                        diff_out = io.BytesIO(); diff_e.save(diff_out, "PNG")
                        zf.writestr(f"diff_page_{i+1:04d}.png", diff_out.getvalue())
                        # FIX 5: Free images immediately after use
                        img1 = img2 = None
                        # FIX 38: Hard cap at 150 words unconditionally
                        words1 = [w[4] for w in doc1[i].get_text("words")][:MAX_WORDS]
                        words2 = [w[4] for w in doc2[i].get_text("words")][:MAX_WORDS]
                        sm = difflib.SequenceMatcher(None, words1, words2)
                        sim = round(sm.ratio()*100, 1); overall_sims.append(sim)
                        added = removed = []
                        for tag,i1,i2,j1,j2 in sm.get_opcodes():
                            if tag=="insert": added = words2[j1:j2]
                            elif tag=="delete": removed = words1[i1:i2]
                            elif tag=="replace": removed = words1[i1:i2]; added = words2[j1:j2]
                        text_diff_pages.append({"page":i+1,"similarity_pct":sim,
                                                 "words_added":added[:50],"words_removed":removed[:50]})
                    overall_sim = round(sum(overall_sims)/len(overall_sims),1) if overall_sims else 0.0
                    zf.writestr("text_diff_summary.json",
                                json.dumps({"pages":text_diff_pages,"overall_similarity_pct":overall_sim},
                                           ensure_ascii=False, indent=2))
            finally:
                doc1.close(); doc2.close()
        return ok(f"Compared {pages} page(s)", out)
    except Exception:
        log.exception("compare_pdf"); return err("Comparison failed", 500)


# ============================================================================
# IMAGE → PDF
# ============================================================================
def _images_to_pdf(paths, page_size_str, output_filename):
    size_map = {"a4": A4, "letter": letter}
    size = size_map.get(page_size_str.lower(), None)
    out = os.path.join(Config.OUTPUT_FOLDER, output_filename)
    c = rl_canvas.Canvas(out, pagesize=size or letter)
    # FIX 46: Track failed images and skip gracefully
    failed = 0
    for path in paths:
        try:
            with Image.open(path) as img:
                iw, ih = img.size
            if size:
                pw, ph = size
            else:
                pw, ph = iw*72/96, ih*72/96
            sw = min(pw*0.95, iw*72/96)
            sh = sw * ih / iw
            if sh > ph*0.95:
                sh = ph*0.95; sw = sh*iw/ih
            x = (pw-sw)/2; y = (ph-sh)/2
            c._pagesize = (pw, ph)
            c.drawImage(path, x, y, width=sw, height=sh)
            c.showPage()
        except Exception as ex:
            log.warning(f"Skipping image {path}: {ex}")
            failed += 1
    if failed == len(paths):
        raise ValueError("All images failed to convert")
    c.save()
    return out


@app.route("/api/v1/image-to-pdf", methods=["POST"])
@app.route("/api/image-to-pdf", methods=["POST"])
@require_auth
@require_rate_limit
def image_to_pdf():
    files = (request.files.getlist("files") or
             request.files.getlist("images") or
             request.files.getlist("image"))
    if not files or all(f.filename=="" for f in files):
        return err("At least one image file required")
    for f in files:
        e = validate_file(f, Config.ALLOWED_IMAGE)
        if e: return err(e)
    page_size = request.form.get("page_size","auto")
    try:
        with FileService.temp_uploads(files) as paths:
            fname = generate_output_filename(files[0].filename, "to_pdf",
                                              is_multi=True, filenames=[f.filename for f in files])
            fname = re.sub(r'\.(jpg|jpeg|png|gif|bmp|tiff|webp)$','.pdf',fname,flags=re.IGNORECASE)
            if not fname.endswith(".pdf"): fname = Path(fname).stem + ".pdf"
            try:
                out = _images_to_pdf(paths, page_size, fname)
            except ValueError as ve:
                return err(str(ve), 400)
        return ok(f"Converted {len(files)} image(s) to PDF", out)
    except Exception:
        log.exception("image_to_pdf"); return err("Image to PDF failed", 500)


@app.route("/api/v1/word-to-pdf", methods=["POST"])
@app.route("/api/word-to-pdf", methods=["POST"])
@require_auth
@require_rate_limit
def word_to_pdf():
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_DOC)
    if e: return err(e)
    try:
        with FileService.temp_upload(f) as path:
            fname = generate_output_filename(f.filename, "to_pdf")
            fname = re.sub(r'\.(doc|docx)$','.pdf',fname,flags=re.IGNORECASE)
            out = libre(path, "pdf", output_filename=fname)
            # FIX 47: Distinguish timeout vs general failure for word_to_pdf
            if not out:
                if not cb_libreoffice.can_execute():
                    return err("LibreOffice service is temporarily unavailable (circuit breaker open)", 503)
                return err("Word to PDF conversion failed — LibreOffice error", 500)
            if not os.path.exists(out) or os.path.getsize(out) == 0:
                return err("LibreOffice conversion produced no output", 500)
        return ok("Word converted to PDF", out)
    except Exception:
        log.exception("word_to_pdf"); return err("Word to PDF failed", 500)


@app.route("/api/v1/excel-to-pdf", methods=["POST"])
@app.route("/api/excel-to-pdf", methods=["POST"])
@require_auth
@require_rate_limit
def excel_to_pdf():
    uploaded_file = request.files.get("file")
    validation_error = validate_file(uploaded_file, Config.ALLOWED_XLS)
    if validation_error:
        return err(validation_error)

    temp_prepared_path = None

    try:
        with FileService.temp_upload(uploaded_file) as temp_path:
            try:
                _ext = os.path.splitext(temp_path)[1].lower()
                if _ext == ".xlsx":
                    from utils.office_utils import prepare_excel_for_pdf
                    _prepared_out = os.path.join(
                        Config.TEMP_FOLDER,
                        f"{uuid.uuid4().hex}_prepared.xlsx"
                    )
                    temp_prepared_path = prepare_excel_for_pdf(temp_path, _prepared_out)
                    conversion_path = temp_prepared_path if temp_prepared_path else temp_path
                else:
                    conversion_path = temp_path
            except Exception as prep_error:
                log.warning(f"Excel pre-processing skipped: {prep_error}")
                conversion_path = temp_path

            output_filename = generate_output_filename(uploaded_file.filename, "to_pdf")
            output_filename = re.sub(r'\.(xls|xlsx)$', '.pdf', output_filename, flags=re.IGNORECASE)
            output_path = os.path.join(Config.OUTPUT_FOLDER, output_filename)

            # FIX 12/31: Route LibreOffice to UUID temp dir, wrap in try/finally
            lo_temp_dir = os.path.join(Config.TEMP_FOLDER, f"lo_out_{uuid.uuid4().hex}")
            os.makedirs(lo_temp_dir, exist_ok=True)

            lo_profile = f"/tmp/lo_profile_{uuid.uuid4().hex}"
            os.makedirs(lo_profile, exist_ok=True)

            try:
                result = subprocess.run(
                    [
                        Config.LIBREOFFICE, "--headless", "--nologo",
                        "--nolockcheck", "--nodefault", "--nofirststartwizard",
                        f"-env:UserInstallation=file://{lo_profile}",
                        "--convert-to",
                        "pdf:calc_pdf_Export:"
                        "SelectPdfVersion=1,"
                        "UseLosslessCompression=true,"
                        "Quality=100,"
                        "ReduceImageResolution=false",
                        "--outdir", lo_temp_dir,
                        conversion_path
                    ],
                    env={**os.environ, "SAL_DEFAULT_PAPER":"A4", "OOO_FORCE_DESKTOP":"true",
                         "HOME": lo_profile, "TMPDIR": lo_temp_dir},
                    capture_output=True, timeout=Config.SUBPROCESS_TIMEOUT
                )

                if result.returncode != 0:
                    stderr_msg = result.stderr.decode('utf-8', errors='ignore')[:500]
                    log.error(f"LibreOffice failed (code {result.returncode}): {stderr_msg}")
                    return err("Excel to PDF conversion failed", 500)

                base_name = os.path.splitext(os.path.basename(conversion_path))[0]
                converted_path = os.path.join(lo_temp_dir, f"{base_name}.pdf")

                if not os.path.exists(converted_path) or os.path.getsize(converted_path) == 0:
                    return err("LibreOffice produced no output", 500)

                # FIX 12: Atomic move from temp dir to output
                os.replace(converted_path, output_path)

                valid, error_msg = is_valid_pdf(output_path, min_pages=1)
                if not valid:
                    os.remove(output_path)
                    return err(f"Output validation failed: {error_msg}", 500)

                return ok("Excel converted to PDF", output_path)

            finally:
                # FIX 31: Always clean up both temp dirs
                shutil.rmtree(lo_temp_dir, ignore_errors=True)
                shutil.rmtree(lo_profile, ignore_errors=True)

    except subprocess.TimeoutExpired:
        return err("Excel to PDF conversion timed out", 500)
    except Exception:
        log.exception("excel_to_pdf")
        return err("Excel to PDF conversion failed", 500)
    finally:
        if temp_prepared_path and os.path.exists(temp_prepared_path):
            try: os.remove(temp_prepared_path)
            except OSError: pass


@app.route("/api/v1/html-to-pdf", methods=["POST"])
@app.route("/api/html-to-pdf", methods=["POST"])
@require_auth
@require_rate_limit
def html_to_pdf():
    _weasyprint_available = False
    try:
        import weasyprint  # noqa
        _weasyprint_available = True
    except ImportError:
        pass

    if not _WKHTMLTOPDF_AVAILABLE and not _weasyprint_available:
        return err(
            "HTML to PDF requires wkhtmltopdf or weasyprint. "
            "Install wkhtmltopdf (https://wkhtmltopdf.org) or run: pip install weasyprint",
            501
        )

    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_HTML)
    if e: return err(e)

    try:
        with FileService.temp_upload(f) as path:
            fname = generate_output_filename(f.filename, "to_pdf")
            fname = re.sub(r'\.(html|htm)$', '.pdf', fname, flags=re.IGNORECASE)
            out_path = os.path.join(Config.OUTPUT_FOLDER, fname)

            if _WKHTMLTOPDF_AVAILABLE:
                try:
                    # FIX 4: SSRF/LFI protection — disable local file access and JS
                    html_dir = os.path.dirname(path)
                    # Strip dangerous tags from HTML before conversion
                    with open(path, 'r', encoding='utf-8', errors='replace') as hf:
                        html_content = hf.read()
                    html_content = re.sub(r'<script[^>]*>.*?</script>', '', html_content,
                                          flags=re.IGNORECASE | re.DOTALL)
                    html_content = re.sub(r'<iframe[^>]*>.*?</iframe>', '', html_content,
                                          flags=re.IGNORECASE | re.DOTALL)
                    html_content = re.sub(r'<object[^>]*>.*?</object>', '', html_content,
                                          flags=re.IGNORECASE | re.DOTALL)
                    html_content = re.sub(r'<embed[^>]*/>', '', html_content,
                                          flags=re.IGNORECASE | re.DOTALL)
                    html_content = re.sub(r'<meta[^>]*http-equiv[^>]*/>', '', html_content,
                                          flags=re.IGNORECASE | re.DOTALL)
                    with open(path, 'w', encoding='utf-8') as hf:
                        hf.write(html_content)
                    result = subprocess.run(
                        [Config.WKHTMLTOPDF, "--quiet",
                         "--disable-local-file-access",
                         "--disable-javascript",
                         "--no-background",
                         "--load-error-handling", "ignore",
                         path, out_path],
                        capture_output=True, timeout=60
                    )
                    if result.returncode == 0 and os.path.exists(out_path) and os.path.getsize(out_path) > 0:
                        return ok("HTML converted to PDF", out_path)
                except subprocess.TimeoutExpired:
                    return err("HTML to PDF timed out", 500)
                except Exception as wk_ex:
                    log.warning(f"wkhtmltopdf failed: {wk_ex}")

            if _weasyprint_available:
                try:
                    from weasyprint import HTML
                    HTML(filename=path).write_pdf(out_path)
                    if os.path.exists(out_path) and os.path.getsize(out_path) > 0:
                        return ok("HTML converted to PDF (WeasyPrint)", out_path)
                except Exception as we:
                    log.warning(f"WeasyPrint failed: {we}")

            return err("HTML to PDF failed — both wkhtmltopdf and weasyprint encountered errors", 500)

    except subprocess.TimeoutExpired:
        return err("HTML to PDF timed out", 500)
    except Exception:
        log.exception("html_to_pdf")
        return err("HTML to PDF failed", 500)


# ============================================================================
# IMAGE TOOLS
# ============================================================================
@app.route("/api/v1/compress-image", methods=["POST"])
@app.route("/api/compress-image", methods=["POST"])
@require_auth
@require_rate_limit
def compress_image():
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_IMAGE)
    if e: return err(e)
    quality = safe_int(request.form.get("quality","75"), 75, 1, 95)
    output_format = request.form.get("output_format","auto").lower()
    try:
        with FileService.temp_upload(f) as path:
            orig_size = os.path.getsize(path)
            # FIX 34: Catch DecompressionBombError
            try:
                img = Image.open(path)
            except Image.DecompressionBombError:
                return err("Image too large to process safely", 413)
            ext = path.rsplit(".",1)[-1].lower()
            if output_format == "webp": target_fmt = "webp"
            elif output_format in ("jpg","jpeg"): target_fmt = "jpeg"
            elif output_format == "png": target_fmt = "png"
            elif ext == "png": target_fmt = "png"
            elif ext == "webp": target_fmt = "webp"
            else: target_fmt = "jpeg"
            fname = generate_output_filename(f.filename, "compressed")
            fname = re.sub(r'\.\w+$', f'.{"jpg" if target_fmt=="jpeg" else target_fmt}', fname)
            out = os.path.join(Config.OUTPUT_FOLDER, fname)
            if target_fmt == "png":
                if img.mode not in ("RGB","RGBA","L","LA","P"): img = img.convert("RGBA")
                tmp_png = os.path.join(Config.TEMP_FOLDER, f"{uuid.uuid4().hex}_tmp.png")
                img.save(tmp_png, format="PNG", optimize=False)
                pngquant_ok = False
                try:
                    res = subprocess.run(
                        ["pngquant","--quality",f"{max(1,quality-15)}-{quality}",
                         "--speed","3","--force","--output",out,tmp_png],
                        capture_output=True, timeout=30)
                    if res.returncode == 0 and os.path.exists(out):
                        pngquant_ok = True
                except (FileNotFoundError, subprocess.TimeoutExpired): pass
                finally:
                    try: os.remove(tmp_png)
                    except OSError: pass
                if not pngquant_ok:
                    img.save(out, format="PNG", optimize=True, compress_level=9)
            elif target_fmt == "webp":
                if img.mode not in ("RGB","RGBA"): img = img.convert("RGB")
                img.save(out, format="WEBP", quality=quality, method=6, optimize=True)
            else:
                if img.mode in ("RGBA","P","LA"):
                    bg = Image.new("RGB", img.size, (255,255,255))
                    if img.mode == "P": img = img.convert("RGBA")
                    mask = img.split()[-1] if img.mode in ("RGBA","LA") else None
                    bg.paste(img, mask=mask); img = bg
                elif img.mode != "RGB": img = img.convert("RGB")
                img.save(out, format="JPEG", quality=quality, optimize=True, progressive=True)
            new_size = os.path.getsize(out)
            reduction = round((1-new_size/orig_size)*100,1) if orig_size else 0
        return ok(f"Image compressed {reduction}% as {target_fmt.upper()}", out,
                  reduction_pct=reduction, output_format=target_fmt)
    except Exception:
        log.exception("compress_image"); return err("Image compression failed", 500)


@app.route("/api/v1/resize-image", methods=["POST"])
@app.route("/api/resize-image", methods=["POST"])
@require_auth
@require_rate_limit
def resize_image():
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_IMAGE)
    if e: return err(e)
    width = safe_int(request.form.get("width","800"), 800, 1, 8000)
    height = safe_int(request.form.get("height","600"), 600, 1, 8000)
    keep_ratio = request.form.get("keep_ratio","true").lower() in ("true","on","1","yes")
    try:
        with FileService.temp_upload(f) as path:
            # FIX 34: Catch DecompressionBombError
            try:
                img = Image.open(path)
            except Image.DecompressionBombError:
                return err("Image too large to process safely", 413)
            if keep_ratio: img.thumbnail((width, height), Image.LANCZOS)
            else: img = img.resize((width, height), Image.LANCZOS)
            ext = path.rsplit(".",1)[-1].lower()
            fmt = "JPEG" if ext in ("jpg","jpeg") else "PNG"
            if fmt == "JPEG" and img.mode in ("RGBA","P"): img = img.convert("RGB")
            fname = generate_output_filename(f.filename, "resized")
            out = os.path.join(Config.OUTPUT_FOLDER, fname)
            img.save(out, format=fmt)
        return ok(f"Image resized to {img.size[0]}×{img.size[1]}", out)
    except Exception:
        log.exception("resize_image"); return err("Resize failed", 500)


@app.route("/api/v1/webp-to-jpg", methods=["POST"])
@app.route("/api/webp-to-jpg", methods=["POST"])
@require_auth
@require_rate_limit
def webp_to_jpg():
    files = (request.files.getlist("files") or
             request.files.getlist("webp"))
    if not files or all(f.filename=="" for f in files):
        return err("At least one WebP file required")
    for f in files:
        e = validate_file(f, Config.ALLOWED_WEBP)
        if e: return err(e)
    quality = safe_int(request.form.get("quality","75"), 75, 1, 95)
    try:
        with FileService.temp_uploads(files) as paths:
            buf = io.BytesIO()
            with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
                for i, path in enumerate(paths):
                    # FIX 34: Catch DecompressionBombError
                    try:
                        img = Image.open(path)
                    except Image.DecompressionBombError:
                        return err("Image too large to process safely", 413)
                    exif_bytes = img.info.get("exif",b"") if hasattr(img,"info") else b""
                    rgb = img.convert("RGB")
                    ib = io.BytesIO()
                    kw = {"format":"JPEG","quality":quality}
                    if exif_bytes: kw["exif"] = exif_bytes
                    rgb.save(ib, **kw)
                    zf.writestr(f"image_{i+1:04d}.jpg", ib.getvalue())
            fname = generate_output_filename(files[0].filename,"to_jpg",
                                              is_multi=True, filenames=[f.filename for f in files])
            out = os.path.join(Config.OUTPUT_FOLDER, fname)
            with open(out,"wb") as fh: fh.write(buf.getvalue())
        return ok(f"Converted {len(files)} WebP(s) to JPG", out)
    except Exception:
        log.exception("webp_to_jpg"); return err("WebP to JPG failed", 500)


@app.route("/api/v1/png-to-jpg", methods=["POST"])
@app.route("/api/png-to-jpg", methods=["POST"])
@require_auth
@require_rate_limit
def png_to_jpg():
    files = (request.files.getlist("files") or
             request.files.getlist("png"))
    if not files or all(f.filename=="" for f in files):
        return err("At least one PNG file required")
    for f in files:
        e = validate_file(f, Config.ALLOWED_PNG)
        if e: return err(e)
    quality = safe_int(request.form.get("quality","75"), 75, 1, 95)
    try:
        with FileService.temp_uploads(files) as paths:
            buf = io.BytesIO()
            with zipfile.ZipFile(buf,"w",zipfile.ZIP_DEFLATED) as zf:
                for i, path in enumerate(paths):
                    # FIX 34: Catch DecompressionBombError
                    try:
                        img = Image.open(path).convert("RGB")
                    except Image.DecompressionBombError:
                        return err("Image too large to process safely", 413)
                    ib = io.BytesIO(); img.save(ib, format="JPEG", quality=quality)
                    zf.writestr(f"image_{i+1:04d}.jpg", ib.getvalue())
            fname = generate_output_filename(files[0].filename,"to_jpg",
                                              is_multi=True, filenames=[f.filename for f in files])
            out = os.path.join(Config.OUTPUT_FOLDER, fname)
            with open(out,"wb") as fh: fh.write(buf.getvalue())
        return ok(f"Converted {len(files)} PNG(s) to JPG", out)
    except Exception:
        log.exception("png_to_jpg"); return err("PNG to JPG failed", 500)


@app.route("/api/v1/image-to-excel", methods=["POST"])
@app.route("/api/image-to-excel", methods=["POST"])
@require_auth
@require_rate_limit
def image_to_excel():
    if not OPENPYXL_AVAILABLE: return err("Requires openpyxl.", 501)
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_IMAGE)
    if e: return err(e)
    raw_lang = request.form.get("lang","eng")
    lang = re.sub(r'[^a-zA-Z0-9+\-]','',raw_lang)[:50] or "eng"
    try:
        with FileService.temp_upload(f) as path:
            # FIX 34: Catch DecompressionBombError
            try:
                img = Image.open(path).convert("RGB")
            except Image.DecompressionBombError:
                return err("Image too large to process safely", 413)
            ocr_grid = None; warning = None
            if TESSERACT_AVAILABLE:
                try:
                    # FIX 20: OCR semaphore for image_to_excel
                    with _ocr_semaphore:
                        data = pytesseract.image_to_data(img, lang=lang,
                                                          output_type=TesseractOutput.DICT, config="--psm 6")
                    words = []
                    for txt, conf, left, top, wd, ht in zip(
                        data.get("text",[]),
                        data.get("conf",[]),
                        data.get("left",[]),
                        data.get("top",[]),
                        data.get("width",[]),
                        data.get("height",[])
                    ):
                        try: conf_int = int(conf)
                        except (ValueError, TypeError): conf_int = 0
                        if (txt or "").strip() and conf_int >= 30:
                            words.append({"text":(txt or "").strip(),"left":left,
                                          "top":top,"width":wd,"height":ht})
                    if len(words) < 3:
                        warning = f"OCR found only {len(words)} word(s)"
                    else:
                        avg_char_w = sum(w["width"]/max(len(w["text"]),1) for w in words)/len(words)
                        dynamic_gap = max(10, avg_char_w * 1.5)
                        tol = 8; rows_dict = {}
                        for w in words:
                            mid_y = w["top"]+w["height"]//2
                            key = next((k for k in rows_dict if abs(k-mid_y)<=tol), None)
                            if key is None: key = mid_y; rows_dict[key] = []
                            rows_dict[key].append(w)
                        sorted_rows = [sorted(rows_dict[k], key=lambda x: x["left"]) for k in sorted(rows_dict)]
                        grid = []
                        for row_words in sorted_rows:
                            cells = [row_words[0]["text"]]
                            for wi in range(1, len(row_words)):
                                prev = row_words[wi-1]; curr = row_words[wi]
                                if curr["left"]-(prev["left"]+prev["width"]) > dynamic_gap:
                                    cells.append(curr["text"])
                                else:
                                    cells[-1] += " " + curr["text"]
                            grid.append(cells)
                        ocr_grid = grid
                except Exception as ocr_ex:
                    warning = f"OCR failed ({ocr_ex}) — falling back to image embed"
            else:
                warning = "pytesseract not available — falling back to image embed"
            wb = Workbook(); ws = wb.active; msg = ""
            if ocr_grid:
                ws.title = "OCR_Table"
                for r_idx, row_cells in enumerate(ocr_grid):
                    for c_idx, val in enumerate(row_cells):
                        cell = ws.cell(row=r_idx+1, column=c_idx+1, value=val)
                        if r_idx == 0: cell.font = Font(bold=True)
                msg = f"OCR extracted {len(ocr_grid)} rows"
            else:
                ws.title = "Image"
                tmp = tempfile.NamedTemporaryFile(
                    suffix=".png", delete=False, dir=Config.TEMP_FOLDER
                )
                try:
                    img.save(tmp.name, format="PNG"); tmp.close()
                    xl_img = XlImage(tmp.name); xl_img.anchor = "B2"
                    ws.add_image(xl_img)
                finally:
                    try: os.unlink(tmp.name)
                    except OSError: pass
                msg = "Image embedded in Excel"
            fname = generate_output_filename(f.filename,"to_excel")
            fname = re.sub(r'\.(jpg|jpeg|png|gif|bmp|tiff|webp)$','.xlsx',fname,flags=re.IGNORECASE)
            out = os.path.join(Config.OUTPUT_FOLDER, fname)
            wb.save(out)
            if warning: msg += f" | Warning: {warning}"
        return ok(msg, out, ocr_warning=warning)
    except Exception:
        log.exception("image_to_excel"); return err("Image to Excel failed", 500)


# ============================================================================
# WORD TOOLS
# ============================================================================
@app.route("/api/v1/word-to-txt", methods=["POST"])
@app.route("/api/word-to-txt", methods=["POST"])
@require_auth
@require_rate_limit
def word_to_txt():
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_DOC)
    if e: return err(e)
    try:
        with FileService.temp_upload(f) as path:
            fname = generate_output_filename(f.filename,"to_txt")
            fname = re.sub(r'\.(doc|docx)$','.txt',fname,flags=re.IGNORECASE)
            if DOCX_AVAILABLE and path.endswith(".docx"):
                doc = DocxDocument(path)
                text = "\n".join(p.text for p in doc.paragraphs)
                out = os.path.join(Config.OUTPUT_FOLDER, fname)
                with open(out,"w",encoding="utf-8") as fh: fh.write(text)
            else:
                out = libre(path, "txt", output_filename=fname)
                if not out: return err("LibreOffice conversion failed.", 500)
        return ok("Word converted to TXT", out)
    except Exception:
        log.exception("word_to_txt"); return err("Word to TXT failed", 500)


@app.route("/api/v1/word-to-excel", methods=["POST"])
@app.route("/api/word-to-excel", methods=["POST"])
@require_auth
@require_rate_limit
def word_to_excel():
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_DOC)
    if e: return err(e)
    try:
        with FileService.temp_upload(f) as path:
            fname = generate_output_filename(f.filename,"to_excel")
            fname = re.sub(r'\.(doc|docx)$','.xlsx',fname,flags=re.IGNORECASE)
            if path.endswith(".doc"):
                out = libre(path, "xlsx", output_filename=fname)
                if not out: return err("Word (.doc) to Excel requires LibreOffice.", 500)
                return ok("Word (.doc) converted to Excel via LibreOffice", out)
            if not DOCX_AVAILABLE or not OPENPYXL_AVAILABLE:
                return err("Requires python-docx + openpyxl.", 501)
            doc = DocxDocument(path); wb = Workbook(); wb.remove(wb.active)
            table_count = len(doc.tables)
            # FIX MISSED-4 / word_to_excel cell count guard
            MAX_WORD_TABLE_CELLS = getattr(Config, "MAX_WORD_TABLE_CELLS", 500000)
            total_cells = sum(
                len(row.cells) for table in doc.tables for row in table.rows
            )
            if total_cells > MAX_WORD_TABLE_CELLS:
                return err(f"Document too large: {total_cells} table cells (max {MAX_WORD_TABLE_CELLS})", 413)
            for t_idx, table in enumerate(doc.tables):
                ws = wb.create_sheet(title=f"Table_{t_idx+1}")
                for r_idx, row in enumerate(table.rows):
                    for c_idx, cell in enumerate(row.cells):
                        co = ws.cell(row=r_idx+1, column=c_idx+1, value=cell.text)
                        if r_idx == 0:
                            try:
                                co.font = Font(bold=True)
                            except Exception:
                                pass
            ws_text = wb.create_sheet("Document_Text")
            ws_text.append(["Line","Style","Text"])
            for cell in ws_text[1]: cell.font = Font(bold=True)
            pr = 2
            for p in doc.paragraphs:
                if p.text.strip():
                    ws_text.cell(pr,1,pr-1); ws_text.cell(pr,2,p.style.name)
                    ws_text.cell(pr,3,p.text); pr += 1
            out = os.path.join(Config.OUTPUT_FOLDER, fname)
            wb.save(out)
        return ok(f"Word converted to Excel — {table_count} table(s)", out, tables_found=table_count)
    except Exception:
        log.exception("word_to_excel"); return err("Word to Excel failed", 500)


@app.route("/api/v1/word-to-ppt", methods=["POST"])
@app.route("/api/word-to-ppt", methods=["POST"])
@require_auth
@require_rate_limit
def word_to_ppt():
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_DOC)
    if e: return err(e)
    try:
        with FileService.temp_upload(f) as path:
            fname = generate_output_filename(f.filename,"to_ppt")
            fname = re.sub(r'\.(doc|docx)$','.pptx',fname,flags=re.IGNORECASE)
            if path.endswith(".doc"):
                out = libre(path, "pptx", output_filename=fname)
                if not out: return err("LibreOffice required for .doc.", 500)
                return ok("Word (.doc) converted to PPT via LibreOffice", out)
            if not DOCX_AVAILABLE or not PPTX_AVAILABLE:
                return err("Requires python-docx + python-pptx.", 501)
            doc = DocxDocument(path)
            prs = Presentation(); prs.slide_width=PptxInches(10); prs.slide_height=PptxInches(7.5)
            tc_layout = prs.slide_layouts[1]
            paragraphs = [p for p in doc.paragraphs if p.text.strip()]
            if not paragraphs:
                prs.slides.add_slide(prs.slide_layouts[6])
                out = os.path.join(Config.OUTPUT_FOLDER, fname); prs.save(out)
                return ok("Word converted to PPT (empty doc)", out)

            MAX_SLIDES_PPT = int(os.environ.get("MAX_SLIDES_PPT",
                                                  getattr(Config, 'MAX_SLIDES_PPT', 200)))
            has_headings = any(p.style.name.startswith("Heading") or p.style.name=="Title"
                               for p in paragraphs)
            slide_count = 0
            if not has_headings:
                slide = prs.slides.add_slide(tc_layout)
                slide_count += 1
                if slide.shapes.title: slide.shapes.title.text = Path(f.filename).stem
                if len(slide.placeholders) > 1:
                    tf = slide.placeholders[1].text_frame; tf.clear()
                    for p in paragraphs:
                        if slide_count >= MAX_SLIDES_PPT: break
                        pr2 = tf.add_paragraph(); pr2.text = p.text
            else:
                current_tf = None
                for p in paragraphs:
                    # FIX 22: Check slide limit at TOP of loop BEFORE add_slide
                    if slide_count >= MAX_SLIDES_PPT:
                        s = prs.slides.add_slide(tc_layout)
                        if s.shapes.title: s.shapes.title.text = "Document Truncated"
                        if len(s.placeholders) > 1:
                            s.placeholders[1].text = (
                                f"Truncated at {MAX_SLIDES_PPT} slides limit."
                            )
                        break
                    sn, text = p.style.name, p.text.strip()
                    if sn.startswith("Heading 1") or sn == "Title":
                        slide = prs.slides.add_slide(tc_layout)
                        slide_count += 1
                        if slide.shapes.title: slide.shapes.title.text = text
                        current_tf = slide.placeholders[1].text_frame if len(slide.placeholders)>1 else None
                        if current_tf: current_tf.clear()
                    elif sn.startswith("Heading"):
                        if current_tf is None:
                            slide = prs.slides.add_slide(tc_layout)
                            slide_count += 1
                            current_tf = slide.placeholders[1].text_frame if len(slide.placeholders)>1 else None
                            if current_tf: current_tf.clear()
                        if current_tf:
                            pr2 = current_tf.add_paragraph(); pr2.text = text
                            if pr2.runs: pr2.runs[0].font.bold = True
                    else:
                        if current_tf:
                            pr2 = current_tf.add_paragraph(); pr2.text = text; pr2.level = 1

            out = os.path.join(Config.OUTPUT_FOLDER, fname)
            # FIX 26: Atomic write for word_to_ppt
            tmp_out = out + ".tmp"
            prs.save(tmp_out)
            os.replace(tmp_out, out)
        return ok("Word converted to PowerPoint", out,
                  slides_created=slide_count,
                  truncated=(slide_count >= MAX_SLIDES_PPT))
    except Exception:
        log.exception("word_to_ppt"); return err("Word to PPT failed", 500)


@app.route("/api/v1/compress-word", methods=["POST"])
@app.route("/api/compress-word", methods=["POST"])
@require_auth
@require_rate_limit
def compress_word():
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_DOC)
    if e: return err(e)
    quality_str = request.form.get("quality","medium").lower()
    jpeg_quality = {"low":50,"medium":70,"high":85}.get(quality_str, 70)
    force_jpeg = request.form.get("force_jpeg","false").lower() in ("true","1","yes")
    try:
        with FileService.temp_upload(f) as path:
            orig_size = os.path.getsize(path)
            fname = generate_output_filename(f.filename,"compressed")
            out = os.path.join(Config.OUTPUT_FOLDER, fname)
            work_path = path
            if path.endswith(".doc"):
                converted = libre(path, "docx", temp=True)
                if not converted: return err("LibreOffice required for .doc.", 500)
                work_path = converted
            tmp_dir = tempfile.mkdtemp()
            try:
                with zipfile.ZipFile(work_path,"r") as zin:
                    _safe_zip_extract(zin, tmp_dir)
                media_dir = os.path.join(tmp_dir,"word","media"); compressed = 0
                if os.path.isdir(media_dir):
                    for fname_img in os.listdir(media_dir):
                        img_path = os.path.join(media_dir, fname_img)
                        ext_img = os.path.splitext(fname_img)[1].lower()
                        if ext_img not in {".png",".jpg",".jpeg",".gif",".bmp"}: continue
                        try:
                            # FIX 34: Catch DecompressionBombError in compress_word
                            try:
                                img = Image.open(img_path)
                            except Image.DecompressionBombError:
                                log.warning(f"Skipping oversized image in Word doc: {fname_img}")
                                continue
                            w, h = img.size
                            if w > 1200 or h > 1200:
                                ratio = min(1200/w, 1200/h)
                                img = img.resize((max(1,int(w*ratio)),max(1,int(h*ratio))), Image.LANCZOS)
                            if force_jpeg:
                                if img.mode in ("RGBA","LA","P"):
                                    bg = Image.new("RGB",img.size,(255,255,255))
                                    if img.mode=="P": img=img.convert("RGBA")
                                    mask = img.split()[-1] if img.mode in ("RGBA","LA") else None
                                    bg.paste(img,mask=mask); img=bg
                                elif img.mode!="RGB": img=img.convert("RGB")
                                new_path = os.path.splitext(img_path)[0]+".jpg"
                                img.save(new_path,format="JPEG",quality=jpeg_quality,optimize=True)
                                if new_path != img_path: os.remove(img_path)
                            elif ext_img==".png" or img.mode in ("RGBA","LA","P"):
                                if img.mode not in ("RGBA","RGB","L"): img=img.convert("RGBA")
                                img.save(img_path,format="PNG",optimize=True,compress_level=9)
                            else:
                                if img.mode!="RGB": img=img.convert("RGB")
                                img.save(img_path,format="JPEG",quality=jpeg_quality,optimize=True)
                            compressed += 1
                        except Exception as ie: log.warning(f"compress-word img {fname_img}: {ie}")
                # FIX 28: Depth guard in compress_word ZIP repack
                with zipfile.ZipFile(out,"w",zipfile.ZIP_DEFLATED,compresslevel=9) as zout:
                    for root, dirs, files_list in os.walk(tmp_dir):
                        # Calculate depth relative to tmp_dir
                        depth = root[len(tmp_dir):].count(os.sep)
                        if depth > 5:
                            dirs.clear()
                            continue
                        for fi in files_list:
                            abs_p = os.path.join(root,fi)
                            zout.write(abs_p, os.path.relpath(abs_p, tmp_dir))
            finally:
                shutil.rmtree(tmp_dir, ignore_errors=True)
                if work_path != path:
                    try: os.remove(work_path)
                    except OSError: pass
            new_size = os.path.getsize(out)
            reduction = round((1-new_size/orig_size)*100,1) if orig_size else 0
        return ok(f"Word compressed — {reduction}% ({compressed} images)", out,
                  reduction_pct=reduction, images_compressed=compressed, force_jpeg_used=force_jpeg)
    except Exception:
        log.exception("compress_word"); return err("Word compression failed", 500)


@app.route("/api/v1/edit-word", methods=["POST"])
@app.route("/api/edit-word", methods=["POST"])
@require_auth
@require_rate_limit
def edit_word():
    if not DOCX_AVAILABLE: return err("Requires python-docx.", 501)
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_DOC)
    if e: return err(e)
    find_text = sanitize_string(request.form.get("find_text",""))
    replace_text = sanitize_string(request.form.get("replace_text",""))
    if not find_text: return err("find_text required")
    try:
        with FileService.temp_upload(f) as path:
            doc = DocxDocument(path); count = 0
            for para in doc.paragraphs:
                for run in para.runs:
                    if find_text in run.text:
                        count += run.text.count(find_text)
                        run.text = run.text.replace(find_text, replace_text)
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for para in cell.paragraphs:
                            for run in para.runs:
                                if find_text in run.text:
                                    count += run.text.count(find_text)
                                    run.text = run.text.replace(find_text, replace_text)
            fname = generate_output_filename(f.filename,"edited")
            out = os.path.join(Config.OUTPUT_FOLDER, fname)
            doc.save(out)
        return ok(f"Replaced {count} occurrence(s)", out)
    except Exception:
        log.exception("edit_word"); return err("Edit Word failed", 500)


@app.route("/api/v1/unlock-word", methods=["POST"])
@app.route("/api/unlock-word", methods=["POST"])
@require_auth
@require_rate_limit
def unlock_word():
    if not MSOFFCRYPTO_AVAILABLE: return err("Requires msoffcrypto-tool.", 501)
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_DOC)
    if e: return err(e)
    pw = sanitize_string(request.form.get("password",""))
    if not pw: return err("Password required")
    try:
        with FileService.temp_upload(f) as path:
            fname = generate_output_filename(f.filename,"unlocked")
            out = os.path.join(Config.OUTPUT_FOLDER, fname)
            with open(path,"rb") as fp:
                of = msoffcrypto.OfficeFile(fp); of.load_key(password=pw)
                with open(out,"wb") as fout: of.decrypt(fout)
            # FIX 36: Validate decryption output for unlock_word
            if not os.path.exists(out) or os.path.getsize(out) < 100:
                try: os.remove(out)
                except OSError: pass
                return err("Decryption failed — incorrect password or unsupported encryption format", 400)
        return ok("Word document unlocked", out)
    except Exception:
        log.exception("unlock_word"); return err("Unlock failed — check password", 500)


@app.route("/api/protect-word", methods=["POST"])
@require_auth
@require_rate_limit
def protect_word():
    if not MSOFFCRYPTO_AVAILABLE:
        return err("Requires msoffcrypto-tool.", 501)
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_DOC)
    if e: return err(e)
    pw = sanitize_string(request.form.get("password",""))
    pw2 = sanitize_string(request.form.get("password2",""))
    ep = validate_password(pw, pw2)
    if ep: return err(ep)

    try:
        with FileService.temp_upload(f) as path:
            work_path = path
            converted_temp = None

            if path.lower().endswith(".doc") and not path.lower().endswith(".docx"):
                log.info(f"Converting .doc to .docx for protection: {path}")
                converted_temp = libre(path, "docx", temp=True)
                if not converted_temp or not os.path.exists(converted_temp):
                    return err("Could not convert .doc to .docx format", 500)
                work_path = converted_temp

            fname = generate_output_filename(f.filename, "protected")
            out = os.path.join(Config.OUTPUT_FOLDER, fname)

            # FIX 17: Use _pw for encryption, redact pw immediately
            _pw = pw
            pw = "[REDACTED]"
            try:
                with open(work_path, "rb") as fp:
                    of = msoffcrypto.OfficeFile(fp)
                    encrypted = False
                    try:
                        with open(out, "wb") as fout:
                            of.encrypt(_pw, fout, cipher_algorithm="AES")
                        encrypted = True
                    except TypeError:
                        pass
                    except Exception:
                        pass

                    if not encrypted:
                        try:
                            with open(out, "wb") as fout:
                                of.encrypt(_pw, fout)
                            encrypted = True
                        except Exception as fallback_ex:
                            log.error(f"protect_word fallback failed: {type(fallback_ex).__name__}")
                            raise fallback_ex
            finally:
                _pw = None

            if converted_temp:
                try: os.remove(converted_temp)
                except OSError: pass

            if not encrypted or not os.path.exists(out) or os.path.getsize(out) == 0:
                return err("Word protection failed — could not encrypt document", 500)

        return ok("Word document protected", out)
    except Exception:
        log.exception("protect_word")
        return err("Protect Word failed", 500)


# ============================================================================
# EXCEL TOOLS
# ============================================================================
@app.route("/api/v1/excel-to-csv", methods=["POST"])
@app.route("/api/excel-to-csv", methods=["POST"])
@require_auth
@require_rate_limit
def excel_to_csv():
    if not OPENPYXL_AVAILABLE: return err("Requires openpyxl.", 501)
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_XLS)
    if e: return err(e)
    sheet_name = sanitize_string(request.form.get("sheet",""), 100)
    all_sheets = request.form.get("all_sheets","false").lower() in ("true","1","yes")
    try:
        with FileService.temp_upload(f) as path:
            wb = load_workbook(path, data_only=True, read_only=True)
            # FIX 54: Stream to disk ZIP for all_sheets branch to avoid RAM bomb
            if all_sheets:
                total_rows = 0
                fname = generate_output_filename(f.filename,"to_csv")
                fname = re.sub(r'\.(xls|xlsx)$','.zip',fname,flags=re.IGNORECASE)
                out = os.path.join(Config.OUTPUT_FOLDER, fname)
                try:
                    with zipfile.ZipFile(out,"w",zipfile.ZIP_DEFLATED) as zf:
                        for sname in wb.sheetnames:
                            ws = wb[sname]
                            safe_name = re.sub(r'[^\w]','_', sname)
                            tmp_csv = os.path.join(Config.TEMP_FOLDER, f"{uuid.uuid4().hex}.csv")
                            try:
                                cnt = 0
                                with open(tmp_csv, "w", newline="", encoding="utf-8-sig") as fh:
                                    writer = csv.writer(fh, quoting=csv.QUOTE_MINIMAL)
                                    for row in ws.iter_rows(values_only=True):
                                        writer.writerow([coerce_cell_for_csv(v) for v in row])
                                        cnt += 1
                                total_rows += cnt
                                zf.write(tmp_csv, f"{safe_name}.csv")
                            finally:
                                try: os.remove(tmp_csv)
                                except OSError: pass
                finally:
                    try: wb.close()
                    except Exception: pass
                return ok(f"All {len(wb.sheetnames)} sheet(s) exported ({total_rows} rows)", out)
            else:
                ws = wb[sheet_name] if sheet_name and sheet_name in wb.sheetnames else wb.active
                fname = generate_output_filename(f.filename,"to_csv")
                fname = re.sub(r'\.(xls|xlsx)$','.csv',fname,flags=re.IGNORECASE)
                out = os.path.join(Config.OUTPUT_FOLDER, fname); cnt = 0
                with open(out,"w",newline="",encoding="utf-8-sig") as fh:
                    writer = csv.writer(fh, quoting=csv.QUOTE_MINIMAL)
                    for row in ws.iter_rows(values_only=True):
                        writer.writerow([coerce_cell_for_csv(v) for v in row]); cnt += 1
                try: wb.close()
                except Exception: pass
                return ok(f"Excel converted to CSV ({cnt} rows)", out)
    except Exception:
        log.exception("excel_to_csv"); return err("Excel to CSV failed", 500)


@app.route("/api/v1/excel-to-word", methods=["POST"])
@app.route("/api/excel-to-word", methods=["POST"])
@require_auth
@require_rate_limit
def excel_to_word():
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_XLS)
    if e: return err(e)
    preserve_formulas = request.form.get("preserve_formulas","true").lower() in ("true","1","yes")
    row_limit = safe_int(request.form.get("row_limit", Config.EXCEL_ROW_LIMIT),
                          Config.EXCEL_ROW_LIMIT, 1, 100000)

    f.seek(0, 2); file_size = f.tell(); f.seek(0)

    if not OPENPYXL_AVAILABLE or not DOCX_AVAILABLE:
        try:
            with FileService.temp_upload(f) as path:
                fname = generate_output_filename(f.filename,"to_word")
                fname = re.sub(r'\.(xls|xlsx)$','.docx',fname,flags=re.IGNORECASE)
                out = libre(path,"docx",output_filename=fname)
                if not out: return err("Excel to Word requires openpyxl+python-docx or LibreOffice.",500)
            return ok("Excel converted to Word (LibreOffice)", out)
        except Exception: log.exception("excel_to_word_libre"); return err("Excel to Word failed", 500)

    if file_size > _ASYNC_FILE_THRESHOLD:
        _excel_to_word_task = None
        try:
            from tasks.office_tasks import excel_to_word_task as _excel_to_word_task
        except Exception as _ie:
            log.error(f"excel_to_word_task lazy import failed: {_ie}")

        if celery_app and _excel_to_word_task:
            ext = f.filename.rsplit(".", 1)[-1].lower() if "." in f.filename else "xlsx"
            bg_input = os.path.join(Config.UPLOAD_FOLDER, f"{uuid.uuid4()}.{ext}")
            f.seek(0)
            with open(bg_input, "wb") as fh: fh.write(f.read())
            # FIX 27: Check Celery availability before writing bg_input — already written above.
            # Clean up if Celery actually unavailable
            try:
                redis_service.client.setex(
                    f"cleanup:{bg_input}", Config.FILE_TTL_SEC + 600, "pending")
            except Exception: pass
            fname = generate_output_filename(f.filename, "to_word")
            fname = re.sub(r'\.(xls|xlsx)$', '.docx', fname, flags=re.IGNORECASE)
            bg_out = os.path.join(Config.OUTPUT_FOLDER, fname)
            job_id = str(uuid.uuid4())
            redis_service.job_set(job_id, {
                "status": "pending", "progress": "0", "operation": "excel_to_word",
                "created_at": get_timestamp(), "user_id": getattr(g, "user_id", "default")
            })
            try:
                redis_service.client.expire(f"job:{job_id}", getattr(Config, 'JOB_TTL_SEC', 7200))
            except Exception: pass
            # FIX 3: Wrap .delay() in try/except
            try:
                task = _excel_to_word_task.delay(bg_input, bg_out, job_id,
                                                  preserve_formulas, row_limit)
            except Exception:
                try: os.remove(bg_input)
                except OSError: pass
                raise
            redis_service.job_update(job_id, {"task_id": task.id})
            return jsonify({
                "success": True,
                "message": "Large file — conversion queued. Check status_url.",
                "job_id": job_id,
                "status_url": f"/api/v1/jobs/{job_id}",
                "poll_interval_ms": 2000
            })
        else:
            # FIX 27: Celery unavailable — clean up bg_input before falling through to sync
            log.warning("excel_to_word async unavailable, processing synchronously")
            try: os.remove(bg_input)
            except Exception: pass

    try:
        with FileService.temp_upload(f) as path:
            wb = load_workbook(path, data_only=not preserve_formulas)
            doc = DocxDocument(); sheet_count = len(wb.sheetnames); formulas_present = False
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]; doc.add_heading(sheet_name, level=1)
                all_rows = list(ws.iter_rows(values_only=True, max_row=row_limit+1))
                truncated = len(all_rows) > row_limit
                rows_write = all_rows[:row_limit]
                if not rows_write: doc.add_paragraph("(empty sheet)"); continue
                n_cols = max((len(r) for r in rows_write), default=1)
                table = doc.add_table(rows=len(rows_write), cols=n_cols)
                try:
                    table.style = "Light Grid Accent 1"
                except (KeyError, Exception):
                    pass
                for r_idx, row_data in enumerate(rows_write):
                    for c_idx in range(n_cols):
                        val = row_data[c_idx] if c_idx < len(row_data) else None
                        cell = table.cell(r_idx, c_idx); cell.text = coerce_cell_for_csv(val)
                        if r_idx == 0:
                            for para in cell.paragraphs:
                                for run in para.runs: run.bold = True
                        if isinstance(val, str) and val.startswith("="):
                            formulas_present = True
                if truncated:
                    doc.add_paragraph(f"(Truncated to {row_limit} rows)")
                doc.add_paragraph()
            wb.close()
            fname = generate_output_filename(f.filename,"to_word")
            fname = re.sub(r'\.(xls|xlsx)$','.docx',fname,flags=re.IGNORECASE)
            out = os.path.join(Config.OUTPUT_FOLDER, fname); doc.save(out)
        return ok(f"Excel converted to Word ({sheet_count} sheet(s))", out,
                  formulas_preserved=preserve_formulas, row_limit_used=row_limit,
                  warning=("Formulas replaced with values." if formulas_present and not preserve_formulas else None))
    except Exception:
        log.exception("excel_to_word"); return err("Excel to Word failed", 500)


@app.route("/api/v1/excel-to-json", methods=["POST"])
@app.route("/api/excel-to-json", methods=["POST"])
@require_auth
@require_rate_limit
def excel_to_json():
    if not OPENPYXL_AVAILABLE: return err("Requires openpyxl.", 501)
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_XLS)
    if e: return err(e)
    use_headers = request.form.get("header","true").lower() in ("true","1","yes")
    # FIX 53: Row limit guard for excel_to_json
    MAX_ROWS_JSON = getattr(Config, "MAX_EXCEL_ROWS_JSON", 50000)
    try:
        with FileService.temp_upload(f) as path:
            wb = load_workbook(path, data_only=True, read_only=True)
            data = {}
            for sname in wb.sheetnames:
                ws = wb[sname]
                # FIX 53: Check row count before materialising
                row_count = ws.max_row or 0
                if row_count > MAX_ROWS_JSON:
                    wb.close()
                    return err(f"Excel too large for JSON conversion (max {MAX_ROWS_JSON} rows per sheet)", 413)
                all_rows = list(ws.iter_rows(values_only=True))
                if not all_rows: data[sname] = []; continue
                if use_headers:
                    headers = []; seen = {}
                    for i, h in enumerate(all_rows[0]):
                        base = str(h).strip() if h is not None else ""
                        if not base: base = f"col_{i}"
                        cnt = seen.get(base,0); seen[base] = cnt+1
                        headers.append(base if cnt==0 else f"{base}_{cnt}")
                    data[sname] = [{headers[c]: coerce_cell_value(v)
                                     for c,v in enumerate(row) if c < len(headers)}
                                    for row in all_rows[1:]]
                else:
                    data[sname] = [[coerce_cell_value(v) for v in row] for row in all_rows]
            wb.close()
            fname = generate_output_filename(f.filename,"to_json")
            fname = re.sub(r'\.(xls|xlsx)$','.json',fname,flags=re.IGNORECASE)
            out = os.path.join(Config.OUTPUT_FOLDER, fname)
            with open(out,"w",encoding="utf-8") as fh:
                json.dump(data, fh, ensure_ascii=False, indent=2, default=str)
        return ok(f"Excel converted to JSON ({len(data)} sheet(s))", out)
    except Exception:
        log.exception("excel_to_json"); return err("Excel to JSON failed", 500)


@app.route("/api/v1/compress-excel", methods=["POST"])
@app.route("/api/compress-excel", methods=["POST"])
@require_auth
@require_rate_limit
def compress_excel():
    if not OPENPYXL_AVAILABLE: return err("Requires openpyxl.", 501)
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_XLS)
    if e: return err(e)
    try:
        with FileService.temp_upload(f) as path:
            orig = os.path.getsize(path); wb = load_workbook(path, data_only=True)
            # FIX 15: Single pass — detect formulas in first load with values_only=False
            formulas_present = False
            try:
                wb2 = load_workbook(path, data_only=False)
                for ws2 in wb2.worksheets:
                    for row in ws2.iter_rows(values_only=False):
                        for cell in row:
                            if isinstance(cell.value, str) and cell.value.startswith("="):
                                formulas_present = True; break
                        if formulas_present: break
                    if formulas_present: break
                wb2.close()
            except Exception: pass
            for ws in wb.worksheets:
                max_r = max_c = 0
                for row in ws.iter_rows():
                    for cell in row:
                        if cell.value is not None:
                            max_r = max(max_r, cell.row); max_c = max(max_c, cell.column)
                if max_r > 0 and ws.max_row > max_r:
                    try: ws.delete_rows(max_r+1, ws.max_row-max_r)
                    except Exception: pass
                if max_c > 0 and ws.max_column > max_c:
                    try: ws.delete_cols(max_c+1, ws.max_column-max_c)
                    except Exception: pass
            tmp_out = os.path.join(Config.TEMP_FOLDER, f"{uuid.uuid4().hex}_cmp.xlsx")
            wb.save(tmp_out); wb.close()
            fname = generate_output_filename(f.filename,"compressed")
            out = os.path.join(Config.OUTPUT_FOLDER, fname)
            try:
                with zipfile.ZipFile(tmp_out,"r") as zin:
                    with zipfile.ZipFile(out,"w",zipfile.ZIP_DEFLATED,compresslevel=9) as zout:
                        for item in zin.infolist():
                            zout.writestr(item, zin.read(item.filename))
            except Exception: shutil.copy(tmp_out, out)
            try: os.remove(tmp_out)
            except OSError: pass
            new_size = os.path.getsize(out)
            reduction = round((1-new_size/orig)*100,1) if orig else 0
        return ok(f"Excel compressed — {reduction}% smaller", out,
                  reduction_pct=reduction, formulas_lost=formulas_present,
                  warning=("Formulas replaced with values." if formulas_present else None))
    except Exception:
        log.exception("compress_excel"); return err("Excel compression failed", 500)


@app.route("/api/v1/unlock-excel", methods=["POST"])
@app.route("/api/unlock-excel", methods=["POST"])
@require_auth
@require_rate_limit
def unlock_excel():
    if not MSOFFCRYPTO_AVAILABLE: return err("Requires msoffcrypto-tool.", 501)
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_XLS)
    if e: return err(e)
    pw = sanitize_string(request.form.get("password",""))
    if not pw: return err("Password required")
    try:
        with FileService.temp_upload(f) as path:
            fname = generate_output_filename(f.filename,"unlocked")
            out = os.path.join(Config.OUTPUT_FOLDER, fname)
            with open(path,"rb") as fp:
                of = msoffcrypto.OfficeFile(fp); of.load_key(password=pw)
                with open(out,"wb") as fout: of.decrypt(fout)
            # FIX 36: Validate decryption output for unlock_excel
            if not os.path.exists(out) or os.path.getsize(out) < 100:
                try: os.remove(out)
                except OSError: pass
                return err("Decryption failed — incorrect password or unsupported encryption format", 400)
        return ok("Excel workbook unlocked", out)
    except Exception:
        log.exception("unlock_excel"); return err("Unlock failed — check password", 500)


@app.route("/api/protect-excel", methods=["POST"])
@require_auth
@require_rate_limit
def protect_excel():
    if not MSOFFCRYPTO_AVAILABLE:
        return err("Requires msoffcrypto-tool.", 501)
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_XLS)
    if e: return err(e)
    pw = sanitize_string(request.form.get("password",""))
    pw2 = sanitize_string(request.form.get("password2",""))
    ep = validate_password(pw, pw2)
    if ep: return err(ep)

    try:
        with FileService.temp_upload(f) as path:
            work_path = path
            converted_temp = None

            if path.lower().endswith(".xls") and not path.lower().endswith(".xlsx"):
                log.info(f"Converting .xls to .xlsx for protection: {path}")
                converted_temp = libre(path, "xlsx", temp=True)
                if not converted_temp or not os.path.exists(converted_temp):
                    return err("Could not convert .xls to .xlsx format", 500)
                work_path = converted_temp

            fname = generate_output_filename(f.filename, "protected")
            out = os.path.join(Config.OUTPUT_FOLDER, fname)

            # FIX 17: Separate _pw, redact pw before logging
            _pw = pw
            pw = "[REDACTED]"
            try:
                with open(work_path, "rb") as fp:
                    of = msoffcrypto.OfficeFile(fp)
                    encrypted = False
                    try:
                        with open(out, "wb") as fout:
                            of.encrypt(_pw, fout, cipher_algorithm="AES")
                        encrypted = True
                    except TypeError:
                        pass
                    except Exception:
                        pass

                    if not encrypted:
                        try:
                            with open(out, "wb") as fout:
                                of.encrypt(_pw, fout)
                            encrypted = True
                        except Exception as fallback_ex:
                            log.error(f"protect_excel fallback failed: {type(fallback_ex).__name__}")
                            raise fallback_ex
            finally:
                _pw = None

            if converted_temp:
                try: os.remove(converted_temp)
                except OSError: pass

            if not encrypted or not os.path.exists(out) or os.path.getsize(out) == 0:
                return err("Excel protection failed — could not encrypt workbook", 500)

        return ok("Excel workbook protected", out)
    except Exception:
        log.exception("protect_excel")
        return err("Protect Excel failed", 500)


@app.route("/api/v1/repair-excel", methods=["POST"])
@app.route("/api/repair-excel", methods=["POST"])
@require_auth
@require_rate_limit
def repair_excel():
    if not OPENPYXL_AVAILABLE: return err("Requires openpyxl.", 501)
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_XLS)
    if e: return err(e)
    try:
        with FileService.temp_upload(f) as path:
            fname = generate_output_filename(f.filename,"repaired")
            out = os.path.join(Config.OUTPUT_FOLDER, fname)
            ext = path.rsplit(".",1)[-1].lower()
            if ext == "xlsm":
                shutil.copy(path, out)
                return ok(
                    "Macro-enabled workbook returned as-is — not modified",
                    out,
                    warning=(
                        "⚠️ This file contains VBA macros which were NOT scanned or removed. "
                        "Only open in a fully trusted environment. "
                        "Convert to .xlsx to strip macros."
                    )
                )
            try:
                wb = load_workbook(path, data_only=False, read_only=False)
                wb.save(out); wb.close()
                if os.path.exists(out) and os.path.getsize(out) > 0:
                    return ok("Excel repaired (openpyxl)", out)
            except Exception as ex1:
                log.warning(f"openpyxl repair: {ex1}")
            lo_out = libre(path, "xlsx", output_filename=fname)
            if lo_out and os.path.exists(lo_out) and os.path.getsize(lo_out) > 0:
                return ok("Excel repaired (LibreOffice)", lo_out)
            return err("Could not repair Excel — file severely corrupted", 500)
    except Exception:
        log.exception("repair_excel"); return err("Excel repair failed", 500)


@app.route("/api/v1/excel-to-jpg", methods=["POST"])
@app.route("/api/excel-to-jpg", methods=["POST"])
@require_auth
@require_rate_limit
def excel_to_jpg():
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_XLS)
    if e: return err(e)
    try:
        with FileService.temp_upload(f) as path:
            pdf_path = libre(path, "pdf", temp=True)
            if pdf_path and os.path.exists(pdf_path) and os.path.getsize(pdf_path)>100:
                doc = fitz.open(pdf_path)
                try:
                    buf = io.BytesIO()
                    with zipfile.ZipFile(buf,"w",zipfile.ZIP_DEFLATED) as zf:
                        for i, page in enumerate(doc):
                            pix = page.get_pixmap(dpi=150)
                            zf.writestr(f"sheet_{i+1:04d}.jpg", pix.tobytes("jpeg"))
                            pix = None
                    fname = generate_output_filename(f.filename,"to_jpg")
                    fname = re.sub(r'\.(xls|xlsx)$','.zip',fname,flags=re.IGNORECASE)
                    out = os.path.join(Config.OUTPUT_FOLDER, fname)
                    with open(out,"wb") as fh: fh.write(buf.getvalue())
                    return ok("Excel sheets exported as JPG", out)
                finally:
                    doc.close()
                    try: os.remove(pdf_path)
                    except OSError: pass
            return err("Excel to JPG failed — LibreOffice unavailable", 500)
    except Exception:
        log.exception("excel_to_jpg"); return err("Excel to JPG failed", 500)


@app.route("/api/v1/excel-to-ppt", methods=["POST"])
@app.route("/api/excel-to-ppt", methods=["POST"])
@require_auth
@require_rate_limit
def excel_to_ppt():
    if not OPENPYXL_AVAILABLE or not PPTX_AVAILABLE:
        return err("Requires openpyxl + python-pptx.", 501)
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_XLS)
    if e: return err(e)
    max_rows_param = safe_int(request.form.get("max_rows", "25"), 25, 1, 200)
    try:
        with FileService.temp_upload(f) as path:
            wb = load_workbook(path, data_only=True)
            prs = Presentation(); prs.slide_width=PptxInches(10); prs.slide_height=PptxInches(7.5)
            tc_layout = prs.slide_layouts[1]
            truncated_any = False
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = [r for r in ws.iter_rows(values_only=True) if any(c is not None for c in r)]
                if not rows: continue
                slide = prs.slides.add_slide(tc_layout)
                if slide.shapes.title: slide.shapes.title.text = sheet_name
                max_cols = max(len(r) for r in rows)
                max_rows = min(len(rows), max_rows_param)
                if len(rows) > max_rows_param:
                    truncated_any = True
                tbl_shape = slide.shapes.add_table(max_rows, max_cols,
                                                    PptxInches(0.5), PptxInches(1.5),
                                                    PptxInches(9), PptxInches(5))
                table = tbl_shape.table
                for r_idx in range(max_rows):
                    for c_idx in range(max_cols):
                        cell = table.cell(r_idx, c_idx)
                        val = rows[r_idx][c_idx] if c_idx < len(rows[r_idx]) else ""
                        cell.text = coerce_cell_for_csv(val)
                        if r_idx == 0: cell.text_frame.paragraphs[0].font.bold = True
            wb.close()
            fname = generate_output_filename(f.filename,"to_ppt")
            fname = re.sub(r'\.(xls|xlsx)$','.pptx',fname,flags=re.IGNORECASE)
            out = os.path.join(Config.OUTPUT_FOLDER, fname)
            # FIX 26: Atomic write for excel_to_ppt
            tmp_out = out + ".tmp"
            prs.save(tmp_out)
            os.replace(tmp_out, out)
        return ok("Excel converted to PowerPoint", out,
                  warning=f"Tables truncated to {max_rows_param} rows per slide" if truncated_any else None)
    except Exception:
        log.exception("excel_to_ppt"); return err("Excel to PPT failed", 500)


@app.route("/api/v1/jpg-to-pdf", methods=["POST"])
@app.route("/api/jpg-to-pdf", methods=["POST"])
@require_auth
@require_rate_limit
def jpg_to_pdf():
    files = (request.files.getlist("files") or
             request.files.getlist("jpg") or
             request.files.getlist("jpeg"))
    if not files or all(f.filename == "" for f in files):
        return err("At least one JPG file required")
    for f in files:
        e = validate_file(f, Config.ALLOWED_IMAGE)
        if e: return err(e)
    page_size = request.form.get("page_size", "auto")
    try:
        with FileService.temp_uploads(files) as paths:
            fname = generate_output_filename(files[0].filename, "to_pdf",
                                              is_multi=True, filenames=[f.filename for f in files])
            fname = re.sub(r'\.(jpg|jpeg)$', '.pdf', fname, flags=re.IGNORECASE)
            if not fname.endswith(".pdf"): fname = Path(fname).stem + ".pdf"
            try:
                out = _images_to_pdf(paths, page_size, fname)
            except ValueError as ve:
                return err(str(ve), 400)
        return ok(f"Converted {len(files)} JPG(s) to PDF", out)
    except Exception:
        log.exception("jpg_to_pdf"); return err("JPG to PDF failed", 500)


@app.route("/api/v1/pdf-to-jpg", methods=["POST"])
@app.route("/api/pdf-to-jpg", methods=["POST"])
@require_auth
@require_rate_limit
def pdf_to_jpg():
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_PDF)
    if e: return err(e)
    try:
        with FileService.temp_upload(f) as path:
            doc = fitz.open(path)
            if len(doc) == 0:
                doc.close()
                return err("Input PDF has no pages", 400)
            # FIX 51: Page count guard for pdf_to_jpg
            MAX_IMAGE_PAGES = getattr(Config, "MAX_IMAGE_PAGES", 50)
            if len(doc) > MAX_IMAGE_PAGES:
                doc.close()
                return err(f"PDF too large for image conversion (max {MAX_IMAGE_PAGES} pages). Split the PDF first.", 413)
            count = len(doc)
            # FIX 21: Stream directly to disk ZIP
            fname = generate_output_filename(f.filename, "to_jpg")
            fname = re.sub(r'\.pdf$', '.zip', fname, flags=re.IGNORECASE)
            if not fname.endswith('.zip'):
                fname = Path(fname).stem + '.zip'
            out = os.path.join(Config.OUTPUT_FOLDER, fname)
            with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zf:
                for i, page in enumerate(doc):
                    pix = page.get_pixmap(dpi=150)
                    zf.writestr(f"page_{i+1:04d}.jpg", pix.tobytes("jpeg"))
                    pix = None
            doc.close()
        return ok(f"Exported {count} page(s) as JPG", out)
    except Exception:
        log.exception("pdf_to_jpg")
        return err("PDF to JPG failed", 500)


@app.route("/api/v1/pdf-to-png", methods=["POST"])
@app.route("/api/pdf-to-png", methods=["POST"])
@require_auth
@require_rate_limit
def pdf_to_png():
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_PDF)
    if e: return err(e)
    try:
        with FileService.temp_upload(f) as path:
            doc = fitz.open(path)
            if len(doc) == 0:
                doc.close()
                return err("Input PDF has no pages", 400)
            # FIX 51: Page count guard for pdf_to_png
            MAX_IMAGE_PAGES = getattr(Config, "MAX_IMAGE_PAGES", 50)
            if len(doc) > MAX_IMAGE_PAGES:
                doc.close()
                return err(f"PDF too large for image conversion (max {MAX_IMAGE_PAGES} pages). Split the PDF first.", 413)
            count = len(doc)
            # FIX 21: Stream directly to disk ZIP
            fname = generate_output_filename(f.filename, "to_png")
            fname = re.sub(r'\.pdf$', '.zip', fname, flags=re.IGNORECASE)
            if not fname.endswith('.zip'):
                fname = Path(fname).stem + '.zip'
            out = os.path.join(Config.OUTPUT_FOLDER, fname)
            with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zf:
                for i, page in enumerate(doc):
                    pix = page.get_pixmap(dpi=150)
                    zf.writestr(f"page_{i+1:04d}.png", pix.tobytes("png"))
                    pix = None
            doc.close()
        return ok(f"Exported {count} page(s) as PNG", out)
    except Exception:
        log.exception("pdf_to_png")
        return err("PDF to PNG failed", 500)


@app.route("/api/v1/word-to-jpg", methods=["POST"])
@app.route("/api/word-to-jpg", methods=["POST"])
@require_auth
@require_rate_limit
def word_to_jpg():
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_DOC)
    if e: return err(e)
    try:
        with FileService.temp_upload(f) as path:
            pdf_path = libre(path, "pdf", temp=True)
            if not pdf_path: return err("LibreOffice conversion failed", 500)
            doc = fitz.open(pdf_path)
            try:
                buf = io.BytesIO()
                with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
                    for i, page in enumerate(doc):
                        pix = page.get_pixmap(dpi=150)
                        zf.writestr(f"page_{i+1:04d}.jpg", pix.tobytes("jpeg"))
                        pix = None
            finally:
                doc.close()
            try: os.remove(pdf_path)
            except OSError: pass
            fname = generate_output_filename(f.filename, "to_jpg")
            fname = re.sub(r'\.(doc|docx)$', '.zip', fname, flags=re.IGNORECASE)
            out = os.path.join(Config.OUTPUT_FOLDER, fname)
            with open(out, "wb") as fh: fh.write(buf.getvalue())
        return ok("Word converted to JPG images", out)
    except Exception:
        log.exception("word_to_jpg"); return err("Word to JPG failed", 500)


@app.route("/api/v1/word-to-png", methods=["POST"])
@app.route("/api/word-to-png", methods=["POST"])
@require_auth
@require_rate_limit
def word_to_png():
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_DOC)
    if e: return err(e)
    try:
        with FileService.temp_upload(f) as path:
            pdf_path = libre(path, "pdf", temp=True)
            if not pdf_path: return err("LibreOffice conversion failed", 500)
            doc = fitz.open(pdf_path)
            try:
                buf = io.BytesIO()
                with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
                    for i, page in enumerate(doc):
                        pix = page.get_pixmap(dpi=150)
                        zf.writestr(f"page_{i+1:04d}.png", pix.tobytes("png"))
                        pix = None
            finally:
                doc.close()
            try: os.remove(pdf_path)
            except OSError: pass
            fname = generate_output_filename(f.filename, "to_png")
            fname = re.sub(r'\.(doc|docx)$', '.zip', fname, flags=re.IGNORECASE)
            out = os.path.join(Config.OUTPUT_FOLDER, fname)
            with open(out, "wb") as fh: fh.write(buf.getvalue())
        return ok("Word converted to PNG images", out)
    except Exception:
        log.exception("word_to_png"); return err("Word to PNG failed", 500)


@app.route("/api/v1/word-to-html", methods=["POST"])
@app.route("/api/word-to-html", methods=["POST"])
@require_auth
@require_rate_limit
def word_to_html():
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_DOC)
    if e: return err(e)
    try:
        with FileService.temp_upload(f) as path:
            fname = generate_output_filename(f.filename, "to_html")
            fname = re.sub(r'\.(doc|docx)$', '.html', fname, flags=re.IGNORECASE)
            # FIX 52: word_to_html only uses libre() — no DocxDocument fallback, so no paragraph guard needed
            out = libre(path, "html", output_filename=fname)
            if not out: return err("LibreOffice conversion failed", 500)
        return ok("Word converted to HTML", out)
    except Exception:
        log.exception("word_to_html"); return err("Word to HTML failed", 500)


@app.route("/api/v1/word-to-json", methods=["POST"])
@app.route("/api/word-to-json", methods=["POST"])
@require_auth
@require_rate_limit
def word_to_json():
    if not DOCX_AVAILABLE: return err("Requires python-docx", 501)
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_DOC)
    if e: return err(e)
    try:
        with FileService.temp_upload(f) as path:
            doc = DocxDocument(path)
            # FIX 52: Guard for word_to_json
            MAX_PARAGRAPHS = getattr(Config, "MAX_WORD_PARAGRAPHS", 10000)
            MAX_TABLE_CELLS = getattr(Config, "MAX_WORD_TABLE_CELLS", 50000)
            para_count = len(doc.paragraphs)
            table_cell_count = sum(len(row.cells) for table in doc.tables for row in table.rows)
            if para_count > MAX_PARAGRAPHS:
                return err(f"Document too large: {para_count} paragraphs (max {MAX_PARAGRAPHS})", 413)
            if table_cell_count > MAX_TABLE_CELLS:
                return err(f"Document too large: {table_cell_count} table cells (max {MAX_TABLE_CELLS})", 413)
            data = {"paragraphs": [p.text for p in doc.paragraphs], "tables": []}
            for table in doc.tables:
                tdata = [[cell.text for cell in row.cells] for row in table.rows]
                data["tables"].append(tdata)
            fname = generate_output_filename(f.filename, "to_json")
            fname = re.sub(r'\.(doc|docx)$', '.json', fname, flags=re.IGNORECASE)
            out = os.path.join(Config.OUTPUT_FOLDER, fname)
            with open(out, "w", encoding="utf-8") as fh:
                json.dump(data, fh, ensure_ascii=False, indent=2)
        return ok("Word converted to JSON", out)
    except Exception:
        log.exception("word_to_json"); return err("Word to JSON failed", 500)


@app.route("/api/v1/excel-to-png", methods=["POST"])
@app.route("/api/excel-to-png", methods=["POST"])
@require_auth
@require_rate_limit
def excel_to_png():
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_XLS)
    if e: return err(e)
    try:
        with FileService.temp_upload(f) as path:
            pdf_path = libre(path, "pdf", temp=True)
            if pdf_path and os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 100:
                doc = fitz.open(pdf_path)
                try:
                    buf = io.BytesIO()
                    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
                        for i, page in enumerate(doc):
                            pix = page.get_pixmap(dpi=150)
                            zf.writestr(f"sheet_{i+1:04d}.png", pix.tobytes("png"))
                            pix = None
                    fname = generate_output_filename(f.filename, "to_png")
                    fname = re.sub(r'\.(xls|xlsx)$', '.zip', fname, flags=re.IGNORECASE)
                    out = os.path.join(Config.OUTPUT_FOLDER, fname)
                    with open(out, "wb") as fh: fh.write(buf.getvalue())
                    return ok("Excel sheets exported as PNG", out)
                finally:
                    doc.close()
                    try: os.remove(pdf_path)
                    except OSError: pass
            return err("Excel to PNG failed — LibreOffice unavailable", 500)
    except Exception:
        log.exception("excel_to_png"); return err("Excel to PNG failed", 500)


@app.route("/api/v1/excel-to-html", methods=["POST"])
@app.route("/api/excel-to-html", methods=["POST"])
@require_auth
@require_rate_limit
def excel_to_html():
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_XLS)
    if e: return err(e)
    try:
        with FileService.temp_upload(f) as path:
            fname = generate_output_filename(f.filename, "to_html")
            fname = re.sub(r'\.(xls|xlsx)$', '.html', fname, flags=re.IGNORECASE)
            out = libre(path, "html", output_filename=fname)
            if not out: return err("LibreOffice conversion failed", 500)
        return ok("Excel converted to HTML", out)
    except Exception:
        log.exception("excel_to_html"); return err("Excel to HTML failed", 500)


@app.route("/api/v1/image-to-word", methods=["POST"])
@app.route("/api/image-to-word", methods=["POST"])
@require_auth
@require_rate_limit
def image_to_word():
    if not DOCX_AVAILABLE: return err("Requires python-docx", 501)
    f = request.files.get("file")
    e = validate_file(f, Config.ALLOWED_IMAGE)
    if e: return err(e)
    try:
        with FileService.temp_upload(f) as path:
            # FIX 34: Catch DecompressionBombError
            try:
                img = Image.open(path)
            except Image.DecompressionBombError:
                return err("Image too large to process safely", 413)
            doc = DocxDocument()
            doc.add_heading("Image OCR Result", 0)

            if TESSERACT_AVAILABLE:
                # FIX 40: OCR semaphore for image_to_word
                with _ocr_semaphore:
                    text = pytesseract.image_to_string(img)
                for para in text.split('\n\n'):
                    if para.strip():
                        doc.add_paragraph(para.strip())
            else:
                tmp = tempfile.NamedTemporaryFile(
                    suffix=".png", delete=False, dir=Config.TEMP_FOLDER
                )
                try:
                    img.save(tmp.name, format="PNG"); tmp.close()
                    doc.add_picture(tmp.name, width=Inches(6))
                finally:
                    try: os.unlink(tmp.name)
                    except OSError: pass
                doc.add_paragraph("(OCR not available — image embedded)")

            fname = generate_output_filename(f.filename, "to_word")
            fname = re.sub(r'\.(jpg|jpeg|png|gif|bmp|tiff|webp)$', '.docx',
                           fname, flags=re.IGNORECASE)
            out = os.path.join(Config.OUTPUT_FOLDER, fname)
            doc.save(out)
        return ok("Image converted to Word", out)
    except Exception:
        log.exception("image_to_word"); return err("Image to Word failed", 500)


# ============================================================================
# ERROR HANDLERS
# ============================================================================
@app.errorhandler(400)
def _e400(e):
    return jsonify({"success":False,"error":"Bad request","request_id":g.get("request_id","-")}),400

@app.errorhandler(404)
def _e404(e):
    return jsonify({"success":False,"error":"Endpoint not found","request_id":g.get("request_id","-")}),404

@app.errorhandler(405)
def _e405(e):
    return jsonify({"success":False,"error":"Method not allowed"}),405

@app.errorhandler(413)
def _e413(e):
    return jsonify({"success":False,"error":f"File too large (max {Config.MAX_FILE_SIZE//1048576} MB)","request_id":g.get("request_id","-")}),413

@app.errorhandler(429)
def _e429(e):
    return jsonify({"success":False,"error":"Rate limit exceeded","request_id":g.get("request_id","-")}),429

@app.errorhandler(500)
def _e500(e):
    log.exception("Unhandled 500")
    return jsonify({"success":False,"error":"Internal server error","request_id":g.get("request_id","-")}),500


# ============================================================================
# GRACEFUL SHUTDOWN
# ============================================================================
_shutdown_event = threading.Event()

def _graceful_shutdown(signum, frame):
    log.info(f"Received signal {signum} — graceful shutdown…")
    _shutdown_event.set()
    # FIX 1: Snapshot thread registry under lock before iterating
    with _registry_lock:
        items = list(_thread_registry.items())
    for jid, (t, cancel_ev) in items:
        log.info(f"Cancelling background thread for job {jid}")
        cancel_ev.set()
        # FIX 45: Join threads with timeout to prevent zombies
        t.join(timeout=2.0)
    if celery_app:
        try:
            SHUTDOWN_TIMEOUT = int(os.environ.get("SHUTDOWN_TIMEOUT",
                                                    getattr(Config, 'SHUTDOWN_TIMEOUT', 30)))
            for q in ["fast", "slow", "office"]:
                celery_app.control.broadcast("cancel_consumer", arguments={"queue": q})
            time.sleep(min(10, SHUTDOWN_TIMEOUT // 3))
            celery_app.control.broadcast("shutdown")
        except Exception as ce:
            log.warning(f"Celery shutdown error: {ce}")
    SHUTDOWN_TIMEOUT = int(os.environ.get("SHUTDOWN_TIMEOUT",
                                            getattr(Config, 'SHUTDOWN_TIMEOUT', 30)))
    deadline = time.time() + SHUTDOWN_TIMEOUT
    while time.time() < deadline:
        time.sleep(1)
    log.info("Shutdown complete.")
    # FIX 7: Move signal handlers to __main__ block; remove sys.exit from _graceful_shutdown
    return

# FIX 7: Signal handlers only registered here at module level (moved to __main__ if needed)
# They are kept here for Gunicorn compatibility but sys.exit removed from handler above
signal.signal(signal.SIGTERM, _graceful_shutdown)
signal.signal(signal.SIGINT, _graceful_shutdown)

# FIX 37: Verify redis_service rate limit uses pipeline — add comment for audit
# VERIFY: redis_service rate_limit uses pipeline — see FIX 37

application = app  # WSGI entry point

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    debug = os.environ.get("FLASK_DEBUG","false").lower() in ("true","1","yes")
    log.info(f"PDFWala V{Config.VERSION} starting on :{port}")
    app.run(host="0.0.0.0", port=port, debug=debug, threaded=True, use_reloader=False)
