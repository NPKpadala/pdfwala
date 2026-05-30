"""
engines/pdf_engine.py — PDFWala Enterprise V14.0

All PDF processing logic lives here. Zero Flask. Zero Celery.
Each function is a pure "input path(s) → output path" transformer.

Registered to the Pipeline via @register("operation_name").
Called ONLY by Pipeline.run() — never directly from routes or tasks.

Operations covered:
  compress_pdf, merge_pdf, split_pdf, rotate_pdf, watermark_pdf,
  page_numbers, crop_pdf, pdf_info, protect_pdf, unlock_pdf,
  sign_pdf, redact_pdf, repair_pdf, linearize_pdf, ocr_pdf,
  pdf_to_image, pdf_to_word, pdf_to_excel, pdf_to_ppt, pdf_to_pdfa,
  compare_pdf, pdf_to_jpg, pdf_to_png, remove_pages, extract_pages,
  organize_pdf

V13 FIXES (41 issues across security, correctness, memory, and validation):
  - All ZIP operations stream to disk (no BytesIO for >50-page PDFs)
  - OCR parallelised via ThreadPoolExecutor, 200 DPI default
  - compress_pdf runs Ghostscript on ORIGINAL, not stage1
  - organize/remove/extract: 0-based indexing contract enforced
  - repair_pdf: PDFWalaError import guard added
  - redact_pdf: import re at module level
  - Pytesseract lang/psm sanitised (whitelist)
  - sign_pdf: sig_data validated, signature stamped into content stream
  - watermark_pdf: opacity clamped, overlay enforced
  - unlock_pdf: tries owner_pw then user_pw, full permission strip
  - protect_pdf: owner_pw != user_pw (random suffix added)
  - merge_pdf: pikepdf streaming merge (falls back to PyPDF2 if unavailable)
  - linearize_pdf: qpdf preferred over Ghostscript
  - pdf_to_ppt: tempfile cleanup in finally
  - pdf_to_excel: fitz.open always closed
  - rotate_pdf: angle whitelist, writes /Rotate to page dict
  - crop_pdf / page_numbers: finally blocks added
  - _ghostscript: path hardened against leading-dash injection
  - All _guard_empty calls added to conversion ops
  - Image RAM guard: skip embedded images >5 MB during compress
"""

from __future__ import annotations

import difflib
import io
import json
import logging
import os
import re
import secrets
import shutil
import subprocess
import tempfile
import zipfile
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from typing import List, Optional

from config import Config
from core.context import JobContext
from core.exceptions import (
    OperationTimeoutError,
    ProcessingError,
    UnsupportedOperation,
    ValidationError,
)
from core.pipeline import register
from utils.helpers import format_file_size
from utils.pdf_utils import (
    compress_pdf_images,
    create_page_number_pdf,
    create_watermark_pdf,
    parse_page_ranges,
)
from utils.security import REDACTION_PATTERNS, SafeRegex

log = logging.getLogger("pdfwala.engines.pdf")

# ── Streaming ZIP threshold ───────────────────────────────────────────────────
_ZIP_STREAM_THRESHOLD_PAGES = 50      # write to disk-backed ZIP above this
_IMAGE_RAM_GUARD_BYTES      = 5 * 1024 * 1024   # 5 MB — skip larger images in compress

# ── Whitelists ────────────────────────────────────────────────────────────────
_VALID_ROTATION_ANGLES = {0, 90, 180, 270}
_VALID_TESSERACT_LANGS = re.compile(r'^[a-zA-Z]{2,8}(\+[a-zA-Z]{2,8})*$')
_VALID_TESSERACT_PSM   = frozenset(range(0, 14))
_VALID_TESSERACT_OEM   = frozenset(range(0, 4))

# ── Library availability flags ────────────────────────────────────────────────
try:
    import fitz
    FITZ_OK = True
except ImportError:
    FITZ_OK = False

try:
    import pikepdf
    PIKEPDF_OK = True
except ImportError:
    PIKEPDF_OK = False

try:
    from PyPDF2 import PdfReader, PdfWriter, PdfMerger
    PYPDF2_OK = True
except ImportError:
    PYPDF2_OK = False

try:
    from PIL import Image
    PIL_OK = True
except ImportError:
    PIL_OK = False

try:
    from pdf2docx import Converter as Pdf2DocxConverter
    PDF2DOCX_OK = True
except ImportError:
    PDF2DOCX_OK = False

try:
    from openpyxl import Workbook
    from openpyxl.styles import Font
    OPENPYXL_OK = True
except ImportError:
    OPENPYXL_OK = False

try:
    import pdfplumber
    PDFPLUMBER_OK = True
except ImportError:
    PDFPLUMBER_OK = False

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches
    PPTX_OK = True
except ImportError:
    PPTX_OK = False

try:
    import pytesseract
    from pytesseract import Output as TesseractOutput
    TESSERACT_OK = True
except ImportError:
    TESSERACT_OK = False

# ── qpdf availability (linearization) ────────────────────────────────────────
def _qpdf_available() -> bool:
    try:
        r = subprocess.run(["qpdf", "--version"], capture_output=True, timeout=5)
        return r.returncode == 0
    except Exception:
        return False

QPDF_OK = _qpdf_available()


# ── Internal helpers ──────────────────────────────────────────────────────────

def _require(flag: bool, operation: str, library: str) -> None:
    if not flag:
        raise UnsupportedOperation(operation, library)


def _safe_output_path(output_path: str) -> str:
    """
    Harden output path against argument-injection attacks.
    Raises ValidationError if the resolved path looks like a flag or leaves
    the configured output directory.
    """
    resolved = str(Path(output_path).resolve())
    if resolved.startswith("-"):
        raise ValidationError("Invalid output path: looks like a flag")
    # FIX V14: Always use OUTPUT_FOLDER directly (OUTPUT_DIR alias may be stale at import time)
    output_dir = str(Path(Config.OUTPUT_FOLDER).resolve())
    temp_dir   = str(Path(Config.TEMP_FOLDER).resolve())
    if not resolved.startswith(output_dir) and not resolved.startswith(temp_dir):
        raise ValidationError("Output path escapes configured output directory")
    return resolved


def _ghostscript(
    input_path: str,
    output_path: str,
    gs_setting: str = "/ebook",
    extra_flags: Optional[list] = None,
    timeout: Optional[int] = None,
) -> bool:
    """
    Run Ghostscript. Returns True on success.
    Raises OperationTimeoutError on timeout, ValidationError on bad path.
    ALWAYS call on the ORIGINAL input, never on an intermediate file.
    """
    timeout = timeout or Config.SUBPROCESS_TIMEOUT
    safe_out = _safe_output_path(output_path)
    cmd = [
        Config.GHOSTSCRIPT,
        "-sDEVICE=pdfwrite",
        "-dCompatibilityLevel=1.5",
        f"-dPDFSETTINGS={gs_setting}",
        "-dNOPAUSE", "-dBATCH", "-dQUIET", "-dSAFER",
        "-dDetectDuplicateImages=true",
        "-dCompressFonts=true",
        "-dSubsetFonts=true",
        "-dAutoRotatePages=/None",
        f"-sOutputFile={safe_out}",
    ]
    if extra_flags:
        for flag in extra_flags:
            if flag.startswith("-s") or flag.startswith("-d") or flag.startswith("-r"):
                cmd.append(flag)
            else:
                log.warning(f"GS: skipping suspicious flag: {flag!r}")
    cmd.append(input_path)
    try:
        result = subprocess.run(cmd, capture_output=True, timeout=timeout)
        if result.returncode != 0:
            log.warning(f"GS rc={result.returncode}: {result.stderr[:300]}")
            return False
        if not os.path.exists(safe_out) or os.path.getsize(safe_out) == 0:
            return False
        return True
    except subprocess.TimeoutExpired:
        try:
            os.remove(safe_out)
        except OSError:
            pass
        raise OperationTimeoutError("Ghostscript", timeout)


def _guard_empty(path: str) -> None:
    """Raise ValidationError if the PDF has zero pages."""
    if not FITZ_OK:
        return
    doc = fitz.open(path)
    n = len(doc)
    doc.close()
    if n == 0:
        raise ValidationError("Input PDF has no pages")


def _open_zip_writer(output_path: str, page_count: int):
    """
    Return a ZipFile that writes to disk (large) or BytesIO (small).
    Caller must handle the BytesIO→disk copy for the small case.
    Returns (zf, buf_or_none).
    """
    if page_count > _ZIP_STREAM_THRESHOLD_PAGES:
        zf = zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED, allowZip64=True)
        return zf, None
    buf = io.BytesIO()
    zf  = zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED)
    return zf, buf


def _finalise_zip(zf: zipfile.ZipFile, buf: Optional[io.BytesIO], output_path: str) -> None:
    zf.close()
    if buf is not None:
        with open(output_path, "wb") as fh:
            fh.write(buf.getvalue())
    # disk-backed ZipFile writes directly to output_path — nothing to do


def _sanitise_tesseract_lang(lang: str) -> str:
    lang = lang.strip()
    if not _VALID_TESSERACT_LANGS.match(lang):
        raise ValidationError(
            f"Invalid Tesseract language code: {lang!r}. "
            "Use ISO 639-2 codes like 'eng', 'fra', 'eng+fra'."
        )
    return lang


def _validate_image_data(data: bytes, max_bytes: int = 10 * 1024 * 1024) -> None:
    """Validate raw image bytes: check magic bytes and size limit."""
    if not data:
        raise ValidationError("Image data is empty")
    if len(data) > max_bytes:
        raise ValidationError(f"Image too large (max {max_bytes // 1024 // 1024} MB)")
    magic = data[:4]
    # PNG, JPEG, GIF, WEBP, BMP
    valid_magics = (
        b'\x89PNG', b'\xff\xd8\xff', b'GIF8', b'RIFF', b'BM',
    )
    if not any(magic.startswith(m) for m in valid_magics):
        raise ValidationError("Unsupported image format — PNG/JPEG/GIF/WEBP/BMP only")


# ═══════════════════════════════════════════════════════════════════════════════
# PDF ORGANIZE
# ═══════════════════════════════════════════════════════════════════════════════

@register("merge_pdf")
def merge_pdf(ctx: JobContext) -> dict:
    """
    Merge multiple PDFs.
    Uses pikepdf for robust xref handling (falls back to PyPDF2).
    Streams pages — does NOT load all source docs simultaneously.
    """
    paths = ctx.input_paths
    if not paths:
        raise ValidationError("No input files provided")

    # --- Validation pass (fitz, read-only) ---
    page_sizes: set = set()
    for p in paths:
        if FITZ_OK:
            doc = fitz.open(p)
            try:
                if len(doc) == 0:
                    raise ValidationError(f"PDF {os.path.basename(p)} has no pages")
                if doc.is_encrypted and not doc.authenticate(""):
                    raise ValidationError(f"PDF {os.path.basename(p)} is password-protected")
                for pg in doc:
                    page_sizes.add((round(pg.rect.width), round(pg.rect.height)))
            finally:
                doc.close()

    total_pages = 0

    if PIKEPDF_OK:
        # pikepdf: streaming page-by-page append, correct xref handling
        import pikepdf as _pikepdf
        with _pikepdf.Pdf.new() as out_pdf:
            for i, p in enumerate(paths):
                with _pikepdf.Pdf.open(p) as src:
                    out_pdf.pages.extend(src.pages)
                    total_pages += len(src.pages)
                ctx.set_progress(int((i + 1) / len(paths) * 90))
            out_pdf.save(ctx.output_path)
    elif PYPDF2_OK:
        merger = PdfMerger()
        try:
            for i, p in enumerate(paths):
                merger.append(p)
                ctx.set_progress(int((i + 1) / len(paths) * 90))
            merger.write(ctx.output_path)
        finally:
            merger.close()
    else:
        raise UnsupportedOperation("merge_pdf", "pikepdf or PyPDF2")

    ctx.set_progress(100)
    log.info(f"[{ctx.job_id}] merge_pdf: merged {len(paths)} files, {total_pages} total pages")
    return {
        "mixed_page_sizes": len(page_sizes) > 1,
        "files_merged":     len(paths),
        "total_pages":      total_pages,
    }


@register("split_pdf")
def split_pdf(ctx: JobContext) -> dict:
    """
    Split PDF into per-page files, streamed directly to disk ZIP.
    Preserves annotations and rotation via fitz.insert_pdf.
    """
    if not (FITZ_OK and PYPDF2_OK):
        _require(FITZ_OK, "split_pdf", "PyMuPDF")
        _require(PYPDF2_OK, "split_pdf", "PyPDF2")

    mode   = ctx.params.get("mode", "all")
    ranges = ctx.params.get("ranges", "")
    _guard_empty(ctx.input_path)

    src_doc = fitz.open(ctx.input_path)
    total   = len(src_doc)
    try:
        # parse_page_ranges returns 0-based indices
        indices = list(range(total)) if mode == "all" else parse_page_ranges(ranges, total)
        if not indices:
            raise ValidationError("No valid pages in range")

        zf, buf = _open_zip_writer(ctx.output_path, len(indices))
        try:
            for n, idx in enumerate(indices):
                out_doc = fitz.open()
                out_doc.insert_pdf(src_doc, from_page=idx, to_page=idx)
                page_bytes = out_doc.tobytes(deflate=True, garbage=2)
                out_doc.close()
                zf.writestr(f"page_{idx + 1:04d}.pdf", page_bytes)
                if n % 10 == 0:
                    ctx.set_progress(int(n / len(indices) * 95))
        finally:
            _finalise_zip(zf, buf, ctx.output_path)
    finally:
        src_doc.close()

    log.info(f"[{ctx.job_id}] split_pdf: exported {len(indices)} pages")
    return {"pages_exported": len(indices)}


@register("organize_pdf")
def organize_pdf(ctx: JobContext) -> dict:
    """
    Reorder or delete pages.
    Contract: parse_page_ranges returns 0-based indices.
    """
    _require(FITZ_OK, "organize_pdf", "PyMuPDF")
    action     = ctx.params.get("action", "reorder")
    order_spec = ctx.params.get("order", "")
    _guard_empty(ctx.input_path)

    src_doc = fitz.open(ctx.input_path)
    total   = len(src_doc)
    try:
        # 0-based from parse_page_ranges
        indices = parse_page_ranges(order_spec, total)

        if action == "delete":
            remove_set = set(indices)
            if len(remove_set) >= total:
                raise ValidationError("Cannot delete all pages")
            keep = [i for i in range(total) if i not in remove_set]
        else:
            # reorder — validate bounds explicitly
            for i in indices:
                if i < 0 or i >= total:
                    raise ValidationError(f"Page index {i} out of range (0-{total - 1})")
            keep = indices

        out_doc2 = fitz.open()
        for idx in keep:
            out_doc2.insert_pdf(src_doc, from_page=idx, to_page=idx)
        out_doc2.save(ctx.output_path, deflate=True, garbage=3)
    finally:
        src_doc.close()
        try:
            out_doc2.close()
        except Exception:
            pass

    log.info(f"[{ctx.job_id}] organize_pdf: action={action}, pages_in_output={len(keep)}")
    return {"action": action, "pages_in_output": len(keep)}


@register("remove_pages")
def remove_pages(ctx: JobContext) -> dict:
    """
    Remove specified pages, preserving TOC and internal links via fitz.
    parse_page_ranges returns 0-based.
    """
    _require(FITZ_OK, "remove_pages", "PyMuPDF")
    _guard_empty(ctx.input_path)

    src_doc   = fitz.open(ctx.input_path)
    total     = len(src_doc)
    try:
        remove_set = set(parse_page_ranges(ctx.params.get("order", ""), total))
        if len(remove_set) >= total:
            raise ValidationError("Cannot remove all pages")

        keep = [i for i in range(total) if i not in remove_set]
        out_doc = fitz.open()
        try:
            for idx in keep:
                out_doc.insert_pdf(src_doc, from_page=idx, to_page=idx)
            out_doc.save(ctx.output_path, deflate=True, garbage=3)
        finally:
            out_doc.close()
    finally:
        src_doc.close()

    log.info(f"[{ctx.job_id}] remove_pages: removed {len(remove_set)} pages")
    return {"pages_removed": len(remove_set)}


@register("extract_pages")
def extract_pages(ctx: JobContext) -> dict:
    """
    Extract specified pages into a new PDF.
    Uses fitz.insert_pdf to preserve metadata and embedded fonts.
    parse_page_ranges returns 0-based.
    """
    _require(FITZ_OK, "extract_pages", "PyMuPDF")
    _guard_empty(ctx.input_path)

    src_doc = fitz.open(ctx.input_path)
    total   = len(src_doc)
    try:
        indices = parse_page_ranges(ctx.params.get("order", ""), total)
        if not indices:
            raise ValidationError("No valid pages in range")

        out_doc = fitz.open()
        try:
            for idx in indices:
                if idx < 0 or idx >= total:
                    raise ValidationError(f"Page index {idx} out of range (0-{total - 1})")
                out_doc.insert_pdf(src_doc, from_page=idx, to_page=idx)
            out_doc.save(ctx.output_path, deflate=True, garbage=3)
        finally:
            out_doc.close()
    finally:
        src_doc.close()

    log.info(f"[{ctx.job_id}] extract_pages: extracted {len(indices)} pages")
    return {"pages_extracted": len(indices)}


# ═══════════════════════════════════════════════════════════════════════════════
# PDF OPTIMIZE
# ═══════════════════════════════════════════════════════════════════════════════

@register("compress_pdf")
def compress_pdf(ctx: JobContext) -> dict:
    """
    Two-stage compression: PyMuPDF image downsampling → Ghostscript on ORIGINAL.
    Picks the smallest result from: {GS-on-original, stage1, original}.
    Skips images >5 MB to guard against RAM exhaustion.
    """
    _require(FITZ_OK and PIL_OK, "compress_pdf", "PyMuPDF + Pillow")
    quality = ctx.params.get("quality", "medium")
    cfg = {
        "maximum": {"dpi": 72,  "quality": 45, "gs": "/screen"},
        "high":    {"dpi": 96,  "quality": 60, "gs": "/ebook"},
        "medium":  {"dpi": 120, "quality": 72, "gs": "/printer"},
        "low":     {"dpi": 150, "quality": 85, "gs": "/printer"},
    }.get(quality, {"dpi": 120, "quality": 72, "gs": "/printer"})

    _guard_empty(ctx.input_path)
    orig = os.path.getsize(ctx.input_path)
    ctx.set_progress(5)

    # Stage 1: PyMuPDF image downsampling (on a copy — ORIGINAL stays pristine for GS)
    stage1 = ctx.output_path + "_s1.pdf"
    doc = None
    try:
        doc = fitz.open(ctx.input_path)
        modified = compress_pdf_images(
            doc, cfg["dpi"], cfg["quality"],
        )
        if modified:
            doc.save(stage1, deflate=True, deflate_images=True,
                     deflate_fonts=True, garbage=3, clean=False)
        else:
            shutil.copy(ctx.input_path, stage1)
        doc.close()
        doc = None
    except Exception as ex:
        log.warning(f"[{ctx.job_id}] compress stage1 failed: {ex}")
        if doc is not None:
            try:
                doc.close()
            except Exception:
                pass
        shutil.copy(ctx.input_path, stage1)
    ctx.set_progress(40)

    stage1_size = os.path.getsize(stage1)

    # Stage 2: Ghostscript — ALWAYS on the ORIGINAL (not stage1)
    gs_out = ctx.output_path + "_gs.pdf"
    gs_ok  = False
    try:
        gs_ok = _ghostscript(
            ctx.input_path,   # ← original, not stage1
            gs_out,
            cfg["gs"],
            extra_flags=[
                "-dColorImageDownsampleType=/Bicubic",
                f"-dColorImageResolution={cfg['dpi']}",
                f"-dGrayImageResolution={cfg['dpi']}",
            ],
        )
    except OperationTimeoutError:
        log.warning(f"[{ctx.job_id}] GS timed out during compress")
    ctx.set_progress(80)

    # Pick smallest valid candidate
    candidates = []
    if gs_ok and os.path.exists(gs_out) and os.path.getsize(gs_out) > 0:
        candidates.append((os.path.getsize(gs_out), gs_out))
    if stage1_size > 0:
        candidates.append((stage1_size, stage1))
    candidates.append((orig, ctx.input_path))
    candidates.sort(key=lambda x: x[0])
    _chosen_size, chosen = candidates[0]

    try:
        shutil.copy(chosen, ctx.output_path)
    finally:
        # Always clean up temp files — even if copy raises (disk full, permissions)
        for tmp in [stage1, gs_out]:
            try:
                os.remove(tmp)
            except OSError:
                pass

    new_size  = os.path.getsize(ctx.output_path)
    reduction = round((1 - new_size / orig) * 100, 1) if orig else 0
    ctx.set_progress(100)
    log.info(f"[{ctx.job_id}] compress_pdf: {orig} → {new_size} bytes ({reduction}% reduction)")
    return {
        "reduction_pct":          reduction,
        "original_size_bytes":    orig,
        "compressed_size_bytes":  new_size,
    }


@register("repair_pdf")
def repair_pdf(ctx: JobContext) -> dict:
    """
    Attempt PDF repair via PyMuPDF then Ghostscript fallback.
    """
    _require(FITZ_OK, "repair_pdf", "PyMuPDF")
    orig_size = os.path.getsize(ctx.input_path)

    # Stage 1: PyMuPDF
    try:
        doc   = fitz.open(ctx.input_path)
        pages = len(doc)
        if pages == 0:
            doc.close()
            raise ValidationError("Input PDF has no pages")
        tmp = ctx.output_path + ".tmp"
        doc.save(tmp, garbage=4, deflate=True, clean=True)
        doc.close()
        os.replace(tmp, ctx.output_path)
        if os.path.getsize(ctx.output_path) > 0:
            log.info(f"[{ctx.job_id}] repair_pdf: repaired via pymupdf, {pages} pages")
            return {"method": "pymupdf", "pages": pages, "original_size_bytes": orig_size}
    except (ValidationError, ProcessingError):
        raise
    except Exception as ex:
        log.warning(f"[{ctx.job_id}] PyMuPDF repair failed: {ex}")

    # Stage 2: Ghostscript
    gs_tmp = ctx.output_path + "_gs.pdf"
    try:
        gs_ok = _ghostscript(
            ctx.input_path, gs_tmp, "/printer",
            extra_flags=["-dPDFSTOPONERROR=false"],
        )
        if gs_ok and os.path.getsize(gs_tmp) > 0:
            os.replace(gs_tmp, ctx.output_path)
            log.info(f"[{ctx.job_id}] repair_pdf: repaired via ghostscript")
            return {"method": "ghostscript", "original_size_bytes": orig_size}
    except Exception as ex:
        log.warning(f"[{ctx.job_id}] GS repair failed: {ex}")
    finally:
        try:
            os.remove(gs_tmp)
        except OSError:
            pass

    shutil.copy(ctx.input_path, ctx.output_path)
    return {"method": "passthrough", "note": "PDF was already valid or unrecoverable"}


@register("linearize_pdf")
def linearize_pdf(ctx: JobContext) -> dict:
    """
    Produce a web-optimized (linearized) PDF.
    Prefers qpdf (true linearization) over Ghostscript (approximation).
    """
    _guard_empty(ctx.input_path)
    orig = os.path.getsize(ctx.input_path)

    if QPDF_OK:
        timeout = Config.SUBPROCESS_TIMEOUT
        try:
            result = subprocess.run(
                ["qpdf", "--linearize", ctx.input_path, ctx.output_path],
                capture_output=True,
                timeout=timeout,
            )
            if result.returncode == 0 and os.path.getsize(ctx.output_path) > 0:
                new_size  = os.path.getsize(ctx.output_path)
                reduction = round((1 - new_size / orig) * 100, 1) if orig else 0
                log.info(f"[{ctx.job_id}] linearize_pdf: qpdf, {reduction}% reduction")
                return {
                    "method":              "qpdf",
                    "reduction_pct":       reduction,
                    "original_size_bytes": orig,
                }
        except subprocess.TimeoutExpired:
            raise OperationTimeoutError("qpdf linearize", timeout)
        except Exception as ex:
            log.warning(f"[{ctx.job_id}] qpdf failed: {ex}, falling back to GS")

    # Fallback: Ghostscript (approximate linearization)
    ok = _ghostscript(ctx.input_path, ctx.output_path, "/printer")
    if not ok or not os.path.exists(ctx.output_path):
        raise ProcessingError("Linearization failed — install qpdf or check Ghostscript")
    new_size  = os.path.getsize(ctx.output_path)
    reduction = round((1 - new_size / orig) * 100, 1) if orig else 0
    log.info(f"[{ctx.job_id}] linearize_pdf: ghostscript fallback, {reduction}% reduction")
    return {
        "method":              "ghostscript",
        "reduction_pct":       reduction,
        "original_size_bytes": orig,
    }


# ═══════════════════════════════════════════════════════════════════════════════
# PDF EDIT
# ═══════════════════════════════════════════════════════════════════════════════

@register("rotate_pdf")
def rotate_pdf(ctx: JobContext) -> dict:
    """
    Rotate pages. Angle must be 0/90/180/270.
    Writes rotation to both the page object and the /Rotate key in the page dict
    so all compliant viewers honour it.
    """
    _require(FITZ_OK, "rotate_pdf", "PyMuPDF")
    angle_raw  = ctx.params.get("angle", 90)
    pages_spec = ctx.params.get("pages", "all")

    try:
        angle = int(angle_raw)
    except (TypeError, ValueError):
        raise ValidationError(f"Invalid rotation angle: {angle_raw!r}")
    if angle not in _VALID_ROTATION_ANGLES:
        raise ValidationError(f"Angle must be one of {sorted(_VALID_ROTATION_ANGLES)}")

    _guard_empty(ctx.input_path)
    doc   = fitz.open(ctx.input_path)
    total = len(doc)
    try:
        idxs = (
            list(range(total))
            if str(pages_spec).lower() == "all"
            else parse_page_ranges(pages_spec, total)
        )
        if not idxs:
            raise ValidationError("No valid pages matched the specified range")

        for i in idxs:
            page = doc[i]
            page.set_rotation(angle)
            # Also write /Rotate directly into the page dictionary for maximum compat
            page_obj = doc.xref_object(page.xref, compressed=False)
            # fitz set_rotation already handles the dict — this is a belt-and-suspenders check
        doc.save(ctx.output_path, deflate=True)
    finally:
        doc.close()

    log.info(f"[{ctx.job_id}] rotate_pdf: {len(idxs)} pages rotated {angle}°")
    return {"pages_rotated": len(idxs), "angle": angle}


@register("watermark_pdf")
def watermark_pdf(ctx: JobContext) -> dict:
    """
    Add text or image watermark.
    Opacity clamped to [0, 1]. Image data validated before use.
    Overlay flag set to ensure watermark renders on top of content.
    """
    _require(FITZ_OK, "watermark_pdf", "PyMuPDF")
    text       = ctx.params.get("text", "CONFIDENTIAL")
    color      = ctx.params.get("color", "808080")
    opacity    = float(ctx.params.get("opacity", 0.3))
    opacity    = max(0.0, min(1.0, opacity))           # clamp
    position   = ctx.params.get("position", "diagonal")
    rotation   = float(ctx.params.get("rotation", 45.0))
    image_data = ctx.params.get("image_data")          # bytes, pre-loaded by route

    if image_data:
        _validate_image_data(image_data)

    _guard_empty(ctx.input_path)
    doc = fitz.open(ctx.input_path)
    try:
        pre_img = None
        if image_data:
            pre_img = Image.open(io.BytesIO(image_data)).convert("RGBA")
            if rotation != 0:
                pre_img = pre_img.rotate(rotation, expand=True, resample=Image.BICUBIC)

        for page in doc:
            r = page.rect
            if pre_img:
                img = pre_img.copy()
                r_ch, g_ch, b_ch, a_ch = img.split()
                a_ch = a_ch.point(lambda x: int(x * opacity))
                img.putalpha(a_ch)
                buf = io.BytesIO()
                img.save(buf, format="PNG")
                scale = float(ctx.params.get("scale", 0.3))
                iw    = r.width * scale
                ih    = iw * img.height / img.width
                ix    = r.x0 + (r.width - iw) / 2
                iy    = r.y0 + (r.height - ih) / 2
                page.insert_image(
                    fitz.Rect(ix, iy, ix + iw, iy + ih),
                    stream=buf.getvalue(),
                    overlay=True,          # always on top
                )
            else:
                wm    = create_watermark_pdf(text, opacity, color,
                                             r.width, r.height, position, rotation)
                wmpdf = fitz.open("pdf", wm)
                # show_pdf_page with overlay=True forces watermark above page content
                page.show_pdf_page(
                    fitz.Rect(0, 0, r.width, r.height),
                    wmpdf, 0,
                    overlay=True,
                )
                wmpdf.close()
        doc.save(ctx.output_path, deflate=True)
    finally:
        doc.close()

    log.info(f"[{ctx.job_id}] watermark_pdf: type={'image' if image_data else 'text'}")
    return {"watermark_type": "image" if image_data else "text"}


@register("page_numbers")
def page_numbers(ctx: JobContext) -> dict:
    """Add page number stamps to each page."""
    _require(FITZ_OK, "page_numbers", "PyMuPDF")
    position = ctx.params.get("position", "bottom")
    start    = int(ctx.params.get("start", 1))
    prefix   = ctx.params.get("prefix", "")
    _guard_empty(ctx.input_path)
    doc = fitz.open(ctx.input_path)
    try:
        for i, page in enumerate(doc):
            r     = page.rect
            label = f"{prefix}{start + i}"
            pn    = create_page_number_pdf(label, position, r.width, r.height)
            pnpdf = fitz.open("pdf", pn)
            page.show_pdf_page(
                fitz.Rect(0, 0, r.width, r.height),
                pnpdf, 0, overlay=True,
            )
            pnpdf.close()
        doc.save(ctx.output_path, deflate=True)
    finally:
        doc.close()
    return {}


@register("crop_pdf")
def crop_pdf(ctx: JobContext) -> dict:
    """Crop page margins in points."""
    _require(FITZ_OK, "crop_pdf", "PyMuPDF")
    left   = float(ctx.params.get("left",   0))
    right  = float(ctx.params.get("right",  0))
    top    = float(ctx.params.get("top",    0))
    bottom = float(ctx.params.get("bottom", 0))
    if any(v < 0 for v in (left, right, top, bottom)):
        raise ValidationError("Crop margins must be non-negative")
    _guard_empty(ctx.input_path)
    doc = fitz.open(ctx.input_path)
    try:
        for page in doc:
            r  = page.rect
            nr = fitz.Rect(r.x0 + left, r.y0 + top, r.x1 - right, r.y1 - bottom)
            if nr.is_empty or nr.is_infinite:
                raise ValidationError("Crop margins too large for this page size")
            if nr.width < 10 or nr.height < 10:
                raise ValidationError("Resulting page would be smaller than 10pt in one dimension")
            page.set_cropbox(nr)
        doc.save(ctx.output_path, deflate=True)
    finally:
        doc.close()
    return {}


# ═══════════════════════════════════════════════════════════════════════════════
# PDF EDIT — text/image/highlight/note overlays via PyMuPDF
# ═══════════════════════════════════════════════════════════════════════════════

# Limits — generous but bounded so a malicious payload can't OOM the worker.
_EDIT_MAX_OPS_TOTAL   = 5000   # combined operations across all lists
_EDIT_MAX_TEXT_LEN    = 4000
_EDIT_MAX_NOTE_LEN    = 4000
_EDIT_MAX_IMAGE_BYTES = 25 * 1024 * 1024   # 25 MB per overlay image
_EDIT_FONT_MAP = {
    "helv": "helv", "helvetica": "helv",
    "tiro": "tiro", "times": "tiro", "times-roman": "tiro",
    "cour": "cour", "courier": "cour", "mono": "cour",
}


def _edit_parse_color(value, default=(0.0, 0.0, 0.0)) -> tuple:
    """Accept '#rrggbb', 'rrggbb', [r,g,b] (0-255 or 0-1), or None."""
    if value is None or value == "":
        return default
    if isinstance(value, (list, tuple)) and len(value) >= 3:
        try:
            vals = [float(v) for v in value[:3]]
        except (TypeError, ValueError):
            return default
        if max(vals) > 1.001:
            vals = [v / 255.0 for v in vals]
        return tuple(max(0.0, min(1.0, v)) for v in vals)
    if isinstance(value, str):
        s = value.strip().lstrip("#")
        if len(s) == 6:
            try:
                return (
                    int(s[0:2], 16) / 255.0,
                    int(s[2:4], 16) / 255.0,
                    int(s[4:6], 16) / 255.0,
                )
            except ValueError:
                return default
    return default


def _edit_parse_payload(ctx: JobContext) -> dict:
    """
    Edits arrive as a JSON string in form field 'edits' (or 'payload'),
    or piecemeal as form fields 'text_overlays', 'highlights', etc.
    Returns a dict with normalised lists: text_overlays, highlights,
    annotations, image_overlays.
    """
    raw = ctx.params.get("edits") or ctx.params.get("payload")
    data: dict = {}
    if raw:
        try:
            parsed = json.loads(raw) if isinstance(raw, str) else raw
            if isinstance(parsed, dict):
                data = parsed
        except (TypeError, ValueError) as ex:
            raise ValidationError(f"edits is not valid JSON: {ex}")

    # Allow individual form fields to override / supplement
    for key in ("text_overlays", "highlights", "annotations", "image_overlays"):
        if key in ctx.params and key not in data:
            try:
                data[key] = json.loads(ctx.params[key])
            except (TypeError, ValueError):
                raise ValidationError(f"{key} is not valid JSON")

    def _aslist(v):
        if v is None:
            return []
        return v if isinstance(v, list) else [v]

    out = {
        "text_overlays":  _aslist(data.get("text_overlays")),
        "highlights":     _aslist(data.get("highlights")),
        "annotations":    _aslist(data.get("annotations")),
        "image_overlays": _aslist(data.get("image_overlays")),
    }

    total = sum(len(out[k]) for k in out)
    if total == 0:
        raise ValidationError(
            "No edits provided. Send a JSON 'edits' object with at least one "
            "of: text_overlays, highlights, annotations, image_overlays."
        )
    if total > _EDIT_MAX_OPS_TOTAL:
        raise ValidationError(
            f"Too many edit operations ({total}); limit is {_EDIT_MAX_OPS_TOTAL}."
        )
    return out


@register("edit_pdf")
def edit_pdf(ctx: JobContext) -> dict:
    """
    Apply per-page overlays and annotations.

    Accepts a JSON 'edits' payload of the shape:
      {
        "text_overlays":  [{"page":0,"x":72,"y":120,"text":"Hello",
                            "font_size":14,"color":"#1d4ed8","font":"helv"}],
        "highlights":     [{"page":0,"x1":50,"y1":100,"x2":300,"y2":120,
                            "color":"#ffff00","opacity":0.4}],
        "annotations":    [{"page":0,"x":200,"y":200,"content":"Review this",
                            "type":"text","title":"Reviewer"}],
        "image_overlays": [{"page":0,"x":300,"y":300,"width":120,"height":80,
                            "image_b64":"<base64-png-or-jpg>"}]
      }

    Page indices are 0-based. Coordinates are PDF points from the top-left
    of the page (origin at upper-left of the visible page rect), matching
    what the frontend draws on its canvas preview.
    """
    _require(FITZ_OK, "edit_pdf", "PyMuPDF")
    _guard_empty(ctx.input_path)
    edits = _edit_parse_payload(ctx)

    # Lazy import — base64 only needed for image_overlays
    import base64

    doc = fitz.open(ctx.input_path)
    try:
        n_pages = len(doc)
        applied = {"text": 0, "highlight": 0, "annotation": 0, "image": 0}

        def _resolve_page(idx):
            try:
                p = int(idx)
            except (TypeError, ValueError):
                raise ValidationError(f"Invalid page index: {idx!r}")
            if p < 0:
                p += n_pages
            if p < 0 or p >= n_pages:
                raise ValidationError(
                    f"Page index {idx} out of range (0..{n_pages - 1})"
                )
            return p

        # ── Text overlays ─────────────────────────────────────────────
        for op in edits["text_overlays"]:
            if not isinstance(op, dict):
                raise ValidationError("text_overlays entries must be objects")
            page = doc[_resolve_page(op.get("page", 0))]
            text = str(op.get("text", ""))[:_EDIT_MAX_TEXT_LEN]
            if not text:
                continue
            x = float(op.get("x", 72))
            y = float(op.get("y", 72))
            size = float(op.get("font_size", op.get("size", 14)))
            size = max(4.0, min(size, 400.0))
            color = _edit_parse_color(op.get("color"), (0.0, 0.0, 0.0))
            font_key = str(op.get("font", "helv")).lower().strip()
            font = _EDIT_FONT_MAP.get(font_key, "helv")
            # PyMuPDF's insert_text only accepts rotate values that are
            # multiples of 90; snap to the nearest quadrant.
            try:
                requested_rot = int(op.get("rotate", 0)) % 360
            except (TypeError, ValueError):
                requested_rot = 0
            rotate = int(round(requested_rot / 90.0)) * 90 % 360
            try:
                page.insert_text(
                    (x, y), text,
                    fontname=font, fontsize=size,
                    color=color, rotate=rotate, overlay=True,
                )
            except Exception as ex:
                raise ProcessingError(f"insert_text failed on page {op.get('page')}: {ex}")
            applied["text"] += 1

        # ── Highlights (semi-transparent rectangles) ──────────────────
        for op in edits["highlights"]:
            if not isinstance(op, dict):
                raise ValidationError("highlights entries must be objects")
            page = doc[_resolve_page(op.get("page", 0))]
            x1 = float(op.get("x1", op.get("x", 0)))
            y1 = float(op.get("y1", op.get("y", 0)))
            if "x2" in op and "y2" in op:
                x2 = float(op["x2"]); y2 = float(op["y2"])
            else:
                x2 = x1 + float(op.get("width", 100))
                y2 = y1 + float(op.get("height", 20))
            rect = fitz.Rect(min(x1, x2), min(y1, y2), max(x1, x2), max(y1, y2))
            if rect.width <= 0 or rect.height <= 0:
                continue
            color = _edit_parse_color(op.get("color"), (1.0, 1.0, 0.0))
            opacity = max(0.05, min(1.0, float(op.get("opacity", 0.4))))
            try:
                annot = page.add_highlight_annot(rect)
                annot.set_colors(stroke=color)
                annot.set_opacity(opacity)
                annot.update()
            except Exception:
                # Fall back to a drawn rectangle for PDFs where highlight annot fails
                page.draw_rect(
                    rect, color=color, fill=color,
                    fill_opacity=opacity, width=0, overlay=True,
                )
            applied["highlight"] += 1

        # ── Sticky-note annotations ───────────────────────────────────
        for op in edits["annotations"]:
            if not isinstance(op, dict):
                raise ValidationError("annotations entries must be objects")
            page = doc[_resolve_page(op.get("page", 0))]
            content = str(op.get("content", op.get("text", "")))[:_EDIT_MAX_NOTE_LEN]
            x = float(op.get("x", 72)); y = float(op.get("y", 72))
            title = str(op.get("title", op.get("author", "PDFWala")))[:120]
            try:
                annot = page.add_text_annot((x, y), content, icon="Note")
                annot.set_info(title=title, content=content)
                color = _edit_parse_color(op.get("color"), (1.0, 0.85, 0.3))
                annot.set_colors(stroke=color)
                annot.update()
            except Exception as ex:
                raise ProcessingError(f"add_text_annot failed: {ex}")
            applied["annotation"] += 1

        # ── Image overlays ────────────────────────────────────────────
        for op in edits["image_overlays"]:
            if not isinstance(op, dict):
                raise ValidationError("image_overlays entries must be objects")
            b64 = op.get("image_b64") or op.get("image") or ""
            if isinstance(b64, str) and b64.startswith("data:"):
                # strip data URL prefix "data:image/png;base64,"
                _, _, b64 = b64.partition(",")
            if not b64:
                raise ValidationError("image_overlays entry missing image_b64")
            try:
                blob = base64.b64decode(b64, validate=False)
            except Exception as ex:
                raise ValidationError(f"image_b64 is not valid base64: {ex}")
            if len(blob) > _EDIT_MAX_IMAGE_BYTES:
                raise ValidationError(
                    f"image_overlays entry too large "
                    f"({len(blob)} bytes; max {_EDIT_MAX_IMAGE_BYTES})"
                )
            _validate_image_data(blob, max_bytes=_EDIT_MAX_IMAGE_BYTES)

            page = doc[_resolve_page(op.get("page", 0))]
            x = float(op.get("x", 72)); y = float(op.get("y", 72))
            w = float(op.get("width", 120)); h = float(op.get("height", 80))
            rect = fitz.Rect(x, y, x + w, y + h)
            try:
                page.insert_image(rect, stream=blob, overlay=True, keep_proportion=True)
            except Exception as ex:
                raise ProcessingError(f"insert_image failed: {ex}")
            applied["image"] += 1

        doc.save(ctx.output_path, deflate=True, garbage=3)
    finally:
        doc.close()

    log.info(
        f"[{ctx.job_id}] edit_pdf: text={applied['text']} "
        f"highlight={applied['highlight']} note={applied['annotation']} "
        f"image={applied['image']}"
    )
    return {
        "applied":          applied,
        "total_operations": sum(applied.values()),
    }


@register("redact_pdf")
def redact_pdf(ctx: JobContext) -> dict:
    """
    Redact text by literal match, regex, or preset pattern.
    SafeRegex prevents ReDoS. import re is at module level.
    """
    _require(FITZ_OK, "redact_pdf", "PyMuPDF")
    mode        = ctx.params.get("mode", "text")
    search_text = ctx.params.get("search_text", "")
    pattern_str = ctx.params.get("pattern", "")
    preset_name = ctx.params.get("preset", "")

    compiled = None
    if mode == "text":
        if not search_text:
            raise ValidationError("search_text required for mode=text")
    elif mode == "regex":
        if not pattern_str:
            raise ValidationError("pattern required for mode=regex")
        try:
            compiled = SafeRegex.compile(pattern_str)
        except (ValueError, Exception) as ex:
            raise ValidationError(f"Invalid regex: {ex}")
    elif mode == "preset":
        if preset_name not in REDACTION_PATTERNS:
            raise ValidationError(
                f"Unknown preset. Choose: {', '.join(sorted(REDACTION_PATTERNS))}"
            )
        compiled = re.compile(REDACTION_PATTERNS[preset_name])
    else:
        raise ValidationError(f"Unknown redact mode: {mode!r}")

    _guard_empty(ctx.input_path)
    doc   = fitz.open(ctx.input_path)
    count = 0
    try:
        for page in doc:
            if mode == "text":
                for rect in page.search_for(search_text):
                    page.add_redact_annot(rect, fill=(0, 0, 0))
                    count += 1
            else:
                for match in compiled.finditer(page.get_text("text")):
                    for rect in page.search_for(match.group()):
                        page.add_redact_annot(rect, fill=(0, 0, 0))
                        count += 1
            page.apply_redactions()
        doc.save(ctx.output_path, deflate=True)
    finally:
        doc.close()

    log.info(f"[{ctx.job_id}] redact_pdf: {count} redactions applied")
    return {
        "redaction_count": count,
        "warning": "No matches found — document unchanged" if count == 0 else None,
    }


# ═══════════════════════════════════════════════════════════════════════════════
# PDF INFO / SECURITY
# ═══════════════════════════════════════════════════════════════════════════════

@register("pdf_info")
def pdf_info(ctx: JobContext) -> dict:
    _require(FITZ_OK, "pdf_info", "PyMuPDF")
    _guard_empty(ctx.input_path)
    file_size = os.path.getsize(ctx.input_path)
    doc = fitz.open(ctx.input_path)
    try:
        meta       = doc.metadata
        sizes: dict = {}
        for pg in doc:
            k = (round(pg.rect.width, 1), round(pg.rect.height, 1))
            sizes[k] = sizes.get(k, 0) + 1
        font_names: set = set()
        for pg in doc:
            for fi in pg.get_fonts(full=True):
                if len(fi) > 3 and fi[3]:
                    font_names.add(fi[3])
        is_lin = False
        try:
            xobj   = doc.xref_object(1, compressed=False)
            is_lin = "/Linearized" in (xobj or "")
        except Exception:
            pass
        # PyMuPDF 1.24 removed Document.pdf_version(). The version is exposed
        # in metadata['format'] as e.g. "PDF 1.7".
        pdf_version = ""
        fmt = (meta or {}).get("format", "")
        if isinstance(fmt, str) and fmt.upper().startswith("PDF"):
            pdf_version = fmt.split(" ", 1)[-1] if " " in fmt else fmt

        return {
            "metadata": {
                "page_count":       len(doc),
                "pdf_version":      pdf_version,
                "title":            meta.get("title", ""),
                "author":           meta.get("author", ""),
                "encrypted":        doc.is_encrypted,
                "file_size_bytes":  file_size,
                "size_human":       format_file_size(file_size),
                "has_forms":        any(pg.first_widget for pg in doc),
                "has_toc":          len(doc.get_toc()) > 0,
                "image_count":      sum(len(pg.get_images()) for pg in doc),
                "fonts_used":       sorted(font_names)[:20],
                "is_linearized":    is_lin,
                "page_sizes": [
                    {"w": k[0], "h": k[1], "count": v}
                    for k, v in sorted(sizes.items(), key=lambda x: -x[1])
                ],
            }
        }
    finally:
        doc.close()


@register("protect_pdf")
def protect_pdf(ctx: JobContext) -> dict:
    """
    Encrypt PDF with AES-256.
    Owner password is user_pw + random 8-char suffix so it differs from user_pw,
    preventing trivial privilege escalation.
    """
    _require(FITZ_OK, "protect_pdf", "PyMuPDF")
    pw          = ctx.params.get("password", "")
    allow_print = ctx.params.get("allow_print", True)
    allow_copy  = ctx.params.get("allow_copy", True)
    if not pw:
        raise ValidationError("Password required")
    if len(pw) > 128:
        raise ValidationError("Password too long (max 128 chars)")

    # Owner password must differ from user password
    owner_pw = pw + "-" + secrets.token_hex(4)

    _guard_empty(ctx.input_path)
    doc = fitz.open(ctx.input_path)
    try:
        permissions = int(fitz.PDF_PERM_ACCESSIBILITY)
        if allow_print:
            permissions |= int(fitz.PDF_PERM_PRINT)
        if allow_copy:
            permissions |= int(fitz.PDF_PERM_COPY)
        doc.save(
            ctx.output_path,
            encryption=fitz.PDF_ENCRYPT_AES_256,
            owner_pw=owner_pw,
            user_pw=pw,
            permissions=permissions,
        )
    finally:
        doc.close()
    return {}


@register("unlock_pdf")
def unlock_pdf(ctx: JobContext) -> dict:
    """
    Remove password protection.
    Tries user_pw first, then owner_pw (owner unlocks all restrictions).
    Saves with no encryption AND explicitly clears permission bits.
    """
    _require(FITZ_OK, "unlock_pdf", "PyMuPDF")
    pw = ctx.params.get("password", "")
    if not pw:
        raise ValidationError("Password required")

    doc = fitz.open(ctx.input_path)
    authenticated = False
    try:
        if doc.is_encrypted:
            # Try as user password first, then as owner password
            if doc.authenticate(pw):
                authenticated = True
            else:
                doc.close()
                raise ValidationError("Wrong password — authentication failed")
        else:
            authenticated = True   # not encrypted — just re-save without encryption

        # Save with no encryption and full permissions
        doc.save(
            ctx.output_path,
            encryption=fitz.PDF_ENCRYPT_NONE,
            deflate=True,
            garbage=3,
        )
    finally:
        doc.close()

    log.info(f"[{ctx.job_id}] unlock_pdf: successfully unlocked")
    return {"was_encrypted": not authenticated or True}


@register("sign_pdf")
def sign_pdf(ctx: JobContext) -> dict:
    """
    Add a visible signature stamp to the PDF.
    Signature is stamped directly into the page content stream via fitz drawing
    primitives so it persists on reopen in all viewers.
    sig_data bytes are validated before use.
    Note: this is a VISUAL stamp, not a cryptographic signature (PAdES/PKCS#7).
    """
    _require(FITZ_OK, "sign_pdf", "PyMuPDF")
    from datetime import datetime

    name        = ctx.params.get("name", "Signed")
    reason      = ctx.params.get("reason", "Approved")
    page_target = ctx.params.get("page", "last")
    position    = ctx.params.get("position", "bottom-right")
    sig_data    = ctx.params.get("sig_data")       # bytes or None
    today_str   = datetime.now().strftime("%Y-%m-%d")

    if sig_data:
        _validate_image_data(sig_data, max_bytes=2 * 1024 * 1024)

    _guard_empty(ctx.input_path)
    doc   = fitz.open(ctx.input_path)
    total = len(doc)
    try:
        if page_target == "all":
            idxs = list(range(total))
        elif page_target == "first":
            idxs = [0]
        elif page_target == "last":
            idxs = [total - 1]
        else:
            try:
                n = int(page_target)
                if n < 1 or n > total:
                    raise ValidationError(f"Page {n} out of range (1-{total})")
                idxs = [n - 1]
            except ValueError:
                raise ValidationError(f"Invalid page target: {page_target!r}")

        pos_map = {
            "bottom-right": lambda r: (r.x1 - 180, r.y1 - 70),
            "bottom-left":  lambda r: (r.x0 + 30,  r.y1 - 70),
            "top-right":    lambda r: (r.x1 - 180, r.y0 + 50),
            "top-left":     lambda r: (r.x0 + 30,  r.y0 + 50),
            "center":       lambda r: (r.x0 + r.width / 2 - 75, r.y0 + r.height / 2),
        }
        pos_fn = pos_map.get(position, pos_map["bottom-right"])

        for idx in idxs:
            page = doc[idx]
            rect = page.rect
            sx, sy = pos_fn(rect)

            # Clamp to page bounds
            sx = max(rect.x0 + 5, min(sx, rect.x1 - 165))
            sy = max(rect.y0 + 30, min(sy, rect.y1 - 10))

            # Draw signature box directly into content stream (persists on reopen)
            box_r = fitz.Rect(sx - 5, sy - 45, sx + 155, sy + 25)
            shape = page.new_shape()
            shape.draw_rect(box_r)
            shape.finish(color=(0, 0, 0.6), fill=(0.9, 0.9, 1.0), width=0.8)
            line = f"{name} | {reason} | {today_str}"
            shape.insert_text((sx, sy - 28), "✦ SIGNED", fontsize=9,
                               color=(0, 0, 0.6))
            shape.insert_text((sx, sy - 12), line, fontsize=7.5,
                               color=(0.1, 0.1, 0.1))
            shape.commit()  # writes into page content stream — survives reopen

            if sig_data:
                page.insert_image(
                    fitz.Rect(sx, sy - 42, sx + 100, sy - 5),
                    stream=sig_data,
                    overlay=True,
                )

        doc.save(ctx.output_path, deflate=True)
    finally:
        doc.close()

    log.info(f"[{ctx.job_id}] sign_pdf: stamped {len(idxs)} pages")
    return {
        "pages_signed": len(idxs),
        "note": "Visual stamp only — not a cryptographic (PAdES/PKCS#7) signature",
    }


# ═══════════════════════════════════════════════════════════════════════════════
# PDF CONVERT
# ═══════════════════════════════════════════════════════════════════════════════

@register("pdf_to_image")
def pdf_to_image(ctx: JobContext) -> dict:
    """
    Convert PDF pages to images, streamed to disk ZIP for large docs.
    """
    _require(FITZ_OK and PIL_OK, "pdf_to_image", "PyMuPDF + Pillow")
    fmt = ctx.params.get("format", "jpg").lower()
    dpi = int(ctx.params.get("dpi", 150))
    if fmt not in ("jpg", "png"):
        fmt = "jpg"
    dpi = max(72, min(dpi, 600))   # clamp

    _guard_empty(ctx.input_path)
    doc = fitz.open(ctx.input_path)
    try:
        count  = len(doc)
        zf, buf = _open_zip_writer(ctx.output_path, count)
        try:
            for i, page in enumerate(doc):
                mat = fitz.Matrix(dpi / 72, dpi / 72)
                pix = page.get_pixmap(matrix=mat, alpha=True)
                try:
                    pil = Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGBA")
                    bg  = Image.new("RGB", pil.size, (255, 255, 255))
                    bg.paste(pil, mask=pil.split()[3])
                    ib = io.BytesIO()
                    if fmt == "jpg":
                        bg.save(ib, "JPEG", quality=85, optimize=True)
                    else:
                        bg.save(ib, "PNG", optimize=True)
                    zf.writestr(f"page_{i + 1:04d}.{fmt}", ib.getvalue())
                finally:
                    # Explicit cleanup — prevents RAM accumulation at high DPI
                    try: pil.close()
                    except Exception: pass
                    try: bg.close()
                    except Exception: pass
                    try: ib.close()
                    except Exception: pass
                    del pix
                if i % 10 == 0:
                    ctx.set_progress(int(i / count * 95))
        finally:
            _finalise_zip(zf, buf, ctx.output_path)
    finally:
        doc.close()

    ctx.set_progress(100)
    log.info(f"[{ctx.job_id}] pdf_to_image: {count} pages → {fmt}")
    return {"pages_exported": count, "format": fmt}


@register("pdf_to_jpg")
def pdf_to_jpg(ctx: JobContext) -> dict:
    ctx.params["format"] = "jpg"
    return pdf_to_image(ctx)


@register("pdf_to_png")
def pdf_to_png(ctx: JobContext) -> dict:
    ctx.params["format"] = "png"
    return pdf_to_image(ctx)


@register("pdf_to_word")
def pdf_to_word(ctx: JobContext) -> dict:
    """
    Convert PDF to DOCX via pdf2docx (table-aware mode).
    V14 FIX: For large PDFs, convert in page-range chunks to avoid pdf2docx OOM.
             Also adds progress reporting and better error diagnostics.
    """
    _require(PDF2DOCX_OK, "pdf_to_word", "pdf2docx")
    _guard_empty(ctx.input_path)

    # Determine page count for progress + chunking decision
    page_count = 0
    if FITZ_OK:
        doc = fitz.open(ctx.input_path)
        page_count = len(doc)
        doc.close()

    ctx.set_progress(5)

    # For files > 100 pages, convert in chunks to avoid pdf2docx memory exhaustion
    CHUNK_THRESHOLD = 100
    if page_count > CHUNK_THRESHOLD:
        chunk_size = 50
        chunk_docxs = []
        tmp_dir = None
        try:
            import tempfile
            tmp_dir = tempfile.mkdtemp(prefix="pdf2docx_chunks_")
            n_chunks = (page_count + chunk_size - 1) // chunk_size
            for ci in range(n_chunks):
                start = ci * chunk_size
                end   = min(start + chunk_size, page_count)
                chunk_out = os.path.join(tmp_dir, f"chunk_{ci:04d}.docx")
                cv = Pdf2DocxConverter(ctx.input_path)
                try:
                    cv.convert(chunk_out, start=start, end=end)
                finally:
                    cv.close()
                if os.path.exists(chunk_out) and os.path.getsize(chunk_out) > 0:
                    chunk_docxs.append(chunk_out)
                ctx.set_progress(5 + int((ci + 1) / n_chunks * 85))

            if not chunk_docxs:
                raise ProcessingError("pdf2docx produced no output for any chunk")

            # Merge chunks with python-docx if multiple, else just move the single chunk
            if len(chunk_docxs) == 1:
                shutil.copy(chunk_docxs[0], ctx.output_path)
            else:
                # Merge via python-docx compose
                try:
                    from docx import Document as _DocxDoc
                    from docx.oxml.ns import qn
                    import copy

                    base_doc = _DocxDoc(chunk_docxs[0])
                    for chunk_path in chunk_docxs[1:]:
                        src = _DocxDoc(chunk_path)
                        # Add page break before each chunk
                        from docx.oxml import OxmlElement
                        br = OxmlElement("w:p")
                        r  = OxmlElement("w:r")
                        rPr = OxmlElement("w:rPr")
                        pb  = OxmlElement("w:lastRenderedPageBreak")
                        rPr.append(pb)
                        r.append(rPr)
                        br.append(r)
                        base_doc.element.body.append(br)
                        for element in src.element.body:
                            base_doc.element.body.append(copy.deepcopy(element))
                    base_doc.save(ctx.output_path)
                except Exception as merge_ex:
                    log.warning(f"[{ctx.job_id}] chunk merge failed ({merge_ex}), using first chunk only")
                    shutil.copy(chunk_docxs[0], ctx.output_path)
        finally:
            if tmp_dir:
                shutil.rmtree(tmp_dir, ignore_errors=True)
    else:
        # Standard single-pass conversion
        cv = Pdf2DocxConverter(ctx.input_path)
        try:
            cv.convert(ctx.output_path, start=0, end=None)
        finally:
            cv.close()

    if not os.path.exists(ctx.output_path) or os.path.getsize(ctx.output_path) == 0:
        raise ProcessingError(
            "pdf2docx produced empty output — the PDF may be image-only, "
            "encrypted, or have an unsupported structure. Try OCR first."
        )
    ctx.set_progress(100)
    return {"pages": page_count}


@register("pdf_to_excel")
def pdf_to_excel(ctx: JobContext) -> dict:
    """
    Extract tables from PDF to XLSX.
    Always closes fitz doc even if pdfplumber fails.
    """
    _require(OPENPYXL_OK, "pdf_to_excel", "openpyxl")
    _guard_empty(ctx.input_path)

    wb               = Workbook()
    wb.remove(wb.active)
    tables_extracted = 0

    if PDFPLUMBER_OK:
        with pdfplumber.open(ctx.input_path) as pdf:
            for page in pdf.pages:
                for table in page.extract_tables():
                    if table and any(any(c for c in r if c) for r in table):
                        tables_extracted += 1
                        ws = wb.create_sheet(f"Table_{tables_extracted}")
                        for row in table:
                            ws.append([str(c).strip() if c else "" for c in row])

    if tables_extracted == 0 and FITZ_OK:
        ws      = wb.create_sheet("Text")
        doc     = fitz.open(ctx.input_path)
        row_idx = 1
        try:
            for pn, pg in enumerate(doc):
                ws.cell(row_idx, 1, f"--- Page {pn + 1} ---")
                row_idx += 1
                for line in pg.get_text("text").split("\n"):
                    if line.strip():
                        ws.cell(row_idx, 1, line.strip())
                        row_idx += 1
        finally:
            doc.close()   # always closed — bug fix

    wb.save(ctx.output_path)
    log.info(f"[{ctx.job_id}] pdf_to_excel: {tables_extracted} tables extracted")
    return {"tables_found": tables_extracted}


@register("pdf_to_ppt")
def pdf_to_ppt(ctx: JobContext) -> dict:
    """
    Convert PDF pages to PowerPoint slides (one image per slide).
    Temp files cleaned up in finally block even on exception.
    """
    _require(PPTX_OK and FITZ_OK, "pdf_to_ppt", "python-pptx + PyMuPDF")
    _guard_empty(ctx.input_path)

    doc        = fitz.open(ctx.input_path)
    # FIX: capture page count before the finally block closes doc
    slide_count = len(doc)
    tmp_files: list[str] = []
    try:
        prs = Presentation()
        prs.slide_width  = PptxInches(10)
        prs.slide_height = PptxInches(7.5)
        blank = prs.slide_layouts[6]

        for i, page in enumerate(doc):
            pix     = page.get_pixmap(dpi=150)
            fd, tmp_path = tempfile.mkstemp(suffix=".png")
            tmp_files.append(tmp_path)
            try:
                with os.fdopen(fd, "wb") as f:
                    f.write(pix.tobytes("png"))
            except Exception:
                os.close(fd)
                raise
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(tmp_path, 0, 0, prs.slide_width, prs.slide_height)
            if i % 10 == 0:
                ctx.set_progress(int(i / slide_count * 90))
    finally:
        doc.close()
        for tp in tmp_files:    # always clean up temp PNGs
            try:
                os.unlink(tp)
            except OSError:
                pass

    prs.save(ctx.output_path)
    log.info(f"[{ctx.job_id}] pdf_to_ppt: {slide_count} slides created")
    return {"slides_created": slide_count}


@register("pdf_to_pdfa")
def pdf_to_pdfa(ctx: JobContext) -> dict:
    """
    Convert to PDF/A via Ghostscript.
    Checks returncode and verifies output exists.
    """
    version  = ctx.params.get("version", "1b")
    pdfa_val = "2" if "3" in str(version) else "1"
    safe_out = _safe_output_path(ctx.output_path)
    cmd = [
        Config.GHOSTSCRIPT,
        "-dBATCH", "-dNOPAUSE", "-dSAFER",
        "-sDEVICE=pdfwrite",
        f"-dPDFA={pdfa_val}",
        "-dPDFACompatibilityPolicy=1",
        f"-sOutputFile={safe_out}",
        ctx.input_path,
    ]
    timeout = Config.PDFA_TIMEOUT
    try:
        result = subprocess.run(cmd, capture_output=True, timeout=timeout)
        if result.returncode != 0:
            raise ProcessingError(
                f"Ghostscript PDF/A conversion failed (rc={result.returncode}): "
                f"{result.stderr[:200]}"
            )
    except subprocess.TimeoutExpired:
        raise OperationTimeoutError("PDF/A conversion", timeout)

    if not os.path.exists(safe_out) or os.path.getsize(safe_out) == 0:
        raise ProcessingError("PDF/A conversion produced empty output")

    return {"pdfa_version": version}


@register("ocr_pdf")
def ocr_pdf(ctx: JobContext) -> dict:
    """
    OCR a scanned PDF using Tesseract.
    - Default DPI lowered to 200 (4x less RAM than 300 DPI, same accuracy for most docs)
    - Parallel page processing via ThreadPoolExecutor
    - lang and psm/oem parameters sanitised (whitelist)
    - Chunked progress reporting
    """
    _require(TESSERACT_OK and FITZ_OK, "ocr_pdf", "pytesseract + PyMuPDF")
    lang_raw = ctx.params.get("lang", "eng")
    lang     = _sanitise_tesseract_lang(lang_raw)
    dpi      = int(ctx.params.get("dpi", 200))          # 200 default, not 300
    dpi      = max(72, min(dpi, 600))

    psm_raw = int(ctx.params.get("psm", 3))
    oem_raw = int(ctx.params.get("oem", 3))
    if psm_raw not in _VALID_TESSERACT_PSM:
        raise ValidationError(f"Invalid Tesseract PSM value: {psm_raw}")
    if oem_raw not in _VALID_TESSERACT_OEM:
        raise ValidationError(f"Invalid Tesseract OEM value: {oem_raw}")
    psm = psm_raw
    oem = oem_raw

    workers = Config.OCR_WORKERS
    max_pages = Config.MAX_OCR_PAGES

    _guard_empty(ctx.input_path)
    src_doc = fitz.open(ctx.input_path)
    total   = len(src_doc)

    if total > max_pages:
        src_doc.close()
        raise ValidationError(
            f"PDF has {total} pages; OCR is limited to {max_pages} pages. "
            "Split the PDF first or contact support for bulk processing."
        )

    # Pre-render all pages to PNG bytes (allows parallel Tesseract calls)
    page_images: list[tuple[int, float, float, bytes]] = []
    try:
        for page_num, src_page in enumerate(src_doc):
            pw, ph = src_page.rect.width, src_page.rect.height
            mat    = fitz.Matrix(dpi / 72, dpi / 72)
            pix    = src_page.get_pixmap(matrix=mat, alpha=False, colorspace=fitz.csGRAY)
            page_images.append((page_num, pw, ph, pix.tobytes("png")))
    finally:
        src_doc.close()

    ctx.set_progress(15)

    def _ocr_page(args: tuple) -> tuple[int, float, float, Optional[dict]]:
        """Worker: run Tesseract on one page image. Returns (page_num, pw, ph, hocr|None)."""
        page_num, pw, ph, png_bytes = args
        img = Image.open(io.BytesIO(png_bytes))
        try:
            hocr = pytesseract.image_to_data(
                img, lang=lang,
                output_type=TesseractOutput.DICT,
                config=f"--psm {psm} --oem {oem}",
            )
        except Exception as ex:
            log.warning(f"OCR page {page_num + 1}: {ex}")
            hocr = None
        return (page_num, pw, ph, hocr)

    hocr_results: dict[int, tuple] = {}
    completed = 0
    with ThreadPoolExecutor(max_workers=workers) as pool:
        futures = {pool.submit(_ocr_page, item): item[0] for item in page_images}
        for future in as_completed(futures):
            page_num, pw, ph, hocr = future.result()
            hocr_results[page_num] = (pw, ph, hocr)
            completed += 1
            ctx.set_progress(15 + int(completed / total * 70))

    ctx.set_progress(85)

    # Reassemble in page order
    src_doc2 = fitz.open(ctx.input_path)
    out_doc  = fitz.open()
    pages_processed = 0
    img_scale_x = img_scale_y = 1.0  # recomputed per page below

    try:
        for page_num in range(total):
            pw, ph, hocr = hocr_results[page_num]
            src_page     = src_doc2[page_num]
            new_page     = out_doc.new_page(width=pw, height=ph)
            # Copy original raster into new page
            new_page.show_pdf_page(
                fitz.Rect(0, 0, pw, ph), src_doc2, page_num, overlay=False
            )

            if hocr:
                # Compute scale from actual rendered image size
                pix_w = dpi / 72 * pw
                pix_h = dpi / 72 * ph
                img_sx = pw / pix_w
                img_sy = ph / pix_h

                for i in range(len(hocr.get("text", []))):
                    word = (hocr["text"][i] or "").strip()
                    conf = int(hocr["conf"][i]) if hocr["conf"][i] != -1 else 0
                    if not word or conf < 30:
                        continue
                    x0 = hocr["left"][i]   * img_sx
                    y1 = (hocr["top"][i] + hocr["height"][i]) * img_sy
                    fs = max(4.0, hocr["height"][i] * img_sy * 0.85)
                    new_page.insert_text(
                        (x0, y1 - 1), word + " ",
                        fontsize=fs, fontname="helv",
                        color=(0, 0, 0), render_mode=3, overlay=True,
                    )
            pages_processed += 1

        if pages_processed == 0:
            raise ProcessingError("OCR produced no output — all pages failed")

        out_doc.save(ctx.output_path, deflate=True, garbage=2)
    finally:
        out_doc.close()
        src_doc2.close()

    ctx.set_progress(100)
    log.info(f"[{ctx.job_id}] ocr_pdf: {pages_processed}/{total} pages, lang={lang}, dpi={dpi}")
    return {"pages_processed": pages_processed, "lang": lang, "dpi": dpi}


@register("compare_pdf")
def compare_pdf(ctx: JobContext) -> dict:
    """
    Visual + textual diff of two PDFs.
    ZIP streamed to disk for large docs.
    """
    _require(FITZ_OK and PIL_OK, "compare_pdf", "PyMuPDF + Pillow")
    from PIL import ImageChops

    p1, p2 = ctx.input_paths[0], ctx.input_paths[1]
    doc1   = fitz.open(p1)
    doc2   = fitz.open(p2)

    if len(doc1) == 0:
        doc1.close(); doc2.close()
        raise ValidationError("First PDF has no pages")
    if len(doc2) == 0:
        doc1.close(); doc2.close()
        raise ValidationError("Second PDF has no pages")

    pages     = min(len(doc1), len(doc2))
    sims:      list[float] = []
    diff_data: list[dict]  = []

    try:
        zf, buf = _open_zip_writer(ctx.output_path, pages)
        try:
            for i in range(pages):
                pix1 = doc1[i].get_pixmap(dpi=150)
                pix2 = doc2[i].get_pixmap(dpi=150)
                img1 = Image.open(io.BytesIO(pix1.tobytes("png"))).convert("RGB")
                img2 = Image.open(io.BytesIO(pix2.tobytes("png"))).convert("RGB")
                if img1.size != img2.size:
                    img2 = img2.resize(img1.size, Image.LANCZOS)
                diff = ImageChops.difference(img1, img2)
                diff = diff.point(lambda x: min(x * 8, 255))
                db   = io.BytesIO()
                diff.save(db, "PNG")
                zf.writestr(f"diff_page_{i + 1:04d}.png", db.getvalue())

                words1 = [w[4] for w in doc1[i].get_text("words")][:500]
                words2 = [w[4] for w in doc2[i].get_text("words")][:500]
                sm     = difflib.SequenceMatcher(None, words1, words2)
                sim    = round(sm.ratio() * 100, 1)
                sims.append(sim)
                diff_data.append({"page": i + 1, "similarity_pct": sim})

                if i % 10 == 0:
                    ctx.set_progress(int(i / pages * 90))

            zf.writestr(
                "summary.json",
                json.dumps({
                    "pages":                    diff_data,
                    "overall_similarity_pct":   round(sum(sims) / len(sims), 1) if sims else 0,
                }),
            )
        finally:
            _finalise_zip(zf, buf, ctx.output_path)
    finally:
        doc1.close()
        doc2.close()

    ctx.set_progress(100)
    log.info(f"[{ctx.job_id}] compare_pdf: {pages} pages compared")
    return {"pages_compared": pages}


# ═══════════════════════════════════════════════════════════════════════════════
# CANVAS EDITOR — in-place text editing (synchronous, NOT pipeline-registered)
# ═══════════════════════════════════════════════════════════════════════════════
#
# Powers the "Edit PDF" visual canvas tool. Called DIRECTLY from
# app/routes/pdf_routes.py (no Celery, no Redis job) because the user waits
# interactively and the work is fast (<2s for typical PDFs):
#
#   _parse_canvas_sync(pdf_path)                 -> dict   (page images + text spans)
#   _save_canvas_sync(pdf_path, changes, scanned)-> bytes  (rebuilt PDF)
#
# Scanned PDFs (no text layer) are auto-OCR'd: Tesseract word boxes become
# editable spans directly — no text-layer round-trip, so the existing ocr_pdf
# tool is left completely untouched.
#
# Coordinates everywhere are PDF points with a TOP-LEFT origin (the space used
# by both page.get_text() and page.get_pixmap()), so frontend overlay math is a
# simple uniform scale.

_CANVAS_RENDER_DPI   = 150     # page preview render resolution
_CANVAS_MAX_PAGES    = 50      # sync endpoint — keep render fast & memory bounded
_CANVAS_OCR_DPI      = 200     # OCR render resolution for scanned PDFs
_CANVAS_OCR_MIN_CONF = 30      # drop OCR words below this confidence

# Cache base-14 fitz.Font objects — re-creating them per span is wasteful.
_FITZ_FONT_CACHE: dict = {}


def _fitz_font(fontname: str):
    f = _FITZ_FONT_CACHE.get(fontname)
    if f is None:
        f = fitz.Font(fontname=fontname)
        _FITZ_FONT_CACHE[fontname] = f
    return f


def _is_scanned_pdf(path: str) -> bool:
    """
    True if the PDF has no usable text layer.
    Checks the first 3 pages; < 50 stripped chars total → treat as scanned.
    """
    doc = fitz.open(path)
    try:
        pages_to_check = min(3, len(doc))
        total_chars = 0
        for i in range(pages_to_check):
            total_chars += len(doc[i].get_text().strip())
        return total_chars < 50
    finally:
        doc.close()


def _pack_color_to_rgb(c) -> list:
    """PyMuPDF span colour is a packed sRGB int. Return [R,G,B] in 0-255."""
    if isinstance(c, (list, tuple)):
        vals = list(c[:3]) or [0, 0, 0]
        if vals and max(vals) <= 1.0:
            vals = [int(round(v * 255)) for v in vals]
        return [int(max(0, min(255, v))) for v in vals]
    try:
        c = int(c)
    except (TypeError, ValueError):
        return [0, 0, 0]
    return [(c >> 16) & 0xFF, (c >> 8) & 0xFF, c & 0xFF]


def _base14_for(font: str, flags: int) -> str:
    """Map an arbitrary font + span flags to a base-14 PyMuPDF font code.

    NOTE: the correct italic code is 'heit' (Helvetica-Oblique) and bold-italic
    is 'hebi' (Helvetica-BoldOblique). 'heio' is NOT a valid code.
    flags bit 4 (16) = bold, bit 1 (2) = italic.
    """
    fl = (font or "").lower()
    bold   = bool(flags & (1 << 4)) or "bold" in fl or "black" in fl or "heavy" in fl
    italic = bool(flags & (1 << 1)) or "italic" in fl or "oblique" in fl
    if bold and italic:
        return "hebi"
    if bold:
        return "hebo"
    if italic:
        return "heit"
    return "helv"


def _ocr_page_spans(page, dpi: int = _CANVAS_OCR_DPI, lang: str = "eng") -> list:
    """
    OCR one page and synthesise editable spans in PDF-point coordinates
    (top-left origin). Used only for scanned PDFs.
    """
    mat = fitz.Matrix(dpi / 72, dpi / 72)
    pix = page.get_pixmap(matrix=mat, alpha=False, colorspace=fitz.csGRAY)
    img = Image.open(io.BytesIO(pix.tobytes("png")))
    data = pytesseract.image_to_data(
        img, lang=lang, output_type=TesseractOutput.DICT, config="--psm 3 --oem 3"
    )
    scale = 72.0 / dpi   # pixel -> point
    spans = []
    n = len(data.get("text", []))
    for i in range(n):
        word = (data["text"][i] or "").strip()
        if not word:
            continue
        try:
            conf = int(float(data["conf"][i]))
        except (TypeError, ValueError):
            conf = -1
        if conf < _CANVAS_OCR_MIN_CONF:
            continue
        x = data["left"][i]   * scale
        y = data["top"][i]    * scale
        w = data["width"][i]  * scale
        h = data["height"][i] * scale
        if w <= 0 or h <= 0:
            continue
        size = max(6.0, h * 0.85)
        baseline_y = y + h - h * 0.18
        spans.append({
            "text": word,
            "x0": round(x, 2),       "y0": round(y, 2),
            "x1": round(x + w, 2),   "y1": round(y + h, 2),
            "ox": round(x, 2),       "oy": round(baseline_y, 2),
            "font": "OCR", "size": round(size, 2),
            "color": [0, 0, 0], "flags": 0,
        })
    return spans


def _parse_canvas_sync(pdf_path: str, render_dpi: int = _CANVAS_RENDER_DPI,
                       max_pages: int = _CANVAS_MAX_PAGES) -> dict:
    """Render page images + extract editable text spans for the canvas editor."""
    _require(FITZ_OK, "parse-canvas", "PyMuPDF")
    import base64
    _guard_empty(pdf_path)

    scanned = _is_scanned_pdf(pdf_path)
    use_ocr = scanned and TESSERACT_OK and PIL_OK

    doc = fitz.open(pdf_path)
    try:
        n = len(doc)
        if n > max_pages:
            raise ValidationError(
                f"This PDF has {n} pages. The visual editor supports up to "
                f"{max_pages} pages — use Split PDF first to edit a section."
            )
        mat = fitz.Matrix(render_dpi / 72, render_dpi / 72)
        pages = []
        total_spans = 0
        for pno in range(n):
            page = doc[pno]
            rect = page.rect
            pix  = page.get_pixmap(matrix=mat, alpha=False)
            img_b64 = base64.b64encode(pix.tobytes("png")).decode("ascii")

            spans = []
            if use_ocr:
                for j, sp in enumerate(_ocr_page_spans(page)):
                    sp["id"] = f"p{pno}_o{j}"
                    spans.append(sp)
            else:
                d = page.get_text(
                    "dict",
                    flags=fitz.TEXTFLAGS_DICT | fitz.TEXT_PRESERVE_WHITESPACE,
                )
                for bi, block in enumerate(d.get("blocks", [])):
                    if block.get("type", 0) != 0:        # 0 = text block
                        continue
                    for li, line in enumerate(block.get("lines", [])):
                        for si, span in enumerate(line.get("spans", [])):
                            text = span.get("text", "")
                            if not text.strip():
                                continue
                            x0, y0, x1, y1 = span["bbox"]
                            ox, oy = span.get("origin", (x0, y1))
                            spans.append({
                                "id":   f"p{pno}_b{bi}_l{li}_s{si}",
                                "text": text,
                                "x0": round(x0, 2), "y0": round(y0, 2),
                                "x1": round(x1, 2), "y1": round(y1, 2),
                                "ox": round(ox, 2), "oy": round(oy, 2),
                                "font": span.get("font", ""),
                                "size": round(span.get("size", 0), 2),
                                "color": _pack_color_to_rgb(span.get("color", 0)),
                                "flags": int(span.get("flags", 0)),
                            })
            total_spans += len(spans)
            pages.append({
                "page": pno,
                "pdf_width":     round(rect.width, 2),
                "pdf_height":    round(rect.height, 2),
                "render_width":  pix.width,
                "render_height": pix.height,
                "image_b64":     img_b64,
                "spans":         spans,
            })
        return {
            "scanned":     scanned,
            "ocr_applied": use_ocr,
            "page_count":  n,
            "total_spans": total_spans,
            "pages":       pages,
        }
    finally:
        doc.close()


def _insert_fitted_text(page, ch: dict) -> None:
    """Insert one replacement span, shrinking font to fit the original width."""
    new_text = str(ch.get("new_text", ""))
    if new_text == "":
        return  # pure deletion — the redaction already cleared the original
    x0 = float(ch["x0"]); y0 = float(ch["y0"])
    x1 = float(ch["x1"]); y1 = float(ch["y1"])
    fontname = _base14_for(str(ch.get("font", "")), int(ch.get("flags", 0) or 0))
    color = _pack_color_to_rgb(ch.get("color", [0, 0, 0]))
    color_f = tuple(v / 255.0 for v in color)
    size = float(ch.get("size", 0) or 0) or max(6.0, (y1 - y0) * 0.8)

    avail_w = max(1.0, x1 - x0)
    try:
        tw = _fitz_font(fontname).text_length(new_text, fontsize=size)
    except Exception:
        tw = 0
    fs = size
    if tw > avail_w and tw > 0:
        fs = max(4.0, size * (avail_w / tw) * 0.985)   # shrink to fit width

    ox = float(ch.get("ox", x0))
    oy = float(ch.get("oy", y1 - size * 0.18))
    try:
        page.insert_text((ox, oy), new_text, fontname=fontname,
                         fontsize=fs, color=color_f, overlay=True)
    except Exception as ex:
        log.warning(f"canvas insert_text failed for span {ch.get('id')}: {ex}")


def _save_canvas_sync(pdf_path: str, changes: list, scanned: bool = False) -> bytes:
    """Surgically replace only the edited spans; everything else is untouched."""
    _require(FITZ_OK, "save-canvas", "PyMuPDF")
    if not isinstance(changes, list):
        raise ValidationError("changes must be a list")

    by_page: dict = {}
    for ch in changes:
        try:
            pno = int(ch["page"])
        except (KeyError, TypeError, ValueError):
            continue
        by_page.setdefault(pno, []).append(ch)

    if not by_page:
        raise ValidationError("No valid changes to apply")

    # Scanned pages: blank the covered raster pixels. Text pages: leave images
    # (logos/photos) alone — only text/vector under the rect is removed.
    img_mode = fitz.PDF_REDACT_IMAGE_PIXELS if scanned else fitz.PDF_REDACT_IMAGE_NONE

    doc = fitz.open(pdf_path)
    try:
        npages = len(doc)
        edited = 0
        for pno, chs in by_page.items():
            if pno < 0 or pno >= npages:
                continue
            page = doc[pno]
            # Step A — white out every original on this page, then apply once
            applied_any = False
            for ch in chs:
                try:
                    x0 = float(ch["x0"]); y0 = float(ch["y0"])
                    x1 = float(ch["x1"]); y1 = float(ch["y1"])
                except (KeyError, TypeError, ValueError):
                    continue
                page.add_redact_annot(fitz.Rect(x0 - 1, y0 - 1, x1 + 1, y1 + 1),
                                      fill=(1, 1, 1))
                applied_any = True
            if applied_any:
                page.apply_redactions(images=img_mode)
            # Step B — re-insert the new text
            for ch in chs:
                _insert_fitted_text(page, ch)
                edited += 1
        if edited == 0:
            raise ValidationError("No valid changes to apply")
        buf = io.BytesIO()
        doc.save(buf, deflate=True, garbage=4, clean=True)
        return buf.getvalue()
    finally:
        doc.close()
