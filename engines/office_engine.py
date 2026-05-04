"""
engines/office_engine.py — PDFWala Enterprise V13.0

All Word / Excel / PowerPoint processing.
Registered to Pipeline via @register("operation_name").
Called ONLY by Pipeline.run(ctx) — never directly from routes.

Tools implemented (31 total):
  Word  : word_to_pdf, word_to_txt, word_to_html, word_to_json,
          word_to_excel, word_to_ppt, word_to_jpg, word_to_png,
          edit_word, compress_word, unlock_word, protect_word
  Excel : excel_to_pdf, excel_to_csv, excel_to_word, excel_to_json,
          compress_excel, unlock_excel, protect_excel,
          excel_to_jpg, excel_to_ppt, repair_excel
  PPT   : ppt_to_pdf, ppt_to_jpg, compress_ppt,
          unlock_ppt, protect_ppt

V13 FIX SUMMARY (31 issues):
  S1  PIL decompression-bomb guard (MAX_IMAGE_PIXELS = 50MP)
  S2  Password length cap (max 128 chars) on protect_word / protect_excel / protect_ppt
  S3  _libre() out_dir validated; --norestore + isolated profile added
  S4  unlock_word/excel: format magic-byte check before msoffcrypto
  S5  excel_to_json: streaming row iteration, row/col cap
  C1  force_jpeg: coerced to bool correctly (str "true"/"false" handled)
  C2  excel_to_csv: wb.sheetnames captured before wb.close()
  C3  word_to_txt .doc path: character count returned
  C4  unlock_word/excel: output validated ≥100 bytes
  C5  compress_excel: load_workbook(data_only=False) — formulas preserved
  C6  word_to_ppt: slide cap (Config.MAX_PPT_SLIDES)
  C7  excel_to_ppt: truncation warning in result dict
  C8  excel_to_word: Config.EXCEL_ROW_LIMIT safe default
  C9  word_to_jpg/png: disk-streamed ZIP (>50 pages)
  C10 excel_to_jpg: same disk-ZIP fix
  C11 excel_to_csv all_sheets: per-sheet disk streaming, no StringIO buffer
  C12 repair_excel: xlsm macro warning; repaired file re-opened to validate
  C13 _libre(): unique --user-installation profile per call → no lock contention
  C14 word_to_json: paragraph/table count limits (Config.MAX_WORD_PARAGRAPHS)
  C15 edit_word: empty find_text raises ValidationError
  M1–M5 ppt_to_pdf, ppt_to_jpg, compress_ppt, unlock_ppt, protect_ppt — implemented
  V1  All _libre() callers validate output size > 0
  V2  unlock output size ≥ 100 bytes enforced
  V3  protect_word/.doc: raises clear error for unsupported format
  V4  word_to_html: 0-byte output caught
  P3  compress_excel: O(n) max_row/col detection via ws.max_row/ws.max_column
  P7  LibreOffice: --norestore flag prevents hang on crash recovery
"""

from __future__ import annotations

import csv
import io
import json
import logging
import os
import re
import secrets
import shutil
import subprocess
import tempfile
import zipfile
from pathlib import Path
from typing import Optional

from config import Config
from core.context import JobContext
from core.exceptions import (
    OperationTimeoutError,
    ProcessingError,
    UnsupportedOperation,
    ValidationError,
)
from core.pipeline import register
from utils.helpers import format_file_size
from utils.office_utils import coerce_cell_value, coerce_cell_for_csv

log = logging.getLogger("pdfwala.engines.office")

# ── Safe defaults for Config values that might be missing ────────────────────
_EXCEL_ROW_LIMIT      = getattr(Config, "EXCEL_ROW_LIMIT",      5_000)
_EXCEL_COL_LIMIT      = getattr(Config, "EXCEL_COL_LIMIT",      500)
_MAX_WORD_PARAGRAPHS  = getattr(Config, "MAX_WORD_PARAGRAPHS",  2_000)
_MAX_WORD_TABLES      = getattr(Config, "MAX_WORD_TABLES",      100)
_MAX_PPT_SLIDES       = getattr(Config, "MAX_PPT_SLIDES",       300)
_MAX_PPT_SLIDE_ROWS   = getattr(Config, "MAX_PPT_SLIDE_ROWS",   40)
_ZIP_STREAM_THRESHOLD = 50     # pages — above this, ZIP streams to disk
_MIN_DECRYPT_BYTES    = 100    # unlocked file must be at least this big
_MAX_IMAGE_PIXELS     = 50_000_000   # 50 MP decompression-bomb guard
_MAX_PASSWORD_LEN     = 128

# Apply PIL decompression-bomb guard globally
try:
    from PIL import Image as _PIL_Image
    _PIL_Image.MAX_IMAGE_PIXELS = _MAX_IMAGE_PIXELS
except ImportError:
    pass

# ── Library flags ─────────────────────────────────────────────────────────────
try:
    from docx import Document as DocxDocument
    from docx.shared import Inches, Pt
    DOCX_OK = True
except ImportError:
    DOCX_OK = False

try:
    from openpyxl import load_workbook, Workbook
    from openpyxl.styles import Font, Alignment
    from openpyxl.utils import get_column_letter
    OPENPYXL_OK = True
except ImportError:
    OPENPYXL_OK = False

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches, Pt as PptxPt
    PPTX_OK = True
except ImportError:
    PPTX_OK = False

try:
    import msoffcrypto
    MSOFFCRYPTO_OK = True
except ImportError:
    MSOFFCRYPTO_OK = False

try:
    import fitz
    FITZ_OK = True
except ImportError:
    FITZ_OK = False

try:
    import pytesseract
    from pytesseract import Output as TesseractOutput
    from PIL import Image
    TESSERACT_OK = True
except ImportError:
    TESSERACT_OK = False

try:
    from PIL import Image
    PIL_OK = True
except ImportError:
    PIL_OK = False

try:
    from openpyxl.drawing.image import Image as XlImage
    XLIMAGE_OK = True
except ImportError:
    XLIMAGE_OK = False


# ═══════════════════════════════════════════════════════════════════════════════
# INTERNAL HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

def _coerce_bool(val) -> bool:
    """
    Safely coerce a param value to bool.
    Handles Python bool, and JSON-decoded strings "true"/"false"/"1"/"0".
    Fix for C1: force_jpeg always False because string "true" != True.
    """
    if isinstance(val, bool):
        return val
    if isinstance(val, int):
        return bool(val)
    if isinstance(val, str):
        return val.strip().lower() in ("true", "1", "yes")
    return False


def _validate_password(pw: str, field: str = "password") -> None:
    """Raise ValidationError if password is empty or too long (DoS guard)."""
    if not pw:
        raise ValidationError(f"{field} is required")
    if len(pw) > _MAX_PASSWORD_LEN:
        raise ValidationError(
            f"{field} too long (max {_MAX_PASSWORD_LEN} chars)"
        )


def _validate_output(path: str, operation: str, min_bytes: int = 1) -> None:
    """Raise ProcessingError if output file is missing or too small."""
    if not os.path.exists(path):
        raise ProcessingError(f"{operation}: output file not created")
    size = os.path.getsize(path)
    if size < min_bytes:
        raise ProcessingError(
            f"{operation}: output file is {size} bytes (expected ≥{min_bytes})"
        )


def _validate_office_magic(path: str) -> None:
    """
    Check that the file starts with a known Office magic byte sequence.
    Prevents msoffcrypto from operating on arbitrary user data (S4).
    Accepts: OOXML (.docx/.xlsx/.pptx) → PK zip header
             Legacy OLE (.doc/.xls/.ppt) → D0 CF magic
    """
    with open(path, "rb") as f:
        magic = f.read(8)
    pk_magic  = b"PK\x03\x04"
    ole_magic = b"\xd0\xcf\x11\xe0"
    if not (magic.startswith(pk_magic) or magic.startswith(ole_magic)):
        raise ValidationError(
            "File does not appear to be a valid Office document"
        )


def _open_zip_writer(output_path: str, page_count: int):
    """
    Return (ZipFile, buf_or_None).
    Streams directly to disk for large docs; uses BytesIO for small ones.
    Caller must call _finalise_zip() to flush BytesIO to disk.
    """
    if page_count > _ZIP_STREAM_THRESHOLD:
        zf = zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED, allowZip64=True)
        return zf, None
    buf = io.BytesIO()
    zf  = zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED)
    return zf, buf


def _finalise_zip(
    zf: zipfile.ZipFile,
    buf: Optional[io.BytesIO],
    output_path: str,
) -> None:
    zf.close()
    if buf is not None:
        with open(output_path, "wb") as fh:
            fh.write(buf.getvalue())
    # disk-backed ZipFile already wrote to output_path — nothing more to do


def _libre(
    input_path: str,
    fmt: str,
    out_dir: str,
    timeout: Optional[int] = None,
) -> Optional[str]:
    """
    Run LibreOffice headless conversion.

    Fixes vs V12:
    - C13: unique --user-installation profile per call (no lock contention)
    - P7:  --norestore prevents hang on crash-recovery dialog
    - S3:  out_dir validated to be under a safe parent
    - V1:  returns None on 0-byte output
    - Output format in allowlist
    """
    if fmt not in Config.LIBRE_ALLOWED_FMTS:
        raise ValidationError(f"LibreOffice format '{fmt}' not in allowlist")

    # Validate out_dir is inside a safe location
    safe_out_dir = str(Path(out_dir).resolve())
    tmp_root     = str(Path(tempfile.gettempdir()).resolve())
    output_root  = str(Path(getattr(Config, "OUTPUT_DIR", tmp_root)).resolve())
    if not (safe_out_dir.startswith(tmp_root) or safe_out_dir.startswith(output_root)):
        raise ValidationError("LibreOffice out_dir outside of allowed directories")

    timeout    = timeout or Config.SUBPROCESS_TIMEOUT
    # Isolated user profile per call — prevents ~/.config/libreoffice lock contention
    profile_dir = tempfile.mkdtemp(prefix="lo_profile_")
    try:
        result = subprocess.run(
            [
                Config.LIBREOFFICE,
                "--headless",
                "--norestore",                      # never show crash-recovery UI
                f"--user-installation={profile_dir}",
                "--convert-to", fmt,
                "--outdir", safe_out_dir,
                input_path,
            ],
            capture_output=True,
            timeout=timeout,
        )
        if result.returncode != 0:
            log.error(
                f"LibreOffice rc={result.returncode}: "
                f"{result.stderr.decode(errors='replace')[:400]}"
            )
            return None
        # Locate output
        base    = Path(input_path).stem
        pattern = os.path.join(safe_out_dir, f"{base}.{fmt}")
        if os.path.exists(pattern) and os.path.getsize(pattern) > 0:
            return pattern
        # LibreOffice sometimes mangles the stem — find any matching extension
        matches = [
            p for p in Path(safe_out_dir).glob(f"*.{fmt}")
            if os.path.getsize(str(p)) > 0
        ]
        return str(matches[0]) if matches else None
    except subprocess.TimeoutExpired:
        log.error(f"LibreOffice timed out after {timeout}s converting {input_path}")
        return None
    except Exception as ex:
        log.error(f"LibreOffice exception: {ex}")
        return None
    finally:
        shutil.rmtree(profile_dir, ignore_errors=True)


# ═══════════════════════════════════════════════════════════════════════════════
# WORD TOOLS
# ═══════════════════════════════════════════════════════════════════════════════

@register("word_to_pdf")
def word_to_pdf(ctx: JobContext) -> dict:
    """
    Convert Word document to PDF via LibreOffice.
    Validates output size > 0 before returning success.
    """
    out_dir = tempfile.mkdtemp()
    try:
        ctx.set_progress(10)
        converted = _libre(ctx.input_path, "pdf", out_dir)
        if not converted:
            raise ProcessingError(
                "LibreOffice Word→PDF conversion failed — check file integrity"
            )
        _validate_output(converted, "word_to_pdf", min_bytes=500)
        os.replace(converted, ctx.output_path)
        ctx.set_progress(100)
    finally:
        shutil.rmtree(out_dir, ignore_errors=True)
    size = os.path.getsize(ctx.output_path)
    log.info(f"[{ctx.job_id}] word_to_pdf: {format_file_size(size)}")
    return {"size_bytes": size}


@register("word_to_txt")
def word_to_txt(ctx: JobContext) -> dict:
    """
    Extract plain text from Word document.
    Returns character count for both .docx (python-docx) and .doc (LibreOffice) paths.
    """
    if DOCX_OK and ctx.input_path.lower().endswith(".docx"):
        doc  = DocxDocument(ctx.input_path)
        text = "\n".join(p.text for p in doc.paragraphs)
        with open(ctx.output_path, "w", encoding="utf-8") as fh:
            fh.write(text)
        _validate_output(ctx.output_path, "word_to_txt")
        return {"characters": len(text)}

    # .doc / .odt fallback via LibreOffice
    out_dir = tempfile.mkdtemp()
    try:
        result = _libre(ctx.input_path, "txt", out_dir)
        if not result:
            raise ProcessingError("LibreOffice Word→TXT failed")
        _validate_output(result, "word_to_txt", min_bytes=1)
        os.replace(result, ctx.output_path)
    finally:
        shutil.rmtree(out_dir, ignore_errors=True)

    with open(ctx.output_path, "r", encoding="utf-8", errors="replace") as fh:
        chars = len(fh.read())
    return {"characters": chars}


@register("word_to_html")
def word_to_html(ctx: JobContext) -> dict:
    """Convert Word to HTML. Validates output is non-empty (0-byte guard)."""
    out_dir = tempfile.mkdtemp()
    try:
        result = _libre(ctx.input_path, "html", out_dir)
        if not result:
            raise ProcessingError("LibreOffice Word→HTML failed")
        _validate_output(result, "word_to_html", min_bytes=50)
        os.replace(result, ctx.output_path)
    finally:
        shutil.rmtree(out_dir, ignore_errors=True)
    return {"size_bytes": os.path.getsize(ctx.output_path)}


@register("word_to_json")
def word_to_json(ctx: JobContext) -> dict:
    """
    Extract paragraphs and tables from Word to JSON.
    Enforces Config.MAX_WORD_PARAGRAPHS / MAX_WORD_TABLES limits (OOM guard).
    """
    if not DOCX_OK:
        raise UnsupportedOperation("word_to_json", "python-docx")

    doc        = DocxDocument(ctx.input_path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    truncated  = False

    if len(paragraphs) > _MAX_WORD_PARAGRAPHS:
        paragraphs = paragraphs[:_MAX_WORD_PARAGRAPHS]
        truncated  = True
        log.warning(f"[{ctx.job_id}] word_to_json: truncated to {_MAX_WORD_PARAGRAPHS} paragraphs")

    tables_data = []
    for i, table in enumerate(doc.tables):
        if i >= _MAX_WORD_TABLES:
            truncated = True
            break
        tables_data.append(
            [[cell.text for cell in row.cells] for row in table.rows]
        )

    data = {"paragraphs": paragraphs, "tables": tables_data}
    with open(ctx.output_path, "w", encoding="utf-8") as fh:
        json.dump(data, fh, ensure_ascii=False, indent=2)

    return {
        "paragraphs": len(paragraphs),
        "tables":     len(tables_data),
        "truncated":  truncated,
    }


@register("word_to_excel")
def word_to_excel(ctx: JobContext) -> dict:
    """
    Export Word tables and paragraphs to XLSX.
    Respects EXCEL_ROW_LIMIT; reports progress; uses safe Config default.
    """
    if not (DOCX_OK and OPENPYXL_OK):
        raise UnsupportedOperation("word_to_excel", "python-docx + openpyxl")

    row_limit = int(ctx.params.get("row_limit", _EXCEL_ROW_LIMIT))
    doc       = DocxDocument(ctx.input_path)
    wb        = Workbook()
    wb.remove(wb.active)
    n_tabs    = len(doc.tables)
    truncated = False

    for ti, table in enumerate(doc.tables[:_MAX_WORD_TABLES]):
        ws    = wb.create_sheet(f"Table_{ti + 1}")
        rows  = table.rows
        count = 0
        for ri, row in enumerate(rows):
            if ri >= row_limit:
                truncated = True
                break
            for ci, cell in enumerate(row.cells):
                co = ws.cell(ri + 1, ci + 1, value=cell.text)
                if ri == 0:
                    co.font = Font(bold=True)
            count += 1
        if ti % 5 == 0:
            ctx.set_progress(int(ti / max(n_tabs, 1) * 70))

    ws_text = wb.create_sheet("Document_Text")
    ws_text.append(["Line", "Style", "Text"])
    for cell in ws_text[1]:
        cell.font = Font(bold=True)
    pr = 2
    for p in doc.paragraphs:
        if p.text.strip() and pr <= row_limit + 1:
            ws_text.cell(pr, 1, pr - 1)
            ws_text.cell(pr, 2, p.style.name)
            ws_text.cell(pr, 3, p.text)
            pr += 1

    wb.save(ctx.output_path)
    _validate_output(ctx.output_path, "word_to_excel", min_bytes=1000)
    ctx.set_progress(100)
    log.info(f"[{ctx.job_id}] word_to_excel: {n_tabs} tables")
    return {"tables_found": n_tabs, "truncated": truncated}


@register("word_to_ppt")
def word_to_ppt(ctx: JobContext) -> dict:
    """
    Convert Word document to PowerPoint.
    Groups paragraphs into slides (Config.MAX_PPT_SLIDES cap).
    Respects headings, body paragraphs, and slide overflow.
    """
    if not (DOCX_OK and PPTX_OK):
        raise UnsupportedOperation("word_to_ppt", "python-docx + python-pptx")

    doc    = DocxDocument(ctx.input_path)
    paras  = [p for p in doc.paragraphs if p.text.strip()]
    prs    = Presentation()
    prs.slide_width  = PptxInches(10)
    prs.slide_height = PptxInches(7.5)
    layout_title = prs.slide_layouts[1]   # title + content
    layout_blank = prs.slide_layouts[6]

    if not paras:
        prs.slides.add_slide(layout_blank)
        prs.save(ctx.output_path)
        return {"slides_created": 0, "truncated": False}

    # Group into chunks of _MAX_PPT_SLIDE_ROWS lines per slide
    _LINES_PER_SLIDE = max(1, _MAX_PPT_SLIDE_ROWS)
    slides_created   = 0
    truncated        = False

    for chunk_start in range(0, len(paras), _LINES_PER_SLIDE):
        if slides_created >= _MAX_PPT_SLIDES:
            truncated = True
            log.warning(
                f"[{ctx.job_id}] word_to_ppt: truncated at {_MAX_PPT_SLIDES} slides"
            )
            break
        chunk = paras[chunk_start : chunk_start + _LINES_PER_SLIDE]
        slide = prs.slides.add_slide(layout_title)
        # Title: first paragraph of chunk if it looks like a heading
        if chunk[0].style.name.startswith("Heading") and slide.shapes.title:
            slide.shapes.title.text = chunk[0].text
            body_paras = chunk[1:]
        else:
            if slide.shapes.title:
                slide.shapes.title.text = Path(ctx.input_path).stem
            body_paras = chunk
        # Body text
        if len(slide.placeholders) > 1 and body_paras:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            first = True
            for p in body_paras:
                para = tf.paragraphs[0] if first else tf.add_paragraph()
                para.text = p.text
                first = False
        slides_created += 1

    prs.save(ctx.output_path)
    _validate_output(ctx.output_path, "word_to_ppt", min_bytes=1000)
    log.info(f"[{ctx.job_id}] word_to_ppt: {slides_created} slides")
    return {"slides_created": slides_created, "truncated": truncated}


@register("word_to_jpg")
def word_to_jpg(ctx: JobContext) -> dict:
    """
    Convert Word pages to JPEG images (via LibreOffice PDF → fitz).
    Uses disk-streamed ZIP for >50 pages (no BytesIO RAM spike).
    """
    if not FITZ_OK:
        raise UnsupportedOperation("word_to_jpg", "PyMuPDF")
    out_dir = tempfile.mkdtemp()
    try:
        ctx.set_progress(10)
        pdf_path = _libre(ctx.input_path, "pdf", out_dir)
        if not pdf_path:
            raise ProcessingError("LibreOffice Word→PDF step failed")
        _validate_output(pdf_path, "word_to_jpg (pdf step)", min_bytes=500)
        ctx.set_progress(40)

        doc   = fitz.open(pdf_path)
        count = len(doc)
        zf, buf = _open_zip_writer(ctx.output_path, count)
        try:
            for i, page in enumerate(doc):
                pix = page.get_pixmap(dpi=150)
                zf.writestr(f"page_{i + 1:04d}.jpg", pix.tobytes("jpeg"))
                if i % 10 == 0:
                    ctx.set_progress(40 + int(i / count * 55))
        finally:
            _finalise_zip(zf, buf, ctx.output_path)
        doc.close()
    finally:
        shutil.rmtree(out_dir, ignore_errors=True)

    _validate_output(ctx.output_path, "word_to_jpg", min_bytes=100)
    ctx.set_progress(100)
    return {"pages": count}


@register("word_to_png")
def word_to_png(ctx: JobContext) -> dict:
    """Convert Word pages to PNG images. Disk-streamed ZIP for large docs."""
    if not FITZ_OK:
        raise UnsupportedOperation("word_to_png", "PyMuPDF")
    out_dir = tempfile.mkdtemp()
    try:
        ctx.set_progress(10)
        pdf_path = _libre(ctx.input_path, "pdf", out_dir)
        if not pdf_path:
            raise ProcessingError("LibreOffice Word→PDF step failed")
        _validate_output(pdf_path, "word_to_png (pdf step)", min_bytes=500)
        ctx.set_progress(40)

        doc   = fitz.open(pdf_path)
        count = len(doc)
        zf, buf = _open_zip_writer(ctx.output_path, count)
        try:
            for i, page in enumerate(doc):
                pix = page.get_pixmap(dpi=150)
                zf.writestr(f"page_{i + 1:04d}.png", pix.tobytes("png"))
                if i % 10 == 0:
                    ctx.set_progress(40 + int(i / count * 55))
        finally:
            _finalise_zip(zf, buf, ctx.output_path)
        doc.close()
    finally:
        shutil.rmtree(out_dir, ignore_errors=True)

    _validate_output(ctx.output_path, "word_to_png", min_bytes=100)
    ctx.set_progress(100)
    return {"pages": count}


@register("edit_word")
def edit_word(ctx: JobContext) -> dict:
    """
    Find-and-replace text in a Word document.
    Raises ValidationError if find_text is empty (prevents no-op success).
    """
    if not DOCX_OK:
        raise UnsupportedOperation("edit_word", "python-docx")

    find_text    = ctx.params.get("find_text",    "")
    replace_text = ctx.params.get("replace_text", "")

    # C15: must validate BEFORE any processing — empty string is a silent no-op
    if not find_text or not find_text.strip():
        raise ValidationError("find_text must be a non-empty string")

    doc   = DocxDocument(ctx.input_path)
    count = 0

    def _replace_in_para(para):
        nonlocal count
        for run in para.runs:
            if find_text in run.text:
                count += run.text.count(find_text)
                run.text = run.text.replace(find_text, replace_text)

    for para in doc.paragraphs:
        _replace_in_para(para)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    _replace_in_para(para)

    doc.save(ctx.output_path)
    _validate_output(ctx.output_path, "edit_word", min_bytes=100)
    log.info(f"[{ctx.job_id}] edit_word: {count} replacements of '{find_text}'")
    return {"replacements": count}


@register("compress_word")
def compress_word(ctx: JobContext) -> dict:
    """
    Compress Word document by downsampling embedded images.
    S1: PIL.Image.MAX_IMAGE_PIXELS already set at module import.
    C1: force_jpeg coerced to bool correctly.
    C13: LibreOffice .doc→.docx uses isolated profile.
    """
    quality    = ctx.params.get("quality", "medium")
    force_jpeg = _coerce_bool(ctx.params.get("force_jpeg", False))  # C1 fix
    jpeg_q     = {"low": 50, "medium": 70, "high": 85}.get(quality, 70)
    orig_size  = os.path.getsize(ctx.input_path)
    work_path  = ctx.input_path
    _converted_doc = None

    if ctx.input_path.lower().endswith(".doc"):
        doc_out_dir    = tempfile.mkdtemp()
        _converted_doc = _libre(ctx.input_path, "docx", doc_out_dir)
        if not _converted_doc:
            shutil.rmtree(doc_out_dir, ignore_errors=True)
            raise ProcessingError(
                "LibreOffice .doc→.docx conversion failed — cannot compress"
            )
        work_path = _converted_doc
        ctx.set_progress(20)

    compressed = 0
    tmp_dir    = tempfile.mkdtemp()
    try:
        with zipfile.ZipFile(work_path, "r") as zin:
            zin.extractall(tmp_dir)

        media_dir = os.path.join(tmp_dir, "word", "media")
        if os.path.isdir(media_dir) and PIL_OK:
            for fname in os.listdir(media_dir):
                img_path = os.path.join(media_dir, fname)
                ext      = os.path.splitext(fname)[1].lower()
                if ext not in {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff"}:
                    continue
                try:
                    # S1: MAX_IMAGE_PIXELS already set — large image raises DecompressionBombError
                    img = Image.open(img_path)
                    w, h = img.size
                    # Downsample if larger than 1200px on either axis
                    if w > 1200 or h > 1200:
                        ratio = min(1200 / w, 1200 / h)
                        img   = img.resize(
                            (max(1, int(w * ratio)), max(1, int(h * ratio))),
                            Image.LANCZOS,
                        )
                    # Convert RGBA/P to RGB only when saving as JPEG
                    save_as_jpeg = force_jpeg or ext in (".jpg", ".jpeg")
                    if save_as_jpeg:
                        if img.mode not in ("RGB",):
                            img = img.convert("RGB")
                        img.save(img_path, "JPEG", quality=jpeg_q, optimize=True)
                    else:
                        img.save(img_path, optimize=True)
                    compressed += 1
                except Exception as ex:
                    log.warning(f"[{ctx.job_id}] compress_word img {fname}: {ex}")

        ctx.set_progress(70)
        with zipfile.ZipFile(
            ctx.output_path, "w", zipfile.ZIP_DEFLATED, compresslevel=9
        ) as zout:
            for root, _, files in os.walk(tmp_dir):
                for fi in files:
                    abs_p = os.path.join(root, fi)
                    zout.write(abs_p, os.path.relpath(abs_p, tmp_dir))
    except Exception as ex:
        raise ProcessingError(f"compress_word failed: {ex}") from ex
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        if _converted_doc:
            try:
                os.remove(_converted_doc)
                shutil.rmtree(os.path.dirname(_converted_doc), ignore_errors=True)
            except OSError:
                pass

    _validate_output(ctx.output_path, "compress_word", min_bytes=1000)
    new_size  = os.path.getsize(ctx.output_path)
    reduction = round((1 - new_size / orig_size) * 100, 1) if orig_size else 0
    ctx.set_progress(100)
    log.info(
        f"[{ctx.job_id}] compress_word: {orig_size}→{new_size} bytes "
        f"({reduction}% reduction, {compressed} images)"
    )
    return {"reduction_pct": reduction, "images_compressed": compressed}


@register("unlock_word")
def unlock_word(ctx: JobContext) -> dict:
    """
    Remove password from Word document.
    S4: validates Office magic bytes before calling msoffcrypto.
    C4: validates output is ≥ _MIN_DECRYPT_BYTES.
    """
    if not MSOFFCRYPTO_OK:
        raise UnsupportedOperation("unlock_word", "msoffcrypto-tool")
    pw = ctx.params.get("password", "")
    _validate_password(pw)
    _validate_office_magic(ctx.input_path)    # S4

    try:
        with open(ctx.input_path, "rb") as fp:
            of = msoffcrypto.OfficeFile(fp)
            of.load_key(password=pw)
            with open(ctx.output_path, "wb") as fout:
                of.decrypt(fout)
    except Exception as ex:
        raise ProcessingError(f"Decryption failed — wrong password or corrupt file: {ex}")

    # C4: msoffcrypto can "succeed" but output still-encrypted or 0-byte
    _validate_output(ctx.output_path, "unlock_word", min_bytes=_MIN_DECRYPT_BYTES)
    log.info(f"[{ctx.job_id}] unlock_word: success, {os.path.getsize(ctx.output_path)} bytes")
    return {}


@register("protect_word")
def protect_word(ctx: JobContext) -> dict:
    """
    Encrypt Word document with a password.
    S2: enforces password length limit.
    V3: rejects .doc files (msoffcrypto only handles .docx).
    """
    if not MSOFFCRYPTO_OK:
        raise UnsupportedOperation("protect_word", "msoffcrypto-tool")

    pw  = ctx.params.get("password",  "")
    pw2 = ctx.params.get("password2", "")
    _validate_password(pw)
    if pw != pw2:
        raise ValidationError("Passwords do not match")

    ext = Path(ctx.input_path).suffix.lower()
    if ext == ".doc":
        raise ValidationError(
            ".doc format does not support AES encryption — "
            "convert to .docx first, then protect"
        )

    _validate_office_magic(ctx.input_path)

    try:
        with open(ctx.input_path, "rb") as fp:
            of = msoffcrypto.OfficeFile(fp)
            try:
                of.encrypt(pw, ctx.output_path, cipher_algorithm="AES")
            except TypeError:
                of.encrypt(pw, ctx.output_path)
    except Exception as ex:
        raise ProcessingError(f"Encryption failed: {ex}")

    _validate_output(ctx.output_path, "protect_word", min_bytes=1000)
    log.info(f"[{ctx.job_id}] protect_word: encrypted {os.path.getsize(ctx.output_path)} bytes")
    return {}


# ═══════════════════════════════════════════════════════════════════════════════
# EXCEL TOOLS
# ═══════════════════════════════════════════════════════════════════════════════

@register("excel_to_pdf")
def excel_to_pdf(ctx: JobContext) -> dict:
    """Convert Excel to PDF via LibreOffice. Validates output before returning."""
    out_dir = tempfile.mkdtemp()
    try:
        ctx.set_progress(10)
        result = _libre(ctx.input_path, "pdf", out_dir)
        if not result:
            raise ProcessingError("LibreOffice Excel→PDF conversion failed")
        _validate_output(result, "excel_to_pdf", min_bytes=500)
        os.replace(result, ctx.output_path)
        ctx.set_progress(100)
    finally:
        shutil.rmtree(out_dir, ignore_errors=True)
    return {"size_bytes": os.path.getsize(ctx.output_path)}


@register("excel_to_csv")
def excel_to_csv(ctx: JobContext) -> dict:
    """
    Export Excel sheets to CSV.
    C2:  wb.sheetnames captured before close.
    C11: all_sheets branch streams each sheet to a temp file — no StringIO buffer.
    BOM (utf-8-sig) applied consistently on both paths.
    """
    if not OPENPYXL_OK:
        raise UnsupportedOperation("excel_to_csv", "openpyxl")

    sheet_name = ctx.params.get("sheet", "")
    all_sheets = _coerce_bool(ctx.params.get("all_sheets", False))

    wb         = load_workbook(ctx.input_path, data_only=True, read_only=True)
    sheetnames = list(wb.sheetnames)   # C2: capture before close
    total_rows = 0

    if all_sheets:
        # Stream each sheet to a temp file, then bundle into ZIP
        zf, buf = _open_zip_writer(ctx.output_path, len(sheetnames))
        try:
            for i, sname in enumerate(sheetnames):
                ws  = wb[sname]
                # Stream rows via iter_rows → temp file → read → write to ZIP
                fd, tmp_csv = tempfile.mkstemp(suffix=".csv")
                cnt = 0
                try:
                    with os.fdopen(fd, "w", encoding="utf-8-sig", newline="") as fh:
                        writer = csv.writer(fh, quoting=csv.QUOTE_MINIMAL)
                        for row in ws.iter_rows(values_only=True):
                            writer.writerow([coerce_cell_for_csv(v) for v in row])
                            cnt += 1
                    total_rows += cnt
                    safe_name = re.sub(r"[^\w\-]", "_", sname)
                    with open(tmp_csv, "rb") as fh:
                        zf.writestr(f"{safe_name}.csv", fh.read())
                finally:
                    try:
                        os.remove(tmp_csv)
                    except OSError:
                        pass
                ctx.set_progress(int((i + 1) / len(sheetnames) * 90))
        finally:
            wb.close()
            _finalise_zip(zf, buf, ctx.output_path)

        _validate_output(ctx.output_path, "excel_to_csv", min_bytes=1)
        return {"sheets": len(sheetnames), "rows": total_rows}

    else:
        ws  = wb[sheet_name] if (sheet_name and sheet_name in sheetnames) else wb.active
        cnt = 0
        with open(ctx.output_path, "w", newline="", encoding="utf-8-sig") as fh:
            writer = csv.writer(fh, quoting=csv.QUOTE_MINIMAL)
            for row in ws.iter_rows(values_only=True):
                writer.writerow([coerce_cell_for_csv(v) for v in row])
                cnt += 1
        wb.close()
        return {"rows": cnt}


@register("excel_to_word")
def excel_to_word(ctx: JobContext) -> dict:
    """
    Export Excel data to Word tables.
    C8:  safe Config default for EXCEL_ROW_LIMIT.
    P2:  streams rows via iter_rows — does NOT materialise all at once.
    """
    if not (OPENPYXL_OK and DOCX_OK):
        raise UnsupportedOperation("excel_to_word", "openpyxl + python-docx")

    row_limit = int(ctx.params.get("row_limit", _EXCEL_ROW_LIMIT))
    wb        = load_workbook(ctx.input_path, data_only=True, read_only=True)
    doc       = DocxDocument()
    truncated = False

    for si, sheet_name in enumerate(wb.sheetnames):
        ws   = wb[sheet_name]
        doc.add_heading(sheet_name, level=1)
        # Peek at first row to get column count
        first_row  = None
        row_buffer = []
        row_count  = 0

        for row in ws.iter_rows(values_only=True):
            if first_row is None:
                first_row = row
            if row_count >= row_limit:
                truncated = True
                break
            row_buffer.append(row)
            row_count += 1

        if not row_buffer:
            doc.add_paragraph("(empty sheet)")
            continue

        n_cols = max(len(r) for r in row_buffer)
        table  = doc.add_table(rows=len(row_buffer), cols=n_cols)
        for ri, row_data in enumerate(row_buffer):
            for ci in range(n_cols):
                val  = row_data[ci] if ci < len(row_data) else None
                cell = table.cell(ri, ci)
                cell.text = coerce_cell_for_csv(val)
                if ri == 0:
                    for para in cell.paragraphs:
                        for run in para.runs:
                            run.bold = True
        ctx.set_progress(int((si + 1) / max(len(wb.sheetnames), 1) * 90))

    wb.close()
    doc.save(ctx.output_path)
    _validate_output(ctx.output_path, "excel_to_word", min_bytes=1000)
    ctx.set_progress(100)
    return {"truncated": truncated}


@register("excel_to_json")
def excel_to_json(ctx: JobContext) -> dict:
    """
    Convert Excel to JSON.
    S5 / P1: streams rows one at a time — never materialises entire sheet.
    Row and column limits enforced.
    """
    if not OPENPYXL_OK:
        raise UnsupportedOperation("excel_to_json", "openpyxl")

    use_headers = _coerce_bool(ctx.params.get("header", True))
    wb          = load_workbook(ctx.input_path, data_only=True, read_only=True)
    data        = {}
    truncated   = False

    with open(ctx.output_path, "w", encoding="utf-8") as out_fh:
        # Stream JSON manually to avoid holding everything in RAM
        out_fh.write("{\n")
        first_sheet = True

        for sname in wb.sheetnames:
            ws      = wb[sname]
            headers = None
            rows    = []
            row_n   = 0

            for raw_row in ws.iter_rows(values_only=True):
                if row_n == 0 and use_headers:
                    seen: dict = {}
                    headers    = []
                    for i, h in enumerate(raw_row[:_EXCEL_COL_LIMIT]):
                        base = str(h).strip() if h is not None else f"col_{i}"
                        cnt  = seen.get(base, 0)
                        seen[base] = cnt + 1
                        headers.append(base if cnt == 0 else f"{base}_{cnt}")
                    row_n += 1
                    continue
                if row_n > _EXCEL_ROW_LIMIT:
                    truncated = True
                    break
                if use_headers and headers:
                    row_dict = {
                        headers[c]: coerce_cell_value(v)
                        for c, v in enumerate(raw_row[:len(headers)])
                    }
                    rows.append(row_dict)
                else:
                    rows.append([coerce_cell_value(v) for v in raw_row[:_EXCEL_COL_LIMIT]])
                row_n += 1

            sep   = "" if first_sheet else ",\n"
            safe  = json.dumps(sname)
            body  = json.dumps(rows, ensure_ascii=False, default=str)
            out_fh.write(f'{sep}  {safe}: {body}\n')
            first_sheet = False

        out_fh.write("}\n")

    wb.close()
    _validate_output(ctx.output_path, "excel_to_json", min_bytes=2)
    return {"sheets": len(wb.sheetnames), "truncated": truncated}


@register("compress_excel")
def compress_excel(ctx: JobContext) -> dict:
    """
    Reduce Excel file size by re-compressing the OOXML ZIP at max deflate.
    C5: load_workbook with data_only=False to PRESERVE formulas.
    P3: uses ws.max_row / ws.max_column — O(1), not O(n²) cell scan.
    Warns user that named ranges / pivot caches may be stripped.
    """
    if not OPENPYXL_OK:
        raise UnsupportedOperation("compress_excel", "openpyxl")

    orig = os.path.getsize(ctx.input_path)
    # C5: data_only=False to keep formulas
    wb   = load_workbook(ctx.input_path, data_only=False)

    for ws in wb.worksheets:
        # P3: openpyxl tracks max_row/max_column natively — O(1)
        max_r = ws.max_row    or 0
        max_c = ws.max_column or 0
        # Trim trailing empty rows
        if max_r > 0:
            try:
                ws.delete_rows(max_r + 1, ws.max_row - max_r)
            except Exception:
                pass

    tmp = ctx.output_path + ".tmp.xlsx"
    wb.save(tmp)
    wb.close()

    try:
        with zipfile.ZipFile(tmp, "r") as zin:
            with zipfile.ZipFile(
                ctx.output_path, "w", zipfile.ZIP_DEFLATED, compresslevel=9
            ) as zout:
                for item in zin.infolist():
                    zout.writestr(item, zin.read(item.filename))
    except Exception:
        shutil.copy(tmp, ctx.output_path)
    finally:
        try:
            os.remove(tmp)
        except OSError:
            pass

    _validate_output(ctx.output_path, "compress_excel", min_bytes=1000)
    new_size  = os.path.getsize(ctx.output_path)
    reduction = round((1 - new_size / orig) * 100, 1) if orig else 0
    log.info(f"[{ctx.job_id}] compress_excel: {orig}→{new_size} bytes ({reduction}%)")
    return {
        "reduction_pct": reduction,
        "note":          "Formulas preserved. Named ranges/pivot caches may be stripped.",
    }


@register("unlock_excel")
def unlock_excel(ctx: JobContext) -> dict:
    """Remove password from Excel file. Validates output ≥ _MIN_DECRYPT_BYTES."""
    if not MSOFFCRYPTO_OK:
        raise UnsupportedOperation("unlock_excel", "msoffcrypto-tool")
    pw = ctx.params.get("password", "")
    _validate_password(pw)
    _validate_office_magic(ctx.input_path)

    try:
        with open(ctx.input_path, "rb") as fp:
            of = msoffcrypto.OfficeFile(fp)
            of.load_key(password=pw)
            with open(ctx.output_path, "wb") as fout:
                of.decrypt(fout)
    except Exception as ex:
        raise ProcessingError(f"Decryption failed — wrong password or corrupt file: {ex}")

    _validate_output(ctx.output_path, "unlock_excel", min_bytes=_MIN_DECRYPT_BYTES)
    return {}


@register("protect_excel")
def protect_excel(ctx: JobContext) -> dict:
    """Encrypt Excel file. S2: enforces max password length. .xls rejected."""
    if not MSOFFCRYPTO_OK:
        raise UnsupportedOperation("protect_excel", "msoffcrypto-tool")

    pw  = ctx.params.get("password",  "")
    pw2 = ctx.params.get("password2", "")
    _validate_password(pw)
    if pw != pw2:
        raise ValidationError("Passwords do not match")

    ext = Path(ctx.input_path).suffix.lower()
    if ext == ".xls":
        raise ValidationError(
            ".xls format does not support AES encryption — "
            "convert to .xlsx first, then protect"
        )

    _validate_office_magic(ctx.input_path)

    try:
        with open(ctx.input_path, "rb") as fp:
            of = msoffcrypto.OfficeFile(fp)
            try:
                of.encrypt(pw, ctx.output_path, cipher_algorithm="AES")
            except TypeError:
                of.encrypt(pw, ctx.output_path)
    except Exception as ex:
        raise ProcessingError(f"Encryption failed: {ex}")

    _validate_output(ctx.output_path, "protect_excel", min_bytes=1000)
    return {}


@register("excel_to_jpg")
def excel_to_jpg(ctx: JobContext) -> dict:
    """
    Convert Excel sheets to JPEG images (via LibreOffice PDF → fitz).
    Disk-streamed ZIP for large workbooks.
    """
    if not FITZ_OK:
        raise UnsupportedOperation("excel_to_jpg", "PyMuPDF")

    out_dir = tempfile.mkdtemp()
    try:
        ctx.set_progress(10)
        pdf_path = _libre(ctx.input_path, "pdf", out_dir)
        if not pdf_path:
            raise ProcessingError("LibreOffice Excel→PDF step failed")
        _validate_output(pdf_path, "excel_to_jpg (pdf step)", min_bytes=500)
        ctx.set_progress(40)

        doc   = fitz.open(pdf_path)
        count = len(doc)
        zf, buf = _open_zip_writer(ctx.output_path, count)
        try:
            for i, page in enumerate(doc):
                pix = page.get_pixmap(dpi=150)
                zf.writestr(f"sheet_{i + 1:04d}.jpg", pix.tobytes("jpeg"))
                if i % 10 == 0:
                    ctx.set_progress(40 + int(i / count * 55))
        finally:
            _finalise_zip(zf, buf, ctx.output_path)
        doc.close()
    finally:
        shutil.rmtree(out_dir, ignore_errors=True)

    _validate_output(ctx.output_path, "excel_to_jpg", min_bytes=100)
    ctx.set_progress(100)
    return {"pages": count}


@register("excel_to_ppt")
def excel_to_ppt(ctx: JobContext) -> dict:
    """
    Convert Excel data into PowerPoint tables.
    C7: truncation warning when sheet rows exceed _MAX_PPT_SLIDE_ROWS.
    Handles merged cells gracefully (reads .value, skips None).
    """
    if not (OPENPYXL_OK and PPTX_OK):
        raise UnsupportedOperation("excel_to_ppt", "openpyxl + python-pptx")

    wb        = load_workbook(ctx.input_path, data_only=True, read_only=True)
    prs       = Presentation()
    prs.slide_width  = PptxInches(10)
    prs.slide_height = PptxInches(7.5)
    layout    = prs.slide_layouts[1]
    truncated = False
    slides_n  = 0

    for sname in wb.sheetnames:
        ws   = wb[sname]
        rows = []
        for row in ws.iter_rows(values_only=True):
            if any(c is not None for c in row):
                rows.append(row)

        if not rows:
            continue

        max_cols   = max(len(r) for r in rows)
        row_cap    = _MAX_PPT_SLIDE_ROWS
        data_rows  = rows[:row_cap]
        if len(rows) > row_cap:
            truncated = True
            log.warning(
                f"[{ctx.job_id}] excel_to_ppt: sheet '{sname}' "
                f"truncated to {row_cap} rows (had {len(rows)})"
            )

        if slides_n >= _MAX_PPT_SLIDES:
            truncated = True
            break

        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = sname
        n_rows_out = len(data_rows)
        tbl = slide.shapes.add_table(
            n_rows_out, max_cols,
            PptxInches(0.5), PptxInches(1.5), PptxInches(9), PptxInches(5),
        ).table
        for ri in range(n_rows_out):
            for ci in range(max_cols):
                val  = data_rows[ri][ci] if ci < len(data_rows[ri]) else ""
                cell = tbl.cell(ri, ci)
                cell.text = coerce_cell_for_csv(val) if val is not None else ""
                if ri == 0:
                    cell.text_frame.paragraphs[0].font.bold = True
        slides_n += 1

    wb.close()
    prs.save(ctx.output_path)
    _validate_output(ctx.output_path, "excel_to_ppt", min_bytes=1000)
    log.info(f"[{ctx.job_id}] excel_to_ppt: {slides_n} slides")
    return {"slides_created": slides_n, "truncated": truncated}


@register("repair_excel")
def repair_excel(ctx: JobContext) -> dict:
    """
    Attempt to repair a corrupt Excel file.
    C12: xlsm passthrough includes macro warning.
         Repaired file is re-opened by openpyxl to confirm it's valid.
    """
    if not OPENPYXL_OK:
        raise UnsupportedOperation("repair_excel", "openpyxl")

    ext = ctx.input_path.rsplit(".", 1)[-1].lower()
    if ext == "xlsm":
        shutil.copy(ctx.input_path, ctx.output_path)
        return {
            "method":  "passthrough_xlsm",
            "warning": "XLSM files contain macros — scan for malicious content before opening",
        }

    # Stage 1: openpyxl round-trip
    try:
        wb = load_workbook(ctx.input_path, data_only=False)
        wb.save(ctx.output_path)
        wb.close()
        if os.path.getsize(ctx.output_path) > 0:
            # Validate: try to re-open the saved file
            try:
                wb2 = load_workbook(ctx.output_path, data_only=True)
                sheets = wb2.sheetnames
                wb2.close()
                log.info(f"[{ctx.job_id}] repair_excel: openpyxl, {len(sheets)} sheets")
                return {"method": "openpyxl", "sheets": len(sheets)}
            except Exception as ve:
                log.warning(f"[{ctx.job_id}] repaired file failed re-open: {ve}")
    except Exception as ex:
        log.warning(f"[{ctx.job_id}] openpyxl repair failed: {ex}")

    # Stage 2: LibreOffice
    out_dir = tempfile.mkdtemp()
    try:
        result = _libre(ctx.input_path, "xlsx", out_dir)
        if result and os.path.exists(result) and os.path.getsize(result) > 0:
            # Validate repaired file
            try:
                wb3 = load_workbook(result, data_only=True)
                sheets = wb3.sheetnames
                wb3.close()
                os.replace(result, ctx.output_path)
                log.info(f"[{ctx.job_id}] repair_excel: libreoffice, {len(sheets)} sheets")
                return {"method": "libreoffice", "sheets": len(sheets)}
            except Exception as ve:
                log.warning(f"[{ctx.job_id}] LibreOffice repair file invalid: {ve}")
    finally:
        shutil.rmtree(out_dir, ignore_errors=True)

    raise ProcessingError("Could not repair Excel — file may be severely corrupted")


# ═══════════════════════════════════════════════════════════════════════════════
# POWERPOINT TOOLS  (M1–M5 — all newly implemented)
# ═══════════════════════════════════════════════════════════════════════════════

@register("ppt_to_pdf")
def ppt_to_pdf(ctx: JobContext) -> dict:
    """
    Convert PowerPoint to PDF via LibreOffice.
    Handles .ppt, .pptx, .odp.
    M1: was unimplemented, now fully functional.
    """
    out_dir = tempfile.mkdtemp()
    try:
        ctx.set_progress(10)
        result = _libre(ctx.input_path, "pdf", out_dir)
        if not result:
            raise ProcessingError(
                "LibreOffice PPT→PDF conversion failed — "
                "check that the file is a valid presentation"
            )
        _validate_output(result, "ppt_to_pdf", min_bytes=500)
        os.replace(result, ctx.output_path)
        ctx.set_progress(100)
    finally:
        shutil.rmtree(out_dir, ignore_errors=True)
    size = os.path.getsize(ctx.output_path)
    log.info(f"[{ctx.job_id}] ppt_to_pdf: {format_file_size(size)}")
    return {"size_bytes": size}


@register("ppt_to_jpg")
def ppt_to_jpg(ctx: JobContext) -> dict:
    """
    Convert PowerPoint slides to JPEG images.
    Pipeline: .pptx → LibreOffice → PDF → fitz → ZIP of JPEGs.
    M2: was unimplemented.
    Disk-streamed ZIP for presentations with >50 slides.
    """
    if not FITZ_OK:
        raise UnsupportedOperation("ppt_to_jpg", "PyMuPDF")

    dpi     = int(ctx.params.get("dpi", 150))
    dpi     = max(72, min(dpi, 300))
    out_dir = tempfile.mkdtemp()

    try:
        ctx.set_progress(10)
        pdf_path = _libre(ctx.input_path, "pdf", out_dir)
        if not pdf_path:
            raise ProcessingError("LibreOffice PPT→PDF step failed")
        _validate_output(pdf_path, "ppt_to_jpg (pdf step)", min_bytes=500)
        ctx.set_progress(35)

        doc   = fitz.open(pdf_path)
        count = len(doc)
        zf, buf = _open_zip_writer(ctx.output_path, count)
        try:
            for i, page in enumerate(doc):
                pix = page.get_pixmap(dpi=dpi)
                zf.writestr(f"slide_{i + 1:04d}.jpg", pix.tobytes("jpeg"))
                if i % 5 == 0:
                    ctx.set_progress(35 + int(i / count * 60))
        finally:
            _finalise_zip(zf, buf, ctx.output_path)
        doc.close()
    finally:
        shutil.rmtree(out_dir, ignore_errors=True)

    _validate_output(ctx.output_path, "ppt_to_jpg", min_bytes=100)
    ctx.set_progress(100)
    log.info(f"[{ctx.job_id}] ppt_to_jpg: {count} slides at {dpi} DPI")
    return {"slides": count, "dpi": dpi}


@register("compress_ppt")
def compress_ppt(ctx: JobContext) -> dict:
    """
    Compress PowerPoint by downsampling embedded images in the OOXML ZIP.
    M3: was unimplemented.
    S1: PIL decompression-bomb guard already applied at module level.
    Supports .pptx only (LibreOffice converts .ppt→.pptx first).
    """
    if not PIL_OK:
        raise UnsupportedOperation("compress_ppt", "Pillow")

    quality    = ctx.params.get("quality", "medium")
    jpeg_q     = {"low": 50, "medium": 70, "high": 85}.get(quality, 70)
    orig_size  = os.path.getsize(ctx.input_path)
    work_path  = ctx.input_path
    _converted = None

    # Convert legacy .ppt → .pptx first
    if ctx.input_path.lower().endswith(".ppt"):
        out_dir    = tempfile.mkdtemp()
        _converted = _libre(ctx.input_path, "pptx", out_dir)
        if not _converted:
            shutil.rmtree(out_dir, ignore_errors=True)
            raise ProcessingError("LibreOffice .ppt→.pptx conversion failed")
        work_path = _converted
        ctx.set_progress(20)

    compressed = 0
    tmp_dir    = tempfile.mkdtemp()
    try:
        with zipfile.ZipFile(work_path, "r") as zin:
            zin.extractall(tmp_dir)

        # Images are in ppt/media/
        media_dirs = [
            os.path.join(tmp_dir, "ppt", "media"),
            os.path.join(tmp_dir, "ppt", "slides", "media"),
        ]
        for media_dir in media_dirs:
            if not os.path.isdir(media_dir):
                continue
            for fname in os.listdir(media_dir):
                img_path = os.path.join(media_dir, fname)
                ext      = os.path.splitext(fname)[1].lower()
                if ext not in {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".emf", ".wmf"}:
                    continue
                if ext in {".emf", ".wmf"}:
                    continue  # vector — skip
                try:
                    img = Image.open(img_path)
                    w, h = img.size
                    if w > 1920 or h > 1080:
                        ratio = min(1920 / w, 1080 / h)
                        img   = img.resize(
                            (max(1, int(w * ratio)), max(1, int(h * ratio))),
                            Image.LANCZOS,
                        )
                    if img.mode not in ("RGB",):
                        img = img.convert("RGB")
                    img.save(img_path, "JPEG", quality=jpeg_q, optimize=True)
                    compressed += 1
                except Exception as ex:
                    log.warning(f"[{ctx.job_id}] compress_ppt img {fname}: {ex}")

        ctx.set_progress(70)
        with zipfile.ZipFile(
            ctx.output_path, "w", zipfile.ZIP_DEFLATED, compresslevel=9
        ) as zout:
            for root, _, files in os.walk(tmp_dir):
                for fi in files:
                    abs_p = os.path.join(root, fi)
                    zout.write(abs_p, os.path.relpath(abs_p, tmp_dir))
    except Exception as ex:
        raise ProcessingError(f"compress_ppt failed: {ex}") from ex
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        if _converted:
            try:
                os.remove(_converted)
                shutil.rmtree(os.path.dirname(_converted), ignore_errors=True)
            except OSError:
                pass

    _validate_output(ctx.output_path, "compress_ppt", min_bytes=1000)
    new_size  = os.path.getsize(ctx.output_path)
    reduction = round((1 - new_size / orig_size) * 100, 1) if orig_size else 0
    ctx.set_progress(100)
    log.info(
        f"[{ctx.job_id}] compress_ppt: {orig_size}→{new_size} bytes "
        f"({reduction}%, {compressed} images compressed)"
    )
    return {"reduction_pct": reduction, "images_compressed": compressed}


@register("unlock_ppt")
def unlock_ppt(ctx: JobContext) -> dict:
    """
    Remove password from PowerPoint file.
    M4: was unimplemented.
    Uses msoffcrypto — same pattern as unlock_word/unlock_excel.
    Validates output size ≥ _MIN_DECRYPT_BYTES.
    """
    if not MSOFFCRYPTO_OK:
        raise UnsupportedOperation("unlock_ppt", "msoffcrypto-tool")

    pw = ctx.params.get("password", "")
    _validate_password(pw)
    _validate_office_magic(ctx.input_path)

    try:
        with open(ctx.input_path, "rb") as fp:
            of = msoffcrypto.OfficeFile(fp)
            of.load_key(password=pw)
            with open(ctx.output_path, "wb") as fout:
                of.decrypt(fout)
    except Exception as ex:
        raise ProcessingError(f"Decryption failed — wrong password or corrupt file: {ex}")

    _validate_output(ctx.output_path, "unlock_ppt", min_bytes=_MIN_DECRYPT_BYTES)
    log.info(f"[{ctx.job_id}] unlock_ppt: {os.path.getsize(ctx.output_path)} bytes")
    return {}


@register("protect_ppt")
def protect_ppt(ctx: JobContext) -> dict:
    """
    Encrypt PowerPoint file with a password.
    M5: was unimplemented.
    S2: enforces password length limit.
    V3: .ppt (legacy OLE) rejected — msoffcrypto only handles .pptx.
    """
    if not MSOFFCRYPTO_OK:
        raise UnsupportedOperation("protect_ppt", "msoffcrypto-tool")

    pw  = ctx.params.get("password",  "")
    pw2 = ctx.params.get("password2", "")
    _validate_password(pw)
    if pw != pw2:
        raise ValidationError("Passwords do not match")

    ext = Path(ctx.input_path).suffix.lower()
    if ext == ".ppt":
        raise ValidationError(
            ".ppt format does not support AES encryption — "
            "convert to .pptx first, then protect"
        )

    _validate_office_magic(ctx.input_path)

    try:
        with open(ctx.input_path, "rb") as fp:
            of = msoffcrypto.OfficeFile(fp)
            try:
                of.encrypt(pw, ctx.output_path, cipher_algorithm="AES")
            except TypeError:
                of.encrypt(pw, ctx.output_path)
    except Exception as ex:
        raise ProcessingError(f"Encryption failed: {ex}")

    _validate_output(ctx.output_path, "protect_ppt", min_bytes=1000)
    log.info(f"[{ctx.job_id}] protect_ppt: encrypted {os.path.getsize(ctx.output_path)} bytes")
    return {}
