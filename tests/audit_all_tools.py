"""
Enterprise tool audit — runs every tool, verifies the output is valid
and meaningful, prints PASS / PARTIAL / FAIL with the exact issue.
"""
import os, json, time, base64, urllib.request, urllib.error, io
from typing import Tuple

BASE = "http://localhost:5000"

# ─────────────────────────────────────────────────────────────── fixtures
def build_fixtures():
    """Build a realistic 2-page text PDF, plus a tiny DOCX/XLSX/PPTX/JPG/PNG."""
    import fitz
    d = fitz.open()
    for i in range(3):
        p = d.new_page(width=595, height=842)
        p.insert_text((72, 80),  f"PDFWala Audit Document — Page {i+1}",
                      fontsize=20, fontname="hebo", color=(0,0,0))
        p.insert_text((72, 120), "This is a paragraph of body text with multiple sentences. "
                                 "It contains enough content to make conversion meaningful. "
                                 "Some words include numbers like 42 and 1337.",
                      fontsize=11, fontname="helv")
        p.insert_text((72, 200), "Section heading", fontsize=14, fontname="hebo")
        p.insert_text((72, 230), "Another paragraph here for body content.",
                      fontsize=11, fontname="helv")
    d.save("/tmp/audit/test.pdf"); d.close()

    from docx import Document
    doc = Document(); doc.add_heading("PDFWala Test", level=1)
    doc.add_paragraph("Body paragraph one with several words.")
    doc.add_paragraph("Body paragraph two with more content for testing.")
    tbl = doc.add_table(rows=2, cols=2)
    tbl.rows[0].cells[0].text = "Header A"; tbl.rows[0].cells[1].text = "Header B"
    tbl.rows[1].cells[0].text = "Row 1A";   tbl.rows[1].cells[1].text = "Row 1B"
    doc.save("/tmp/audit/test.docx")

    import openpyxl
    wb = openpyxl.Workbook(); ws = wb.active
    ws.append(["Name", "Score", "Grade"])
    ws.append(["Alice", 95, "A"]); ws.append(["Bob", 72, "C"])
    wb.save("/tmp/audit/test.xlsx")

    from pptx import Presentation
    p = Presentation(); s = p.slides.add_slide(p.slide_layouts[0])
    s.shapes.title.text = "PDFWala"
    s.placeholders[1].text = "Sample slide content"
    p.save("/tmp/audit/test.pptx")

    from PIL import Image
    Image.new("RGB", (800, 600), "red").save("/tmp/audit/test.jpg",  "JPEG", quality=90)
    Image.new("RGB", (800, 600), "blue").save("/tmp/audit/test.png", "PNG")

    # An image-only PDF for OCR
    d = fitz.open()
    pg = d.new_page(width=595, height=842)
    img = Image.new("RGB", (1190, 800), "white")
    from PIL import ImageDraw, ImageFont
    draw = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 60)
    except Exception:
        font = ImageFont.load_default()
    draw.text((60, 80), "SCANNED PDF TEXT", fill="black", font=font)
    draw.text((60, 200), "Hello World 12345", fill="black", font=font)
    buf = io.BytesIO(); img.save(buf, "PNG")
    pg.insert_image(fitz.Rect(0, 0, 595, 421), stream=buf.getvalue())
    d.save("/tmp/audit/scanned.pdf"); d.close()

    # A 2nd PDF for merge/compare
    d = fitz.open(); pg = d.new_page()
    pg.insert_text((72, 80), "Second document for merge/compare", fontsize=18)
    d.save("/tmp/audit/test2.pdf"); d.close()

    # An "edit me" PDF
    d = fitz.open(); pg = d.new_page()
    pg.insert_text((72, 80), "EDIT ME PDF", fontsize=22, fontname="hebo")
    pg.insert_text((72, 120), "Original line text.", fontsize=12)
    d.save("/tmp/audit/edit.pdf"); d.close()


# ─────────────────────────────────────────────────────────────── HTTP
def post(path, fields=None, files=None, multi_field=None):
    boundary = "----A"
    body = io.BytesIO()
    for k, v in (fields or {}).items():
        body.write(f"--{boundary}\r\nContent-Disposition: form-data; name=\"{k}\"\r\n\r\n{v}\r\n".encode())
    if multi_field:
        for fn, data, ct in multi_field:
            body.write(f"--{boundary}\r\nContent-Disposition: form-data; name=\"files\"; filename=\"{fn}\"\r\nContent-Type: {ct}\r\n\r\n".encode())
            body.write(data); body.write(b"\r\n")
    for k, (fn, data, ct) in (files or {}).items():
        body.write(f"--{boundary}\r\nContent-Disposition: form-data; name=\"{k}\"; filename=\"{fn}\"\r\nContent-Type: {ct}\r\n\r\n".encode())
        body.write(data); body.write(b"\r\n")
    body.write(f"--{boundary}--\r\n".encode())
    req = urllib.request.Request(BASE + path, data=body.getvalue(),
                                 headers={"Content-Type": f"multipart/form-data; boundary={boundary}"})
    try:
        r = urllib.request.urlopen(req, timeout=300)
        return r.status, r.read(), r.headers.get_content_type()
    except urllib.error.HTTPError as e:
        return e.code, e.read(), e.headers.get_content_type()
    except Exception as ex:
        return 0, str(ex).encode(), "text/plain"


def get(path):
    try:
        r = urllib.request.urlopen(BASE + path, timeout=300)
        return r.status, r.read(), r.headers.get_content_type()
    except urllib.error.HTTPError as e:
        return e.code, e.read(), e.headers.get_content_type()


# ─────────────────────────────────────────────────────────────── helpers
def upload(path, mime="application/pdf"):
    return (os.path.basename(path), open(path, "rb").read(), mime)


def run_sync(name, route, fields, files):
    st, body, ct = post(route, fields=fields, files=files)
    if st not in (200, 202) or ct != "application/json":
        return None, f"HTTP {st} (ct={ct}) {body[:120]!r}"
    try:
        d = json.loads(body)
    except Exception:
        return None, f"bad JSON: {body[:120]!r}"
    if not d.get("success"):
        return None, f"success=false err={d.get('error','')[:160]}"
    # Async tools (pdf_routes._ALWAYS_ASYNC) reply HTTP 202 + status_url.
    # Poll until completed, then return the final payload as if sync.
    if st == 202 or d.get("status") == "queued":
        status_url = d.get("status_url")
        if not status_url:
            return None, f"async without status_url: {d!r}"
        deadline = time.time() + 600   # 10 min cap per tool
        while time.time() < deadline:
            time.sleep(2)
            ps, pb, _ = get(status_url)
            if ps != 200:
                return None, f"poll HTTP {ps}"
            try:
                pd = json.loads(pb)
            except Exception:
                return None, f"poll bad JSON: {pb[:120]!r}"
            status = pd.get("status")
            if status == "completed":
                return pd, None
            if status in ("failed", "timeout"):
                return None, f"async {status}: {pd.get('error','')[:160]}"
        return None, "async poll timed out after 10 min"
    return d, None


def fetch_output(d):
    """Download the output and return its bytes (or None if missing)."""
    url = d.get("download_url")
    if not url:
        return None
    st, body, ct = get(url)
    if st != 200 or not body:
        return None
    return body


# Verifiers per file type
def verify_pdf(blob):
    if not blob or not blob[:4].startswith(b"%PDF"):
        return False, "not a PDF (missing %PDF magic)"
    try:
        import fitz
        d = fitz.open(stream=blob, filetype="pdf")
        n = len(d); d.close()
        if n < 1: return False, "0 pages"
        return True, f"{n} pages"
    except Exception as ex:
        return False, f"open failed: {ex}"


def verify_docx(blob):
    if not blob or not blob[:2] == b"PK":
        return False, "not a DOCX (PK magic)"
    try:
        import zipfile
        z = zipfile.ZipFile(io.BytesIO(blob))
        names = z.namelist()
        if "word/document.xml" not in names:
            return False, "no word/document.xml"
        xml = z.read("word/document.xml").decode("utf-8", errors="replace")
        z.close()
        return True, f"{len(names)} parts, {xml.count('<w:p')} paragraphs"
    except Exception as ex:
        return False, f"zip parse failed: {ex}"


def verify_xlsx(blob):
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(blob), data_only=False, read_only=True)
        sheets = wb.sheetnames
        # Count any cell with a value to confirm extraction worked
        n_rows = sum(1 for sh in sheets for _ in wb[sh].iter_rows(min_row=1, max_row=10))
        wb.close()
        return True, f"{len(sheets)} sheet(s), at least {n_rows} rows scanned"
    except Exception as ex:
        return False, f"xlsx open failed: {ex}"


def verify_pptx(blob):
    try:
        from pptx import Presentation
        p = Presentation(io.BytesIO(blob))
        return True, f"{len(p.slides)} slides"
    except Exception as ex:
        return False, f"pptx open failed: {ex}"


def verify_zip(blob, expect_min_entries=1):
    try:
        import zipfile
        z = zipfile.ZipFile(io.BytesIO(blob))
        n = len(z.namelist()); z.close()
        if n < expect_min_entries:
            return False, f"only {n} entries, expected ≥{expect_min_entries}"
        return True, f"{n} entries"
    except Exception as ex:
        return False, f"zip open failed: {ex}"


def verify_image(blob, fmt=None):
    try:
        from PIL import Image
        img = Image.open(io.BytesIO(blob)); img.load()
        return True, f"{img.format} {img.width}x{img.height}"
    except Exception as ex:
        return False, f"image open failed: {ex}"


def verify_csv(blob):
    if not blob: return False, "empty"
    txt = blob.decode("utf-8", errors="replace")
    if "," not in txt: return False, "no commas"
    lines = [l for l in txt.splitlines() if l.strip()]
    return True, f"{len(lines)} lines"


def verify_txt(blob):
    if not blob: return False, "empty"
    return True, f"{len(blob)} bytes"


def verify_json(blob):
    try:
        json.loads(blob); return True, "valid JSON"
    except Exception as ex:
        return False, f"bad JSON: {ex}"


def verify_html(blob):
    if not blob: return False, "empty"
    txt = blob.decode("utf-8", errors="replace")
    if "<" not in txt: return False, "no tags"
    return True, f"{len(txt)} chars"


VERIFY = {
    "pdf": verify_pdf, "docx": verify_docx, "xlsx": verify_xlsx,
    "pptx": verify_pptx, "zip": lambda b: verify_zip(b),
    "jpg": verify_image, "jpeg": verify_image, "png": verify_image,
    "csv": verify_csv, "txt": verify_txt, "json": verify_json,
    "html": verify_html,
}


# ─────────────────────────────────────────────────────────────── tests
results = []

def case(name, route, fields, files, output_ext, **extra):
    d, err = run_sync(name, route, fields, files)
    if d is None:
        results.append((name, "FAIL", err))
        return
    fname = d.get("filename", "")
    ext = (fname.rsplit(".", 1)[-1] or output_ext).lower()
    blob = fetch_output(d)
    if not blob:
        results.append((name, "FAIL", f"output download failed for {fname}"))
        return
    verifier = VERIFY.get(ext)
    if not verifier:
        results.append((name, "PASS?", f"no verifier for ext={ext}; size={len(blob)}"))
        return
    ok, detail = verifier(blob)
    extra_str = ""
    # custom extra checks per tool
    cb = extra.get("check")
    if cb:
        ok2, more = cb(blob, d)
        if ok and not ok2: ok = False
        extra_str = " | " + more
    results.append((name, "PASS" if ok else "PARTIAL", f"{detail}{extra_str}"))


def case_canvas():
    """Canvas editor — parse + save round-trip."""
    pdf = open("/tmp/audit/edit.pdf","rb").read()
    st, body, ct = post("/api/pdf/parse-canvas", files={"file": ("edit.pdf", pdf, "application/pdf")})
    if st != 200:
        results.append(("canvas-parse", "FAIL", f"HTTP {st}: {body[:120]!r}")); return
    info = json.loads(body); spans = info["pages"][0]["spans"]
    target = next((s for s in spans if "EDIT" in s["text"]), spans[0])
    change = dict(target); change["page"]=0; change["new_text"]="ENTERPRISE TEST"
    st2, b2, ct2 = post("/api/pdf/save-canvas",
        fields={"changes": json.dumps([change]), "scanned": "false"},
        files={"file": ("edit.pdf", pdf, "application/pdf")})
    if st2 != 200 or ct2 != "application/pdf":
        results.append(("canvas-save", "FAIL", f"HTTP {st2}: {b2[:120]!r}")); return
    import fitz
    d2 = fitz.open(stream=b2, filetype="pdf")
    text = d2[0].get_text(); d2.close()
    ok = "ENTERPRISE TEST" in text and "EDIT ME PDF" not in text
    results.append(("canvas-edit", "PASS" if ok else "PARTIAL",
                    f"edit_applied={ok} text={text[:80]!r}"))


# ─────────────────────────────────────────────────────────────── run
def main():
    os.makedirs("/tmp/audit", exist_ok=True)
    build_fixtures()
    pdf  = upload("/tmp/audit/test.pdf")
    pdf2 = upload("/tmp/audit/test2.pdf")
    scan = upload("/tmp/audit/scanned.pdf")
    docx = upload("/tmp/audit/test.docx",
                  "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
    xlsx = upload("/tmp/audit/test.xlsx",
                  "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
    pptx = upload("/tmp/audit/test.pptx",
                  "application/vnd.openxmlformats-officedocument.presentationml.presentation")
    jpg  = upload("/tmp/audit/test.jpg", "image/jpeg")
    png  = upload("/tmp/audit/test.png", "image/png")

    # PDF tools
    case("merge",         "/api/pdf/merge",      {}, None,  "pdf",
         **{}); # placeholder; do merge with multi field
    # merge needs multi-field "files"
    st, body, ct = post("/api/pdf/merge", multi_field=[("a.pdf", pdf[1], "application/pdf"),
                                                        ("b.pdf", pdf2[1], "application/pdf")])
    if st != 200:
        results[-1] = ("merge", "FAIL", f"HTTP {st}: {body[:120]!r}")
    else:
        d = json.loads(body)
        if not d.get("success"):
            results[-1] = ("merge", "FAIL", f"success=false err={d.get('error','')[:160]}")
        else:
            blob = fetch_output(d)
            ok, det = verify_pdf(blob or b"")
            import fitz
            try:
                pgs = len(fitz.open(stream=blob, filetype="pdf"))
            except Exception:
                pgs = 0
            results[-1] = ("merge", "PASS" if ok and pgs >= 4 else "PARTIAL",
                           f"{det} (expected ≥4 pages, got {pgs})")
    case("split",         "/api/pdf/split",      {"mode":"all"}, {"file":pdf}, "zip",
         check=lambda b,d: (True, f"zip ok"))
    case("compress",      "/api/pdf/compress",   {"level":"medium"}, {"file":pdf}, "pdf")
    case("rotate",        "/api/pdf/rotate",     {"angle":"90"}, {"file":pdf}, "pdf")
    case("watermark",     "/api/pdf/watermark",  {"text":"CONFIDENTIAL","opacity":"30"},
         {"file":pdf}, "pdf",
         check=lambda b,d: ("CONFIDENTIAL" in __import__('fitz').open(stream=b,filetype='pdf')[0].get_text(),
                            "watermark text present" if "CONFIDENTIAL" in __import__('fitz').open(stream=b,filetype='pdf')[0].get_text() else "watermark text MISSING"))
    case("page-numbers",  "/api/pdf/page-numbers",{"position":"bottom","start":"1"},
         {"file":pdf}, "pdf")
    case("extract-pages", "/api/pdf/extract-pages",{"order":"1-2"}, {"file":pdf}, "pdf",
         check=lambda b,d: (len(__import__('fitz').open(stream=b,filetype='pdf'))==2,
                            f"{len(__import__('fitz').open(stream=b,filetype='pdf'))} pages (want 2)"))
    case("remove-pages",  "/api/pdf/remove-pages",{"order":"1"}, {"file":pdf}, "pdf",
         check=lambda b,d: (len(__import__('fitz').open(stream=b,filetype='pdf'))==2,
                            f"{len(__import__('fitz').open(stream=b,filetype='pdf'))} pages (want 2)"))
    case("organize",      "/api/pdf/organize",   {"order":"3,1,2"}, {"file":pdf}, "pdf")
    case("crop",          "/api/pdf/crop",       {"left":"5","right":"5","top":"5","bottom":"5"},
         {"file":pdf}, "pdf")
    case("repair",        "/api/pdf/repair",     {}, {"file":pdf}, "pdf")
    case("linearize",     "/api/pdf/linearize",  {}, {"file":pdf}, "pdf")
    case("info",          "/api/pdf/info",       {}, {"file":pdf}, "json")
    case("protect",       "/api/pdf/protect",    {"password":"secret"}, {"file":pdf}, "pdf",
         check=lambda b,d: (__import__('fitz').open(stream=b,filetype='pdf').is_encrypted,
                            "encrypted" if __import__('fitz').open(stream=b,filetype='pdf').is_encrypted else "NOT encrypted"))
    case("sign",          "/api/pdf/sign",       {"name":"John Doe"}, {"file":pdf}, "pdf")
    case("redact",        "/api/pdf/redact",     {"search_text":"PDFWala"}, {"file":pdf}, "pdf",
         check=lambda b,d: ("PDFWala" not in __import__('fitz').open(stream=b,filetype='pdf')[0].get_text(),
                            "redacted out" if "PDFWala" not in __import__('fitz').open(stream=b,filetype='pdf')[0].get_text() else "STILL VISIBLE"))
    case("ocr",           "/api/pdf/ocr",        {"lang":"eng"}, {"file":scan}, "pdf",
         check=lambda b,d: (any("scanned" in __import__('fitz').open(stream=b,filetype='pdf')[0].get_text().lower() or
                                 "SCANNED" in __import__('fitz').open(stream=b,filetype='pdf')[0].get_text()
                                 for _ in [0]),
                            f"ocr text: {__import__('fitz').open(stream=b,filetype='pdf')[0].get_text()[:80]!r}"))
    case("pdf-to-word",   "/api/pdf/to-word",    {}, {"file":pdf}, "docx",
         check=lambda b,d: (__import__('zipfile').ZipFile(io.BytesIO(b)).read("word/document.xml").decode("utf-8","replace").count("PDFWala") >= 1,
                            "PDFWala text preserved" if __import__('zipfile').ZipFile(io.BytesIO(b)).read("word/document.xml").decode("utf-8","replace").count("PDFWala") >= 1 else "TEXT MISSING"))
    case("pdf-to-excel",  "/api/pdf/to-excel",   {}, {"file":pdf}, "xlsx")
    case("pdf-to-ppt",    "/api/pdf/to-ppt",     {}, {"file":pdf}, "pptx")
    case("pdf-to-image",  "/api/pdf/to-image",   {"format":"jpg","dpi":"150"}, {"file":pdf}, "zip",
         check=lambda b,d: (verify_zip(b, expect_min_entries=3)[0],
                            verify_zip(b, expect_min_entries=3)[1]))
    case("pdf-to-jpg",    "/api/pdf/to-jpg",     {}, {"file":pdf}, "zip")
    case("pdf-to-png",    "/api/pdf/to-png",     {}, {"file":pdf}, "zip")
    case("pdf-to-pdfa",   "/api/pdf/to-pdfa",    {}, {"file":pdf}, "pdf")

    # Canvas editor
    case_canvas()

    # Office tools
    case("word-to-pdf",   "/api/office/word/to-pdf", {}, {"file":docx}, "pdf",
         check=lambda b,d: ("PDFWala Test" in __import__('fitz').open(stream=b,filetype='pdf')[0].get_text(),
                            "heading preserved" if "PDFWala Test" in __import__('fitz').open(stream=b,filetype='pdf')[0].get_text() else "TEXT MISSING"))
    case("excel-to-pdf",  "/api/office/excel/to-pdf", {}, {"file":xlsx}, "pdf")
    case("ppt-to-pdf",    "/api/office/ppt/to-pdf",   {}, {"file":pptx}, "pdf")
    case("word-to-html",  "/api/office/word/to-html", {}, {"file":docx}, "html")
    case("word-to-txt",   "/api/office/word/to-txt",  {}, {"file":docx}, "txt")
    case("word-to-jpg",   "/api/office/word/to-jpg",  {}, {"file":docx}, "zip")
    case("word-to-png",   "/api/office/word/to-png",  {}, {"file":docx}, "zip")
    case("word-to-excel", "/api/office/word/to-excel",{}, {"file":docx}, "xlsx")
    case("word-to-ppt",   "/api/office/word/to-ppt",  {}, {"file":docx}, "pptx")
    case("excel-to-csv",  "/api/office/excel/to-csv", {}, {"file":xlsx}, "csv")
    case("excel-to-word", "/api/office/excel/to-word",{}, {"file":xlsx}, "docx")
    case("excel-to-json", "/api/office/excel/to-json",{}, {"file":xlsx}, "json")
    case("excel-to-jpg",  "/api/office/excel/to-jpg", {}, {"file":xlsx}, "zip")
    case("excel-to-html", "/api/office/excel/to-html",{}, {"file":xlsx}, "html")
    case("ppt-to-jpg",    "/api/office/ppt/to-jpg",   {}, {"file":pptx}, "zip")

    # Image tools
    case("img-compress",  "/api/image/compress", {"quality":"70"}, {"file":jpg}, "jpg")
    case("img-resize",    "/api/image/resize",   {"width":"400","height":"300"}, {"file":jpg}, "jpg",
         check=lambda b,d: (__import__('PIL.Image', fromlist=['Image']).open(io.BytesIO(b)).size == (400,300),
                            f"size={__import__('PIL.Image', fromlist=['Image']).open(io.BytesIO(b)).size}"))
    case("img-convert",   "/api/image/convert",  {"format":"png"}, {"file":jpg}, "png")
    case("png-to-jpg",    "/api/image/png-to-jpg", {}, {"file":png}, "jpg")
    case("img-to-pdf",    "/api/image/to-pdf",   {}, {"file":jpg}, "pdf")

    # Report
    print("=" * 72)
    pass_n = sum(1 for _,s,_ in results if s=="PASS")
    part_n = sum(1 for _,s,_ in results if s=="PARTIAL" or s=="PASS?")
    fail_n = sum(1 for _,s,_ in results if s=="FAIL")
    for n,s,m in results:
        marker = {"PASS":"✓","PARTIAL":"~","PASS?":"?","FAIL":"✗"}[s]
        print(f"  {marker} {s:7s} {n:18s} {m}")
    print("=" * 72)
    print(f"  TOTAL: {len(results)} | PASS: {pass_n} | PARTIAL/PASS?: {part_n} | FAIL: {fail_n}")

main()
